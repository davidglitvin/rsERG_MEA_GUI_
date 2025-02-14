import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math

import ipywidgets as widgets
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed, PPT export will be disabled.")

import mne
from tqdm.notebook import tqdm
from scipy import signal

###############################################################################
# 1) HELPER FUNCTIONS
###############################################################################

def group_keys_by_rows(psd_keys, row_size=4):
    """
    Utility to chunk the psd_keys list into sublists of length row_size,
    for the old approach if no block-based keys are present.
    """
    rows = []
    for i in range(0, len(psd_keys), row_size):
        rows.append(psd_keys[i:i+row_size])
    return rows

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    """
    Identifies which PSD traces (epochs) should be excluded based on:
      1) Low-frequency outlier in [low_band_min, low_band_max].
      2) Repeated outliers in test_bands >= test_band_count_threshold times.
    Returns (kept_indices, excluded_indices).
    """
    if test_band_count_threshold is None:
        test_band_count_threshold = len(test_bands) // 2

    mean_psd = np.mean(psd_array, axis=0)
    excluded_traces = []
    kept_traces = []

    low_band_indices = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0]
        for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # 1) Low-band outlier
        if np.any(trace[low_band_indices] > low_band_threshold * mean_psd[low_band_indices]):
            excluded_traces.append(i)
            continue

        # 2) Repeated band outliers
        suprathreshold_count = 0
        for indices in band_indices:
            if np.any(trace[indices] > test_band_threshold * mean_psd[indices]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces

def plot_psds_with_exclusion(
    ax,
    psd_array,
    freqs,
    kept_traces,
    excluded_traces,
    original_mean_psd,
    title,
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    alpha_excluded=0.05,
    alpha_kept=0.7,
    title_fontsize=8,
    axis_label_fontsize=8,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6
):
    if len(title) > max_title_length:
        title = title[:max_title_length] + "..."

    # Plot kept
    if show_kept:
        for idx_i, idx in enumerate(kept_traces):
            label = "Kept Trace" if idx_i == 0 else None
            ax.plot(freqs, psd_array[idx], color=color_kept, alpha=alpha_kept, label=label)

    # Original mean
    if show_original_mean:
        ax.plot(freqs, original_mean_psd, color=color_old_mean, linewidth=2, label="Original Mean")

    # New mean
    if show_new_mean and kept_traces:
        new_mean_psd = np.mean(psd_array[kept_traces], axis=0)
        ax.plot(freqs, new_mean_psd, color=color_new_mean, linewidth=2, label="New Mean")

    # Excluded
    if show_excluded:
        for idx_j, idx in enumerate(excluded_traces):
            label = "Excluded Trace" if idx_j == 0 else None
            ax.plot(freqs, psd_array[idx], color=color_excluded, alpha=alpha_excluded, zorder=10, label=label)

    # Vertical lines
    if show_vertical_lines and vertical_lines is not None:
        for vfreq in vertical_lines:
            ax.axvline(
                vfreq,
                color=vertical_line_color,
                linestyle=vertical_line_style,
                alpha=vertical_line_alpha
            )

    ax.set_xlabel("Frequency (Hz)", fontsize=axis_label_fontsize)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_label_fontsize)
    ax.set_title(title, fontsize=title_fontsize)
    ax.legend(loc="upper right", fontsize=legend_fontsize)
    ax.tick_params(axis='both', labelsize=tick_label_fontsize)

def exclude_and_plot(
    ax,
    title,
    psd_data,
    freqs,
    low_band=(1,3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None,
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    alpha_excluded=0.05,
    alpha_kept=0.7,
    title_fontsize=8,
    axis_label_fontsize=8,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6
):
    """
    Helper that calls exclude_traces & plot_psds_with_exclusion in one go.
    Returns (kept_indices, excluded_indices).
    """
    if psd_data is None or freqs is None:
        ax.text(0.5, 0.5, f"No PSD data\n{title}", ha='center', va='center')
        return [], []

    original_mean_psd = np.mean(psd_data, axis=0)
    kept, excl = exclude_traces(
        psd_data,
        freqs,
        low_band=low_band,
        low_band_threshold=low_band_threshold,
        test_bands=test_bands,
        test_band_threshold=test_band_threshold,
        test_band_count_threshold=test_band_count_threshold
    )
    plot_psds_with_exclusion(
        ax=ax,
        psd_array=psd_data,
        freqs=freqs,
        kept_traces=kept,
        excluded_traces=excl,
        original_mean_psd=original_mean_psd,
        title=title,
        show_kept=show_kept,
        show_excluded=show_excluded,
        show_original_mean=show_original_mean,
        show_new_mean=show_new_mean,
        show_vertical_lines=show_vertical_lines,
        color_kept=color_kept,
        color_excluded=color_excluded,
        color_old_mean=color_old_mean,
        color_new_mean=color_new_mean,
        alpha_excluded=alpha_excluded,
        alpha_kept=alpha_kept,
        title_fontsize=title_fontsize,
        axis_label_fontsize=axis_label_fontsize,
        legend_fontsize=legend_fontsize,
        tick_label_fontsize=tick_label_fontsize,
        max_title_length=max_title_length,
        vertical_lines=vertical_lines,
        vertical_line_color=vertical_line_color,
        vertical_line_style=vertical_line_style,
        vertical_line_alpha=vertical_line_alpha
    )
    return kept, excl


###############################################################################
# 2) Main GUI
###############################################################################

def build_exportable_plot_psd_gui():
    """
    A revised GUI that:
      - Loads a PSD pickle containing either block-based keys (like "block1:Ch1")
        or normal keys (like "Ch1", "Ch2", etc.)
      - If block-based => row=channels, col=blocks
      - If single => uses old row/col approach with group_keys_by_rows & num_cols
      - Plots threshold-based outlier dropping
      - Exports figures or cleaned PSD
    """

    # ~~~~ UI ELEMENTS ~~~~
    load_psd_button = widgets.Button(description='Load PSD Pickle', button_style='info')
    psd_file_chooser = FileChooser(os.getcwd(), title='Select PSD Pickle File', select_default=False)
    psd_file_chooser.show_only_files = True
    psd_file_chooser.filter_pattern = ['*.pkl']

    loaded_psd = {}
    current_figures = []
    kept_indices_dict = {}

    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description='Select Keys:',
        layout=widgets.Layout(width='320px', height='240px', overflow_y='auto')
    )

    # Booleans for plot
    show_kept_cb = widgets.Checkbox(value=True, description='Show Kept')
    show_excluded_cb = widgets.Checkbox(value=True, description='Show Excl')
    show_original_mean_cb = widgets.Checkbox(value=True, description='Orig Mean')
    show_new_mean_cb = widgets.Checkbox(value=True, description='New Mean')
    show_vertical_lines_cb = widgets.Checkbox(value=True, description='Vertical Lines')

    vertical_lines_text = widgets.Text(value="10,15", description='Lines(Hz):', layout=widgets.Layout(width='150px'))

    # Axis
    x_min_widget = widgets.FloatText(value=None, description='X Min:', layout=widgets.Layout(width='220px'))
    x_max_widget = widgets.FloatText(value=None, description='X Max:', layout=widgets.Layout(width='220px'))
    y_min_widget = widgets.FloatText(value=None, description='Y Min:', layout=widgets.Layout(width='220px'))
    y_max_widget = widgets.FloatText(value=None, description='Y Max:', layout=widgets.Layout(width='220px'))

    # Fonts
    title_fs_widget = widgets.IntText(value=8, description='TitleFS:', layout=widgets.Layout(width='200px'))
    axis_fs_widget = widgets.IntText(value=8, description='AxisFS:', layout=widgets.Layout(width='200px'))
    legend_fs_widget = widgets.IntText(value=8, description='LegendFS:', layout=widgets.Layout(width='200px'))
    tick_fs_widget = widgets.IntText(value=8, description='TickFS:', layout=widgets.Layout(width='200px'))
    max_title_len_widget = widgets.IntText(value=40, description='MaxTitle:', layout=widgets.Layout(width='200px'))

    # Colors
    color_kept_widget = widgets.ColorPicker(value='lightgray', description='Kept:', layout=widgets.Layout(width='220px'))
    color_excl_widget = widgets.ColorPicker(value='red', description='Excl:', layout=widgets.Layout(width='220px'))
    color_oldmean_widget = widgets.ColorPicker(value='blue', description='OrigM:', layout=widgets.Layout(width='220px'))
    color_newmean_widget = widgets.ColorPicker(value='green', description='NewM:', layout=widgets.Layout(width='220px'))

    # Exclusion thresholds
    low_band_min_widget = widgets.FloatText(value=1.0, description='LBMin:', layout=widgets.Layout(width='200px'))
    low_band_max_widget = widgets.FloatText(value=3.0, description='LBMax:', layout=widgets.Layout(width='200px'))
    low_band_thr_widget = widgets.FloatText(value=3.0, description='LBThr:', layout=widgets.Layout(width='200px'))
    test_band_thr_widget = widgets.FloatText(value=10.0, description='TBThr:', layout=widgets.Layout(width='200px'))

    test_band_text = widgets.Textarea(
        value="(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)",
        description='Test Bands:',
        layout=widgets.Layout(width='250px', height='60px')
    )

    # If single => old approach uses num_cols
    num_cols_widget = widgets.IntText(value=4, description='NumCols:', layout=widgets.Layout(width='200px'))

    # Buttons
    plot_psd_button = widgets.Button(description='Plot PSD(s)', button_style='success')
    export_filename_widget = widgets.Text(value='MyPSDExport', description='FileBase:', layout=widgets.Layout(width='200px'))
    export_format_widget = widgets.Dropdown(options=['png','svg','jpeg','ppt'], value='png', description='Fmt:', layout=widgets.Layout(width='200px'))
    export_button = widgets.Button(description='Export Figs', button_style='warning')
    export_cleaned_button = widgets.Button(description='Export Cleaned', button_style='info')

    # Outputs
    load_output_area = widgets.Output()
    plot_output_area = widgets.Output()
    export_output_area = widgets.Output()

    # ~~~~ A) LOAD PSD
    def on_load_psd_clicked(b):
        with load_output_area:
            clear_output()
            chosen_file = psd_file_chooser.selected
            if not chosen_file:
                print("No .pkl file selected.")
                return
            if not os.path.isfile(chosen_file):
                print(f"File not found => {chosen_file}")
                return
            try:
                with open(chosen_file,'rb') as f:
                    data_ = pickle.load(f)
                loaded_psd.clear()
                loaded_psd.update(data_)
                print(f"Loaded PSD => {chosen_file}")
                all_keys = sorted(loaded_psd.keys())
                print("Found keys:", all_keys)
                channels_dropdown.options = all_keys
            except Exception as e:
                print(f"Error loading PSD pickle: {e}")

    load_psd_button.on_click(on_load_psd_clicked)

    # ~~~~ B) PLOT PSD
    def on_plot_psd_clicked(b):
        nonlocal current_figures, kept_indices_dict
        with plot_output_area:
            clear_output()
            current_figures = []
            kept_indices_dict = {}

            selected_keys = list(channels_dropdown.value)
            if not selected_keys:
                print("No keys selected.")
                return
            # Check if block-based
            has_blocks = any(":" in k for k in selected_keys)

            # Gather thresholds
            lb_min = low_band_min_widget.value
            lb_max = low_band_max_widget.value
            lb_thr = low_band_thr_widget.value
            tb_thr = test_band_thr_widget.value
            test_str = test_band_text.value.strip()
            test_bands_list = []
            if test_str:
                pairs = test_str.replace(" ","").split(")")
                for p in pairs:
                    p=p.strip(",(").strip()
                    if not p: continue
                    vals = p.split(",")
                    if len(vals)==2:
                        try:
                            lf=float(vals[0])
                            hf=float(vals[1])
                            test_bands_list.append((lf,hf))
                        except:
                            pass
            if not test_bands_list:
                test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]

            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value
            show_kept = show_kept_cb.value
            show_excl = show_excluded_cb.value
            show_orig = show_original_mean_cb.value
            show_newm = show_new_mean_cb.value
            show_vlines = show_vertical_lines_cb.value

            t_fs = title_fs_widget.value
            a_fs = axis_fs_widget.value
            l_fs = legend_fs_widget.value
            tk_fs = tick_fs_widget.value
            mt_len = max_title_len_widget.value

            c_kept = color_kept_widget.value
            c_excl = color_excl_widget.value
            c_old = color_oldmean_widget.value
            c_new = color_newmean_widget.value

            vlines_str = vertical_lines_text.value.strip()
            vert_lines=[]
            if vlines_str:
                for part in vlines_str.split(","):
                    part=part.strip()
                    if part:
                        try: vert_lines.append(float(part))
                        except: pass

            # If block-based => parse block name & channel name => row=channel, col=block
            if has_blocks:
                # parse set of blocks, set of channel names
                blocks=set()
                chans=set()
                for key_ in selected_keys:
                    if ":" in key_:
                        blk,ch=key_.split(":",1)
                        blocks.add(blk)
                        chans.add(ch)
                    else:
                        chans.add(key_)
                blocks=sorted(blocks)
                chans=sorted(chans)

                n_rows=len(chans)
                n_cols=len(blocks)
                fig, axes = plt.subplots(n_rows, n_cols, figsize=(4*n_cols,3*n_rows), squeeze=False)

                for r_idx, chnm in enumerate(chans):
                    for c_idx, blk_ in enumerate(blocks):
                        ax=axes[r_idx,c_idx]
                        psd_key = f"{blk_}:{chnm}"
                        psd_info = loaded_psd.get(psd_key)
                        if not psd_info or psd_info.get('psd') is None or psd_info.get('freqs') is None:
                            ax.text(0.5,0.5,f"No data\n{psd_key}",ha='center',va='center')
                            continue
                        arr=psd_info['psd']
                        fre=psd_info['freqs']
                        kept_ix, excl_ix = exclude_and_plot(
                            ax=ax,
                            title=psd_key,
                            psd_data=arr,
                            freqs=fre,
                            low_band=(lb_min, lb_max),
                            low_band_threshold=lb_thr,
                            test_bands=test_bands_list,
                            test_band_threshold=tb_thr,
                            show_kept=show_kept,
                            show_excluded=show_excl,
                            show_original_mean=show_orig,
                            show_new_mean=show_newm,
                            show_vertical_lines=show_vlines,
                            color_kept=c_kept,
                            color_excluded=c_excl,
                            color_old_mean=c_old,
                            color_new_mean=c_new,
                            title_fontsize=t_fs,
                            axis_label_fontsize=a_fs,
                            legend_fontsize=l_fs,
                            tick_label_fontsize=tk_fs,
                            max_title_length=mt_len,
                            vertical_lines=vert_lines
                        )
                        kept_indices_dict[psd_key]= kept_ix
                        # set range
                        if x_min is not None or x_max is not None:
                            ax.set_xlim(x_min, x_max)
                        if y_min is not None or y_max is not None:
                            ax.set_ylim(y_min, y_max)
                        if r_idx==0:
                            ax.set_title(blk_, fontsize=t_fs)
                        if c_idx==0:
                            ax.set_ylabel(chnm, fontsize=a_fs)
                plt.tight_layout()
                plt.show()
                current_figures=[fig]
                print(f"Plotted => row={n_rows} channels, col={n_cols} blocks => 1 figure.")
            else:
                # no block-based keys => old approach
                print("No block-based keys detected, using old row/col approach.")
                # Build dictionary for selected
                psd_data_dict={}
                for k in selected_keys:
                    if k in loaded_psd:
                        info=loaded_psd[k]
                        if info.get('psd') is not None:
                            psd_data_dict[k]=info

                row_size=num_cols_widget.value
                # group keys into rows
                row_keys_list = group_keys_by_rows(list(psd_data_dict.keys()), row_size=row_size)

                figs=[]
                fig_count=0
                for row_idx, row_of_keys in enumerate(row_keys_list, start=1):
                    num_plots=len(row_of_keys)
                    if num_plots==0:
                        continue
                    n_rows=math.ceil(num_plots/row_size)
                    fig, axes = plt.subplots(n_rows, row_size, figsize=(4*row_size,3*n_rows), squeeze=False)
                    axes=axes.flatten() if (n_rows*row_size>1) else [axes]
                    for i, key_ in enumerate(row_of_keys):
                        ax=axes[i]
                        info=psd_data_dict[key_]
                        arr=info['psd']
                        fre=info['freqs']
                        kept_ix, excl_ix = exclude_and_plot(
                            ax=ax,
                            title=key_,
                            psd_data=arr,
                            freqs=fre,
                            low_band=(low_band_min_widget.value,low_band_max_widget.value),
                            low_band_threshold=low_band_thr_widget.value,
                            test_bands=test_bands_list,
                            test_band_threshold=test_band_thr_widget.value,
                            show_kept=show_kept_cb.value,
                            show_excluded=show_excluded_cb.value,
                            show_original_mean=show_original_mean_cb.value,
                            show_new_mean=show_new_mean_cb.value,
                            show_vertical_lines=show_vertical_lines_cb.value,
                            color_kept=c_kept,
                            color_excluded=c_excl,
                            color_old_mean=c_old,
                            color_new_mean=c_new,
                            title_fontsize=title_fs_widget.value,
                            axis_label_fontsize=axis_fs_widget.value,
                            legend_fontsize=legend_fs_widget.value,
                            tick_label_fontsize=tick_fs_widget.value,
                            max_title_length=max_title_len_widget.value,
                            vertical_lines=vert_lines
                        )
                        kept_indices_dict[key_] = kept_ix
                        if x_min_widget.value is not None or x_max_widget.value is not None:
                            ax.set_xlim(x_min_widget.value, x_max_widget.value)
                        if y_min_widget.value is not None or y_max_widget.value is not None:
                            ax.set_ylim(y_min_widget.value, y_max_widget.value)
                    # hide extras
                    for ax_ in axes[num_plots:]:
                        ax_.axis('off')

                    plt.tight_layout()
                    plt.show()
                    figs.append(fig)
                    fig_count+=1
                current_figures=figs
                print(f"Plotted => {fig_count} figure(s). Kept indices recorded.")
            print("Done plotting. Kept indices => kept_indices_dict")

    plot_psd_button.on_click(on_plot_psd_clicked)

    # ~~~~ C) EXPORT FIGURES
    def on_export_button_clicked(b):
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No figures to export. Plot first!")
                return
            fname_base = export_filename_widget.value.strip()
            fmt=export_format_widget.value.lower()
            if not fname_base:
                print("Provide a base filename.")
                return
            if fmt=='ppt' and not HAS_PPTX:
                print("Cannot export PPT (python-pptx not installed).")
                return
            print(f"Exporting {len(current_figures)} figure(s) as {fmt}...")

            if fmt=='ppt':
                prs=Presentation()
                for i,fig in enumerate(current_figures, start=1):
                    slide=prs.slides.add_slide(prs.slide_layouts[6])
                    tmp_png=f"{fname_base}_{i}.png"
                    fig.savefig(tmp_png, format='png', dpi=150)
                    left=top=Inches(1)
                    slide.shapes.add_picture(tmp_png,left,top,Inches(8),Inches(4.5))
                    os.remove(tmp_png)
                ppt_file=f"{fname_base}.pptx"
                prs.save(ppt_file)
                print(f"Saved => {ppt_file}")
            else:
                for i,fig in enumerate(current_figures, start=1):
                    out_file=f"{fname_base}_{i}.{fmt}"
                    fig.savefig(out_file, format=fmt, dpi=150)
                    print(f"Saved => {out_file}")
            print("Done exporting figures.")

    export_button.on_click(on_export_button_clicked)

    # ~~~~ D) EXPORT CLEANED PSD
    def on_export_cleaned_clicked(b):
        with export_output_area:
            clear_output()
            if not kept_indices_dict:
                print("No kept indices found. Please plot first.")
                return
            if not loaded_psd:
                print("No PSD data loaded.")
                return

            cleaned_data={}
            for key_, psd_info in loaded_psd.items():
                if key_ not in kept_indices_dict:
                    # skip or keep as is => we'll skip
                    continue
                kept_ix = kept_indices_dict[key_]
                arr = psd_info.get('psd')
                fre= psd_info.get('freqs')
                if arr is None or fre is None:
                    continue
                if len(kept_ix)==0:
                    cleaned_data[key_]={"freqs": fre, "psd": None}
                else:
                    cleaned_data[key_]={"freqs": fre, "psd": arr[kept_ix,:]}

            if not cleaned_data:
                print("No cleaned data (possibly everything excluded?).")
                return

            base_name = export_filename_widget.value.strip()
            if not base_name:
                base_name="cleaned_psd"
            outpkl=f"{base_name}_cleaned.pkl"
            try:
                with open(outpkl,'wb') as f:
                    pickle.dump(cleaned_data, f)
                print(f"Saved cleaned PSD => {outpkl}")
                print("Keys in cleaned PSD =>", list(cleaned_data.keys()))
            except Exception as e:
                print(f"Error saving cleaned PSD: {e}")

    export_cleaned_button.on_click(on_export_cleaned_clicked)

    # ~~~~ Layout ~~~~
    load_box = widgets.VBox([
        widgets.HTML("<h3>Load PSD Pickle</h3>"),
        load_psd_button,
        psd_file_chooser,
        load_output_area
    ])

    channel_sel_box = widgets.VBox([
        widgets.HTML("<h3>Keys (Channels/Blocks)</h3>"),
        channels_dropdown
    ])

    exclusion_params_box = widgets.HBox([
        widgets.VBox([
            widgets.Label("LowBand Range:"),
            low_band_min_widget,
            low_band_max_widget,
            low_band_thr_widget
        ]),
        widgets.VBox([
            widgets.Label("TestBand Thr:"),
            test_band_thr_widget,
            widgets.Label("TestBands:"),
            test_band_text
        ])
    ])

    axis_range_box = widgets.HBox([
        widgets.VBox([x_min_widget, x_max_widget]),
        widgets.VBox([y_min_widget, y_max_widget])
    ])

    font_size_box = widgets.HBox([
        title_fs_widget, axis_fs_widget, legend_fs_widget,
        tick_fs_widget, max_title_len_widget
    ])
    color_box = widgets.HBox([
        color_kept_widget, color_excl_widget,
        color_oldmean_widget, color_newmean_widget
    ])

    plot_box = widgets.VBox([
        widgets.HTML("<h3>Plot & Exclusion</h3>"),
        show_kept_cb,
        show_excluded_cb,
        show_original_mean_cb,
        show_new_mean_cb,
        show_vertical_lines_cb,
        vertical_lines_text,
        widgets.HTML("<b>Exclusion Params</b>"),
        exclusion_params_box,
        widgets.HTML("<b>Axes & Layout</b>"),
        axis_range_box,
        widgets.HTML("<b>Font Sizes</b>"),
        font_size_box,
        widgets.HTML("<b>Colors</b>"),
        color_box,
        widgets.HTML("<b>If No Blocks => #Columns</b>"),
        num_cols_widget,
        plot_psd_button,
        plot_output_area
    ])

    export_box = widgets.VBox([
        widgets.HTML("<h3>Export Tools</h3>"),
        widgets.HBox([export_filename_widget, export_format_widget]),
        widgets.HBox([export_button, export_cleaned_button]),
        export_output_area
    ])

    ui = widgets.VBox([
        widgets.HTML("<h2>PSD Clean GUI (Blocks => row=chan, col=block)</h2>"),
        widgets.HBox([load_box, channel_sel_box]),
        plot_box,
        export_box
    ])

    display(ui)

# Build & display the GUI
build_exportable_plot_psd_gui()

