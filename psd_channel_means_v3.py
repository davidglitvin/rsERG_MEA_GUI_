import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math
import re

import ipywidgets as widgets
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed, PPT export will be disabled.")


###############################################################################
# 1) Helper Functions
###############################################################################

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    """
    Returns (kept_indices, excluded_indices) after applying
    the specified thresholds for outlier detection.
    """
    if test_band_count_threshold is None:
        test_band_count_threshold = len(test_bands) // 2

    mean_psd = np.mean(psd_array, axis=0)
    excluded_traces = []
    kept_traces = []

    low_band_indices = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0]
        for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # 1) Check low-frequency extreme outlier
        if np.any(trace[low_band_indices] > low_band_threshold * mean_psd[low_band_indices]):
            excluded_traces.append(i)
            continue

        # 2) Check repeated suprathreshold events in test bands
        suprathreshold_count = 0
        for indices in band_indices:
            if np.any(trace[indices] > test_band_threshold * mean_psd[indices]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces


def compute_channel_means(
    psds_dict,
    exclude=False,
    low_band=(1,3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    """
    Compute channel means before and after exclusion.
    Returns a dict {channel: (original_mean, new_mean)}.
    """
    channel_means = {}
    for channel, data in psds_dict.items():
        psd = data.get('psd', None)
        freqs = data.get('freqs', None)
        if psd is None or freqs is None:
            print(f"Channel '{channel}' is missing 'psd' or 'freqs'. Skipping.")
            continue

        original_mean = np.mean(psd, axis=0)
        if exclude:
            kept_traces, _ = exclude_traces(
                psd_array=psd,
                freqs=freqs,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold
            )
            if kept_traces:
                new_mean = np.mean(psd[kept_traces], axis=0)
            else:
                new_mean = np.zeros_like(original_mean)
        else:
            new_mean = original_mean.copy()
        
        channel_means[channel] = (original_mean, new_mean)
    
    return channel_means


def compute_group_mean(channel_means, selected_channels):
    """
    Compute the mean of 'new_mean' across selected_channels.
    Returns a single array (the mean across those channels), or None if empty.
    """
    means = []
    for ch in selected_channels:
        if ch in channel_means:
            _, new_m = channel_means[ch]
            means.append(new_m)
    if means:
        return np.mean(means, axis=0)
    else:
        # If nothing is selected, return None
        return None


def plot_individual_channels(
    ax,
    channels,
    channel_means_dict,
    freqs,
    colors_for_channels,
    show_original_mean=True,
    show_new_mean=True,
    title="",
    axis_fs=10,
    legend_fs=8,
    tick_fs=8
):
    """
    Plots each channel in 'channels' onto Axes 'ax',
    using a distinct color from 'colors_for_channels'.
    If show_original_mean=True, it plots original mean in color (solid).
    If show_new_mean=True, it plots outlier-excluded mean in color (dashed).
    """
    for ch in channels:
        if ch not in channel_means_dict:
            continue
        orig_mean, new_mean = channel_means_dict[ch]
        color = colors_for_channels.get(ch, "blue")

        lbl_orig = f"{ch} Orig"
        lbl_new  = f"{ch} New"

        if show_original_mean:
            ax.plot(freqs, orig_mean, color=color, linestyle='-', label=lbl_orig)
        if show_new_mean:
            ax.plot(freqs, new_mean, color=color, linestyle='--', label=lbl_new)
    
    ax.set_title(title, fontsize=axis_fs)
    ax.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
    ax.tick_params(axis='both', labelsize=tick_fs)
    ax.legend(fontsize=legend_fs, loc="upper right")


###############################################################################
# 3) Main Interactive GUI
###############################################################################
def build_exportable_plot_psd_gui():
    """
    An interactive GUI that:
      - Loads PSD data (block or non-block).
      - Merge blocks or keep them separate.
      - Plot PSDs (either block-by-block or merged).
      - "Compute & Plot Final Means" => a single figure with final group means 
        (depending on merged vs non-merged).
      - Export final means data as (Pickle, Excel) or final plot as (png, svg, jpeg, ppt).
    """
    # Data containers
    loaded_raw_data = {}   # the PSD dictionary exactly as loaded from pickle
    merged_psd = {}        # for merged data if user chooses
    blockwise_means_after = {}  # store (block -> channel_means_after) in non-merged mode
    block_map = {}         # store raw block-based data in non-merged mode
    channel_means_after_merged = {}  # store in merged mode
    do_merge = False       # track user choice
    final_data_dict = {}   # store final means for export
    final_fig = None       # store the figure for final means
    current_figures = []   # store subplots from the "screening" or block-based plots

    ############################################################################
    # 1. UI WIDGETS
    ############################################################################
    merge_blocks_cb = widgets.Checkbox(
        value=False, 
        description='Merge blocks?',
        tooltip='If checked, blockN:ChX entries are merged into ChX. Otherwise, each block is separate.'
    )

    load_psd_button = widgets.Button(
        description='Load PSD Pickle',
        button_style='info',
        tooltip='Load a pickled PSD file'
    )
    psd_file_chooser = FileChooser(os.getcwd(), title='Select PSD Pickle File', select_default=False)
    psd_file_chooser.show_only_files = True
    psd_file_chooser.filter_pattern = ['*.pkl']
    
    load_output_area = widgets.Output()
    
    # Channels dropdown (only used in merged mode)
    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description='Select Ch (merged):',
        layout=widgets.Layout(width='300px', height='200px')
    )

    show_eye1_cb = widgets.Checkbox(value=True, description='Show Eye1 (Ch1-8)')
    show_eye2_cb = widgets.Checkbox(value=True, description='Show Eye2 (Ch9-16)')
    show_original_mean_cb = widgets.Checkbox(value=True, description='Show Original Means')
    show_new_mean_cb = widgets.Checkbox(value=True, description='Show New Means')

    x_min_widget = widgets.FloatText(value=None, description='X-min:', layout=widgets.Layout(width='120px'))
    x_max_widget = widgets.FloatText(value=None, description='X-max:', layout=widgets.Layout(width='120px'))
    y_min_widget = widgets.FloatText(value=None, description='Y-min:', layout=widgets.Layout(width='120px'))
    y_max_widget = widgets.FloatText(value=None, description='Y-max:', layout=widgets.Layout(width='120px'))

    title_fs_widget = widgets.IntText(value=10, description='Title FS:', layout=widgets.Layout(width='100px'))
    axis_fs_widget  = widgets.IntText(value=8,  description='Axis FS:',  layout=widgets.Layout(width='100px'))
    legend_fs_widget= widgets.IntText(value=8, description='Legend FS:', layout=widgets.Layout(width='100px'))
    tick_fs_widget  = widgets.IntText(value=8,  description='Tick FS:',  layout=widgets.Layout(width='100px'))

    test_band_text = widgets.Textarea(
        value="(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)",
        description='Test Bands:',
        layout=widgets.Layout(width='250px', height='80px')
    )
    low_band_threshold_widget = widgets.FloatText(value=3.0, description='Low band thresh:', layout=widgets.Layout(width='150px'))
    test_band_threshold_widget = widgets.FloatText(value=10.0, description='Test band thresh:', layout=widgets.Layout(width='150px'))

    plot_psd_button = widgets.Button(
        description='Plot PSDs',
        button_style='success',
        tooltip='Plot the selected PSDs'
    )
    plot_output_area = widgets.Output()

    # Button to compute & plot final means
    compute_final_means_button = widgets.Button(
        description='Compute & Plot Final Means',
        button_style='primary',
        tooltip='Compute the group means across channels/blocks, plot them, and store in final_data_dict.'
    )
    final_means_output_area = widgets.Output()

    export_filename_widget = widgets.Text(value='MyPSDExport', description='File Base Name:', layout=widgets.Layout(width='300px'))
    export_format_widget = widgets.Dropdown(
        options=['png', 'svg', 'jpeg', 'ppt'], value='png', description='Format:',
        layout=widgets.Layout(width='150px')
    )
    export_button = widgets.Button(description='Export Subplots', button_style='warning')

    export_output_area = widgets.Output()

    final_export_basename_widget = widgets.Text(
        value='FinalMeansExport',
        description='Base Name:',
        layout=widgets.Layout(width='200px')
    )
    final_export_format_widget = widgets.Dropdown(
        options=['Pickle','Excel','png','svg','jpeg','ppt'],
        value='Pickle',
        description='Format:',
        layout=widgets.Layout(width='140px')
    )
    final_export_button = widgets.Button(description='Export Final Means', button_style='info')
    final_export_output = widgets.Output()

    ############################################################################
    # 2. LOAD LOGIC
    ############################################################################
    def on_load_psd_clicked(btn):
        with load_output_area:
            clear_output()
            psd_path = psd_file_chooser.selected
            if not psd_path:
                print("No file selected.")
                return
            if not os.path.isfile(psd_path):
                print(f"File does not exist: {psd_path}")
                return
            try:
                with open(psd_path, 'rb') as f:
                    data_in = pickle.load(f)
                loaded_raw_data.clear()
                loaded_raw_data.update(data_in)
                print(f"Loaded PSD data from {psd_path}")
                print("Next: choose 'Merge blocks?' or not, then click 'Plot PSDs'.")
            except Exception as e:
                print(f"Error loading PSD: {e}")

    load_psd_button.on_click(on_load_psd_clicked)

    def merge_block_keys(raw_dict):
        """
        Merge 'blockN:ChX' => 'ChX'.
        Returns dict {ChX: {psd: ..., freqs: ...}}
        """
        block_pattern = re.compile(r'^block(\d+):Ch(\d+)$')
        merged_dict = {}

        for key, val in raw_dict.items():
            m = block_pattern.match(key)
            if m:
                ch_num = m.group(2)
                ch_key = f"Ch{ch_num}"
                if ch_key not in merged_dict:
                    merged_dict[ch_key] = {'psd': [], 'freqs': None}
                merged_dict[ch_key]['psd'].append(val['psd'])
                if merged_dict[ch_key]['freqs'] is None:
                    merged_dict[ch_key]['freqs'] = val['freqs']
            else:
                # Already "ChX" or something
                if key not in merged_dict:
                    merged_dict[key] = {'psd': [], 'freqs': None}
                merged_dict[key]['psd'].append(val['psd'])
                if merged_dict[key]['freqs'] is None:
                    merged_dict[key]['freqs'] = val['freqs']

        # Concatenate
        for ch_key, stuff in merged_dict.items():
            if len(stuff['psd']) > 1:
                merged_dict[ch_key]['psd'] = np.concatenate(stuff['psd'], axis=0)
            else:
                merged_dict[ch_key]['psd'] = stuff['psd'][0]

        return merged_dict

    ############################################################################
    # 3. PLOT LOGIC
    ############################################################################
    def on_plot_psd_clicked(btn):
        nonlocal do_merge
        with plot_output_area:
            clear_output()
            current_figures.clear()

            if not loaded_raw_data:
                print("No data loaded. Please load PSD first.")
                return

            do_merge = merge_blocks_cb.value
            print(f"Merge blocks? {do_merge}")

            # Parse test bands
            test_bands_str = test_band_text.value.strip().replace(" ", "")
            test_bands_list = []
            if test_bands_str:
                pairs = test_bands_str.split(")")
                for p in pairs:
                    p = p.strip(",").strip("(").strip()
                    if not p:
                        continue
                    vals = p.split(",")
                    if len(vals) == 2:
                        try:
                            lowf = float(vals[0])
                            highf = float(vals[1])
                            test_bands_list.append((lowf, highf))
                        except:
                            pass
                if not test_bands_list:
                    test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]
            else:
                test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]

            low_band = (1,3)
            low_band_threshold = low_band_threshold_widget.value
            test_band_threshold = test_band_threshold_widget.value

            # Axis / font
            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value
            title_fs = title_fs_widget.value
            axis_fs  = axis_fs_widget.value
            legend_fs= legend_fs_widget.value
            tick_fs  = tick_fs_widget.value

            # If merging
            if do_merge:
                merged_psd.clear()
                merged_psd.update(merge_block_keys(loaded_raw_data))

                # Sort channels => fill dropdown
                def chnum(k):
                    return int(k.replace("Ch","")) if k.startswith("Ch") else 9999
                sorted_ch = sorted(merged_psd.keys(), key=chnum)
                channels_dropdown.options = sorted_ch

                print("Blocks merged into single dict. Now select channels above (if not already). Then re-click 'Plot PSDs' to see your selected subset.")
                return

            else:
                # Non-merged approach => block-based subplots
                # We'll do the big figure with #blocks rows x 2 columns
                block_pattern = re.compile(r'^block(\d+):Ch(\d+)$')
                block_map.clear()
                # Fill block_map = {blockNum: { key-> val, key->val}}
                for key, val in loaded_raw_data.items():
                    m = block_pattern.match(key)
                    if m:
                        bnum = int(m.group(1))
                        if bnum not in block_map:
                            block_map[bnum] = {}
                        block_map[bnum][key] = val
                    else:
                        # treat as block0 if no block
                        if 0 not in block_map:
                            block_map[0] = {}
                        block_map[0][key] = val

                all_blocks = sorted(block_map.keys())
                if not all_blocks:
                    print("No 'blockN:ChX' keys found. Possibly your data has only 'ChX' keys?")
                    return

                # We'll create subplots => rows=number_of_blocks, cols=2
                n_blocks = len(all_blocks)
                fig, axes = plt.subplots(n_blocks, 2, figsize=(12, 4*n_blocks), squeeze=False)

                import matplotlib.cm as cm
                color_cycle = cm.get_cmap('tab20').colors

                blockwise_means_after.clear()

                row_idx = 0
                for bnum in all_blocks:
                    # Build local dict for that block
                    block_dict = {}
                    for k, v in block_map[bnum].items():
                        block_dict[k] = v
                    # compute means
                    block_means_before = compute_channel_means(
                        block_dict,
                        exclude=False
                    )
                    block_means_after_ = compute_channel_means(
                        block_dict,
                        exclude=True,
                        low_band=low_band,
                        low_band_threshold=low_band_threshold,
                        test_bands=test_bands_list,
                        test_band_threshold=test_band_threshold
                    )
                    # Store for later final means usage
                    blockwise_means_after[bnum] = block_means_after_

                    # separate eye1 vs eye2
                    def parse_block_ch(key):
                        mm = block_pattern.match(key)
                        if mm:
                            return int(mm.group(2))
                        return None

                    eye1_keys = []
                    eye2_keys = []
                    for kk in block_dict.keys():
                        chn = parse_block_ch(kk)
                        if chn is not None:
                            if 1 <= chn <= 8:
                                eye1_keys.append(kk)
                            elif 9 <= chn <=16:
                                eye2_keys.append(kk)

                    # color mapping
                    allkeys = sorted(block_dict.keys())
                    cdict = {}
                    for i, ckey in enumerate(allkeys):
                        cdict[ckey] = color_cycle[i % len(color_cycle)]

                    ax_eye1 = axes[row_idx, 0]
                    if show_eye1_cb.value and len(eye1_keys) > 0:
                        freqs_e1 = block_dict[eye1_keys[0]]['freqs']
                        plot_individual_channels(
                            ax=ax_eye1,
                            channels=eye1_keys,
                            channel_means_dict=block_means_before,
                            freqs=freqs_e1,
                            colors_for_channels=cdict,
                            show_original_mean=show_original_mean_cb.value,
                            show_new_mean=show_new_mean_cb.value,
                            title=f"Block {bnum}: Eye1",
                            axis_fs=axis_fs,
                            legend_fs=legend_fs,
                            tick_fs=tick_fs
                        )
                        gmean_e1 = compute_group_mean(block_means_after_, eye1_keys)
                        if gmean_e1 is not None:
                            ax_eye1.plot(freqs_e1, gmean_e1, color='black', linewidth=2, label='GroupMean')
                            ax_eye1.legend(fontsize=legend_fs)

                        if x_min is not None: ax_eye1.set_xlim(left=x_min)
                        if x_max is not None: ax_eye1.set_xlim(right=x_max)
                        if y_min is not None: ax_eye1.set_ylim(bottom=y_min)
                        if y_max is not None: ax_eye1.set_ylim(top=y_max)
                    else:
                        ax_eye1.set_title(f"Block {bnum}: Eye1 (no data or hidden)")

                    ax_eye2 = axes[row_idx, 1]
                    if show_eye2_cb.value and len(eye2_keys) > 0:
                        freqs_e2 = block_dict[eye2_keys[0]]['freqs']
                        plot_individual_channels(
                            ax=ax_eye2,
                            channels=eye2_keys,
                            channel_means_dict=block_means_before,
                            freqs=freqs_e2,
                            colors_for_channels=cdict,
                            show_original_mean=show_original_mean_cb.value,
                            show_new_mean=show_new_mean_cb.value,
                            title=f"Block {bnum}: Eye2",
                            axis_fs=axis_fs,
                            legend_fs=legend_fs,
                            tick_fs=tick_fs
                        )
                        gmean_e2 = compute_group_mean(block_means_after_, eye2_keys)
                        if gmean_e2 is not None:
                            ax_eye2.plot(freqs_e2, gmean_e2, color='black', linewidth=2, label='GroupMean')
                            ax_eye2.legend(fontsize=legend_fs)

                        if x_min is not None: ax_eye2.set_xlim(left=x_min)
                        if x_max is not None: ax_eye2.set_xlim(right=x_max)
                        if y_min is not None: ax_eye2.set_ylim(bottom=y_min)
                        if y_max is not None: ax_eye2.set_ylim(top=y_max)
                    else:
                        ax_eye2.set_title(f"Block {bnum}: Eye2 (no data or hidden)")

                    row_idx += 1

                plt.tight_layout()
                plt.show()
                current_figures.append(fig)

                print(f"Plotted {n_blocks} block(s). You can now click 'Compute & Plot Final Means' if desired for an overall block-based final figure.")

    ############################################################################
    # 4. COMPUTE & PLOT FINAL MEANS BUTTON
    ############################################################################
    def on_compute_final_means_clicked(btn):
        nonlocal final_data_dict, final_fig
        with final_means_output_area:
            clear_output()
            final_data_dict.clear()
            final_fig = None

            if not loaded_raw_data:
                print("No data loaded. Please load PSD first.")
                return

            if do_merge:
                # We have merged_psd and channel_means_after_merged
                # Let's compute means after exclusion for entire merged set
                if not merged_psd:
                    print("Blocks not merged yet or no data. Please do 'Plot PSDs' after checking 'Merge blocks?'.")
                    return

                # figure out user-selected channels (1..16)
                selected_channels = list(channels_dropdown.value)
                if not selected_channels:
                    print("No channels selected in merged mode.")
                    return

                # Recompute "after" means with user thresholds
                test_bands_str = test_band_text.value.strip().replace(" ", "")
                test_bands_list = []
                if test_bands_str:
                    pairs = test_bands_str.split(")")
                    for p in pairs:
                        p = p.strip(",").strip("(").strip()
                        if not p:
                            continue
                        vals = p.split(",")
                        if len(vals) == 2:
                            try:
                                lowf = float(vals[0])
                                highf = float(vals[1])
                                test_bands_list.append((lowf, highf))
                            except:
                                pass
                    if not test_bands_list:
                        test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]
                else:
                    test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]

                low_band=(1,3)
                low_thresh=low_band_threshold_widget.value
                test_thresh=test_band_threshold_widget.value

                channel_means_after_merged.clear()
                cma = compute_channel_means(
                    merged_psd,
                    exclude=True,
                    low_band=low_band,
                    low_band_threshold=low_thresh,
                    test_bands=test_bands_list,
                    test_band_threshold=test_thresh
                )
                channel_means_after_merged.update(cma)

                # Eye1 = ch1..ch8, Eye2 = ch9..ch16
                def chnum(k):
                    return int(k.replace("Ch","")) if k.startswith("Ch") else None

                eye1_sel = [ch for ch in selected_channels if chnum(ch) and 1 <= chnum(ch) <= 8]
                eye2_sel = [ch for ch in selected_channels if chnum(ch) and 9 <= chnum(ch) <=16]

                # We compute final group means
                gm_e1 = compute_group_mean(channel_means_after_merged, eye1_sel) if eye1_sel else None
                gm_e2 = compute_group_mean(channel_means_after_merged, eye2_sel) if eye2_sel else None

                # Store in final_data_dict for export
                # We'll assume same freqs for all channels
                freq_ref = None
                if eye1_sel:
                    freq_ref = merged_psd[eye1_sel[0]]['freqs']
                elif eye2_sel:
                    freq_ref = merged_psd[eye2_sel[0]]['freqs']

                if freq_ref is not None:
                    final_data_dict["freqs"] = freq_ref

                if gm_e1 is not None:
                    final_data_dict["Eye1_final"] = gm_e1
                if gm_e2 is not None:
                    final_data_dict["Eye2_final"] = gm_e2

                # Plot them on single figure with red/blue lines
                fig, ax = plt.subplots(figsize=(6,4))
                if gm_e1 is not None:
                    ax.plot(freq_ref, gm_e1, color='red', label='Eye1 final mean')
                if gm_e2 is not None:
                    ax.plot(freq_ref, gm_e2, color='blue', label='Eye2 final mean')
                ax.set_xlabel("Frequency (Hz)")
                ax.set_ylabel("PSD (V²/Hz)")
                ax.set_title("Merged Final Means")
                ax.legend()
                plt.show()

                final_fig = fig
                print("Computed & plotted final means in merged mode. You can now Export Final Means below.")

            else:
                # Non-merged => we have blockwise_means_after for each block
                if not block_map:
                    print("No block data found. Please do 'Plot PSDs' with blocks unmerged first.")
                    return
                if not blockwise_means_after:
                    print("No blockwise means stored. Please 'Plot PSDs' first in non-merged mode.")
                    return

                # We'll gather group means for each block x Eye1, Eye2
                # Then plot them all in a single figure with 1 row, 2 columns
                # Eye1 on left => each block is a line
                # Eye2 on right => each block is a line
                fig, axes = plt.subplots(1, 2, figsize=(10,5))
                ax1, ax2 = axes

                # We'll assume freq arrays are consistent across blocks
                # We'll store them in final_data_dict as well
                final_data_dict.clear()

                block_nums = sorted(block_map.keys())
                block_colors = [
                    'red','blue','green','orange','purple','magenta','cyan','brown','lime','navy'
                ]
                c_idx=0

                for bnum in block_nums:
                    ch_means_after = blockwise_means_after[bnum]

                    # separate Eye1 vs Eye2
                    # We'll quickly gather which channels are Eye1/Eye2
                    # from the keys of block_map[bnum]
                    block_dict = block_map[bnum]
                    def parse_block_ch(key):
                        mm = re.match(r'^block(\d+):Ch(\d+)$', key)
                        if mm:
                            return int(mm.group(2))
                        return None
                    eye1_keys = []
                    eye2_keys = []
                    for k in block_dict.keys():
                        chn = parse_block_ch(k)
                        if chn is not None:
                            if 1 <= chn <= 8:
                                eye1_keys.append(k)
                            elif 9 <= chn <= 16:
                                eye2_keys.append(k)

                    gm_e1 = compute_group_mean(ch_means_after, eye1_keys) if eye1_keys else None
                    gm_e2 = compute_group_mean(ch_means_after, eye2_keys) if eye2_keys else None

                    # We'll store them in final_data_dict with keys like "Block1_Eye1"
                    if eye1_keys and gm_e1 is not None:
                        final_data_dict[f"Block{bnum}_Eye1"] = gm_e1
                        # freq ref
                        final_data_dict[f"Block{bnum}_Eye1_freqs"] = block_dict[eye1_keys[0]]['freqs']

                    if eye2_keys and gm_e2 is not None:
                        final_data_dict[f"Block{bnum}_Eye2"] = gm_e2
                        final_data_dict[f"Block{bnum}_Eye2_freqs"] = block_dict[eye2_keys[0]]['freqs']

                    # Plot them
                    color = block_colors[c_idx % len(block_colors)]
                    c_idx+=1

                    if gm_e1 is not None:
                        freqs_e1 = block_dict[eye1_keys[0]]['freqs']
                        ax1.plot(freqs_e1, gm_e1, color=color, label=f"Block{bnum}")
                    if gm_e2 is not None:
                        freqs_e2 = block_dict[eye2_keys[0]]['freqs']
                        ax2.plot(freqs_e2, gm_e2, color=color, label=f"Block{bnum}")

                ax1.set_title("Final Means (Eye1)")
                ax1.set_xlabel("Frequency (Hz)")
                ax1.set_ylabel("PSD (V²/Hz)")
                ax1.legend()
                ax2.set_title("Final Means (Eye2)")
                ax2.set_xlabel("Frequency (Hz)")
                ax2.set_ylabel("PSD (V²/Hz)")
                ax2.legend()

                plt.tight_layout()
                plt.show()
                final_fig = fig
                print("Computed & plotted final means in block-based mode. Stored in final_data_dict as 'BlockX_EyeY' keys.")

    ############################################################################
    # 5. EXPORT PLOTS (SUBPLOTS) - existing
    ############################################################################
    def on_export_button_clicked(b):
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No subplot figures to export. Please plot PSDs first.")
                return

            base_name = export_filename_widget.value.strip()
            fmt = export_format_widget.value.lower()
            if not base_name:
                print("Provide a valid filename base.")
                return
            if fmt == 'ppt' and not HAS_PPTX:
                print("python-pptx not installed, can't export PPT.")
                return

            print(f"Exporting {len(current_figures)} subplot figure(s) as {fmt}...")

            if fmt == 'ppt':
                prs = Presentation()
                for idx, fig in enumerate(current_figures, start=1):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tmpname = f"{base_name}_{idx}.png"
                    fig.savefig(tmpname, dpi=150)
                    left = top = Inches(1)
                    slide.shapes.add_picture(tmpname, left, top, Inches(8), Inches(4.5))
                    os.remove(tmpname)
                prs.save(f"{base_name}.pptx")
                print(f"Saved => {base_name}.pptx")
            else:
                for idx, fig in enumerate(current_figures, start=1):
                    outname = f"{base_name}_{idx}.{fmt}"
                    fig.savefig(outname, format=fmt, dpi=150)
                    print(f"Saved => {outname}")
            print("Done exporting.")

    ############################################################################
    # 6. EXPORT FINAL MEANS (DATA OR FIGURE)
    ############################################################################
    def on_final_export_clicked(b):
        nonlocal final_data_dict, final_fig
        with final_export_output:
            clear_output()
            base_name = final_export_basename_widget.value.strip()
            chosen_fmt = final_export_format_widget.value
            if not base_name:
                print("Please provide a valid base filename.")
                return

            if not final_data_dict and final_fig is None:
                print("No final data or final figure found. Please click 'Compute & Plot Final Means' first.")
                return

            print(f"Exporting final means as {chosen_fmt}...")

            if chosen_fmt == 'Pickle':
                outpkl = f"{base_name}_final_means.pkl"
                try:
                    with open(outpkl, 'wb') as f:
                        pickle.dump(final_data_dict, f)
                    print(f"Saved => {outpkl}")
                except Exception as e:
                    print(f"Error saving pickle: {e}")

            elif chosen_fmt == 'Excel':
                import pandas as pd
                outxlsx = f"{base_name}_final_means.xlsx"
                try:
                    df = pd.DataFrame()
                    # Each key in final_data_dict is e.g. "Eye1_final", "freqs", or "Block1_Eye1"
                    for k,v in final_data_dict.items():
                        arr = np.array(v)
                        df[k] = arr.ravel()  # flatten 1D if needed
                    df.to_excel(outxlsx, index=False)
                    print(f"Saved => {outxlsx}")
                except Exception as e:
                    print(f"Error saving Excel: {e}")

            elif chosen_fmt in ['png','svg','jpeg']:
                if final_fig is None:
                    print("No final figure to export.")
                    return
                outimg = f"{base_name}_final_means.{chosen_fmt}"
                try:
                    final_fig.savefig(outimg, format=chosen_fmt, dpi=150)
                    print(f"Saved => {outimg}")
                except Exception as e:
                    print(f"Error saving final fig: {e}")

            elif chosen_fmt == 'ppt':
                if not HAS_PPTX:
                    print("python-pptx not installed, cannot export PPT.")
                    return
                if final_fig is None:
                    print("No final figure to export.")
                    return
                ppt_name = f"{base_name}_final_means.pptx"
                try:
                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tmp_png = f"{base_name}_temp.png"
                    final_fig.savefig(tmp_png, dpi=150)
                    left = top = Inches(1)
                    slide.shapes.add_picture(tmp_png, left, top, Inches(8), Inches(4.5))
                    os.remove(tmp_png)
                    prs.save(ppt_name)
                    print(f"Exported => {ppt_name}")
                except Exception as e:
                    print(f"Error exporting PPT: {e}")

            print("Done exporting final means.")

    ############################################################################
    # HOOK UP THE BUTTONS
    ############################################################################
    load_psd_button.on_click(on_load_psd_clicked)
    plot_psd_button.on_click(on_plot_psd_clicked)
    compute_final_means_button.on_click(on_compute_final_means_clicked)
    export_button.on_click(on_export_button_clicked)
    final_export_button.on_click(on_final_export_clicked)

    ############################################################################
    # LAYOUT
    ############################################################################
    load_controls = widgets.HBox([merge_blocks_cb, load_psd_button, psd_file_chooser])
    load_section = widgets.VBox([
        widgets.HTML("<h3>Load PSD Data</h3>"),
        load_controls,
        load_output_area
    ])

    plot_controls_top = widgets.HBox([show_eye1_cb, show_eye2_cb, show_original_mean_cb, show_new_mean_cb])
    thresholds_box = widgets.HBox([low_band_threshold_widget, test_band_threshold_widget])
    axis_box = widgets.HBox([x_min_widget, x_max_widget, y_min_widget, y_max_widget])
    font_box = widgets.HBox([title_fs_widget, axis_fs_widget, legend_fs_widget, tick_fs_widget])
    
    plot_section = widgets.VBox([
        widgets.HTML("<h3>Plot Options</h3>"),
        plot_controls_top,
        widgets.Label("Test Bands: (e.g. (7,9),(9,11),...)"),
        test_band_text,
        thresholds_box,
        axis_box,
        font_box,
        widgets.HTML("<b>(Only relevant if merging blocks):</b>"),
        channels_dropdown,
        plot_psd_button,
        plot_output_area
    ])

    final_means_section = widgets.VBox([
        widgets.HTML("<h3>Compute & Plot Final Means</h3>"),
        compute_final_means_button,
        final_means_output_area
    ])

    export_subplots_section = widgets.VBox([
        widgets.HTML("<h3>Export Subplots (screening/block plots)</h3>"),
        widgets.HBox([export_filename_widget, export_format_widget, export_button]),
        export_output_area
    ])

    export_final_section = widgets.VBox([
        widgets.HTML("<h3>Export Final Means Plot/Data</h3>"),
        widgets.HBox([final_export_basename_widget, final_export_format_widget, final_export_button]),
        final_export_output
    ])

    ui = widgets.VBox([
        widgets.HTML("<h2>Block vs Merged PSD GUI</h2>"),
        load_section,
        plot_section,
        final_means_section,
        export_subplots_section,
        export_final_section
    ])

    display(ui)


# Build and display the GUI
build_exportable_plot_psd_gui()
