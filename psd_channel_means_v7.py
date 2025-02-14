import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math
import re

import ipywidgets as widgets
from ipywidgets import VBox, HBox
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. PPT export will be disabled.")

# Check for pandas (for Excel export)
try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    print("Warning: pandas not installed. Excel export will be disabled.")

###############################################################################
# 1) Helper Functions
###############################################################################

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    if test_band_count_threshold is None:
        test_band_count_threshold = len(test_bands) // 2

    mean_psd = np.mean(psd_array, axis=0)
    excluded_traces = []
    kept_traces = []

    low_band_indices = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0]
        for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        if np.any(trace[low_band_indices] > low_band_threshold * mean_psd[low_band_indices]):
            excluded_traces.append(i)
            continue

        suprathreshold_count = 0
        for indices in band_indices:
            if np.any(trace[indices] > test_band_threshold * mean_psd[indices]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces

def compute_channel_means(
    psds_dict,
    exclude=False,
    low_band=(1,3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    channel_means = {}
    for channel, data in psds_dict.items():
        psd = data.get('psd', None)
        freqs = data.get('freqs', None)
        if psd is None or freqs is None:
            print(f"Channel '{channel}' is missing 'psd' or 'freqs'. Skipping.")
            continue

        original_mean = np.mean(psd, axis=0)
        if exclude:
            kept_traces, _ = exclude_traces(
                psd_array=psd,
                freqs=freqs,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold
            )
            if kept_traces:
                new_mean = np.mean(psd[kept_traces], axis=0)
            else:
                new_mean = np.zeros_like(original_mean)
        else:
            new_mean = original_mean.copy()
        
        channel_means[channel] = (original_mean, new_mean)
    
    return channel_means

def compute_combined_means(channel_means, blocks, eye_channels, deselected_channels):
    combined_means = {ch: None for ch in eye_channels}
    for ch in eye_channels:
        for block in blocks:
            channel = f"{block}:{ch}"
            if channel in channel_means and channel not in deselected_channels:
                _, new_mean = channel_means[channel]
                if combined_means[ch] is None:
                    combined_means[ch] = new_mean.copy()
                else:
                    combined_means[ch] += new_mean
        if combined_means[ch] is not None:
            combined_means[ch] /= len(blocks)
    return combined_means

def plot_aggregated_means(
    combined_means,
    freqs,
    eye,
    title_fs=10,
    axis_fs=8,
    legend_fs=8,
    y_max=80,
    channel_colors=None,
    show_channel_means=True,
    show_mean_of_means=True
):
    fig, ax = plt.subplots(figsize=(12, 6))

    mean_of_means = None
    count = 0

    for ch, mean in combined_means.items():
        if mean is None:
            continue
        color = channel_colors.get(ch, "black")
        if show_channel_means:
            ax.plot(freqs, mean, linestyle='--', color=color, label=f"{eye}:{ch} Mean")
        if mean_of_means is None:
            mean_of_means = mean.copy()
        else:
            mean_of_means += mean
        count += 1

    if mean_of_means is not None and count > 0:
        mean_of_means /= count
        if show_mean_of_means:
            ax.plot(freqs, mean_of_means, linestyle='-', color='black', linewidth=2, label="Mean of Means")

    ax.set_title(f"{eye} Aggregated Means", fontsize=title_fs)
    ax.set_xlim([0, 45])
    ax.set_ylim([0, y_max])
    ax.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
    ax.legend(fontsize=legend_fs)
    plt.tight_layout()
    display(fig)  # Use display instead of plt.show() to ensure the figure is shown in the output_area
    return fig  # Return the figure for potential exporting

def plot_block_layout(
    psds_dict,
    channel_means,
    freqs,
    deselected_channels,
    show_original=True,
    show_new=True,
    title_fs=10,
    axis_fs=8,
    legend_fs=8,
    plot_width=12,
    plot_height=18,
    y_max=80,
    channel_colors=None,
    show_combined_means=False,
    show_aggregated_plots=False
):
    if channel_colors is None:
        color_cycle = plt.cm.get_cmap('tab20').colors
        all_channels = sorted({key.split(':Ch')[-1] for key in psds_dict.keys()})
        channel_colors = {f"Ch{ch}": color_cycle[i % len(color_cycle)] for i, ch in enumerate(all_channels)}

    blocks = sorted(set([key.split(':')[0] for key in psds_dict.keys()]))
    fig, axes = plt.subplots(len(blocks), 2, figsize=(plot_width, plot_height))

    if len(blocks) == 1:
        axes = [axes]

    combined_eye1 = compute_combined_means(channel_means, blocks, [f"Ch{i}" for i in range(1, 9)], deselected_channels)
    combined_eye2 = compute_combined_means(channel_means, blocks, [f"Ch{i}" for i in range(9, 17)], deselected_channels)

    for i, block in enumerate(blocks):
        eye1_channels = [key for key in psds_dict.keys() if key.startswith(block) and int(key.split(':Ch')[-1]) <= 8]
        eye2_channels = [key for key in psds_dict.keys() if key.startswith(block) and int(key.split(':Ch')[-1]) > 8]

        # Eye 1 (Ch1-8)
        ax_eye1 = axes[i][0] if len(blocks) > 1 else axes[0]
        for channel in eye1_channels:
            if channel in deselected_channels:
                continue
            original_mean, new_mean = channel_means[channel]
            ch_number = f"Ch{channel.split(':Ch')[-1]}"
            color = channel_colors.get(ch_number, "black")
            if show_original:
                ax_eye1.plot(freqs, original_mean, linestyle='-', color=color, label=f"{channel} Original")
            if show_new:
                ax_eye1.plot(freqs, new_mean, linestyle='--', color=color, label=f"{channel} New")
        # Removed incorrect combined_mean plotting per block
        ax_eye1.set_title(f"{block} Eye 1 (Ch1-8)", fontsize=title_fs)
        ax_eye1.set_xlim([0, 45])
        ax_eye1.set_ylim([0, y_max])
        ax_eye1.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
        ax_eye1.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
        ax_eye1.legend(fontsize=legend_fs)

        # Eye 2 (Ch9-16)
        ax_eye2 = axes[i][1] if len(blocks) > 1 else axes[1]
        for channel in eye2_channels:
            if channel in deselected_channels:
                continue
            original_mean, new_mean = channel_means[channel]
            ch_number = f"Ch{channel.split(':Ch')[-1]}"
            color = channel_colors.get(ch_number, "black")
            if show_original:
                ax_eye2.plot(freqs, original_mean, linestyle='-', color=color, label=f"{channel} Original")
            if show_new:
                ax_eye2.plot(freqs, new_mean, linestyle='--', color=color, label=f"{channel} New")
        # Removed incorrect combined_mean plotting per block
        ax_eye2.set_title(f"{block} Eye 2 (Ch9-16)", fontsize=title_fs)
        ax_eye2.set_xlim([0, 45])
        ax_eye2.set_ylim([0, y_max])
        ax_eye2.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
        ax_eye2.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
        ax_eye2.legend(fontsize=legend_fs)

    plt.tight_layout()
    display(fig)  # Use display instead of plt.show() to ensure the figure is shown in the output_area

    aggregated_fig1 = None
    aggregated_fig2 = None
    if show_aggregated_plots:
        aggregated_fig1 = plot_aggregated_means(
            combined_eye1, freqs, "Eye 1",
            title_fs=title_fs, axis_fs=axis_fs, legend_fs=legend_fs, y_max=y_max,
            channel_colors=channel_colors
        )
        aggregated_fig2 = plot_aggregated_means(
            combined_eye2, freqs, "Eye 2",
            title_fs=title_fs, axis_fs=axis_fs, legend_fs=legend_fs, y_max=y_max,
            channel_colors=channel_colors
        )
        # Return main figure and aggregated figures
        return fig, aggregated_fig1, aggregated_fig2
    else:
        return fig  # Return only the main figure for exporting

def export_figure(fig, export_path, export_format="png"):
    if export_format.lower() in ["png", "jpeg", "svg", "pdf"]:
        fig.savefig(f"{export_path}.{export_format}", format=export_format, dpi=300)
        print(f"Figure exported as {export_path}.{export_format}")
    elif export_format.lower() == "ppt" and HAS_PPTX:
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # Blank slide
        temp_img = f"{export_path}.png"
        fig.savefig(temp_img, format="png", dpi=300)
        slide.shapes.add_picture(temp_img, Inches(1), Inches(1), Inches(8), Inches(5))
        ppt.save(f"{export_path}.pptx")
        os.remove(temp_img)
        print(f"Figure exported as {export_path}.pptx")
    else:
        print(f"Unsupported format or missing dependencies for {export_format}.")

###############################################################################
# 2) Main GUI
###############################################################################

def build_gui():
    psds_dict = {}
    channel_means = {}
    freqs = None
    deselected_channels = []

    file_chooser = FileChooser(os.getcwd())
    load_button = widgets.Button(description="Load PSD", button_style='info')
    plot_button = widgets.Button(description="Plot", button_style='success')
    export_button = widgets.Button(description="Export Figure", button_style='warning')
    output_area = widgets.Output()

    channel_toggles = widgets.VBox()
    plot_width = widgets.FloatText(value=12, description="Plot Width")
    plot_height = widgets.FloatText(value=18, description="Plot Height")
    y_max_input = widgets.FloatText(value=80, description="Y Max")
    export_format = widgets.Dropdown(options=["png", "jpeg", "svg", "pdf", "ppt"], value="png", description="Export Format")
    export_name = widgets.Text(value="plot_export", description="File Name")
    show_original = widgets.Checkbox(value=True, description="Show Original Traces")
    show_new = widgets.Checkbox(value=True, description="Show New Traces")
    show_combined = widgets.Checkbox(value=False, description="Show Combined Means")
    show_aggregated = widgets.Checkbox(value=False, description="Show Aggregated Plots")

    # Last plotted figures to be used in export (list to store multiple figures)
    last_plotted_figures = []

    def update_channel_toggles():
        nonlocal deselected_channels
        deselected_channels.clear()
        toggles = []
        for channel in sorted(psds_dict.keys()):
            toggle = widgets.Checkbox(value=True, description=channel)
            def toggle_change(change, channel=channel):
                if not change['new']:
                    if channel not in deselected_channels:
                        deselected_channels.append(channel)
                else:
                    if channel in deselected_channels:
                        deselected_channels.remove(channel)
            toggle.observe(toggle_change, names='value')
            toggles.append(toggle)
        channel_toggles.children = toggles

    def load_psd_callback(b):
        nonlocal psds_dict, freqs, channel_means
        with output_area:
            clear_output()
            file_path = file_chooser.selected
            if not file_path or not os.path.exists(file_path):
                print("Please select a valid file.")
                return
            try:
                with open(file_path, 'rb') as f:
                    psds_dict = pickle.load(f)
                    if not psds_dict:
                        print("Loaded PSD data is empty.")
                        return
                    freqs = next(iter(psds_dict.values()))['freqs']
                    channel_means = compute_channel_means(psds_dict, exclude=True)
                    print("Loaded PSD data.")
                    update_channel_toggles()
            except Exception as e:
                print(f"Failed to load file: {e}")

    def plot_callback(b):
        nonlocal last_plotted_figures
        with output_area:
            clear_output()
            if not psds_dict:
                print("No PSD data loaded.")
                return
            try:
                # Clear previous figures
                last_plotted_figures.clear()
                # Plot and retrieve figures
                if show_aggregated.value:
                    # plot_block_layout returns main_fig, agg_fig1, agg_fig2
                    main_fig, agg_fig1, agg_fig2 = plot_block_layout(
                        psds_dict,
                        channel_means,
                        freqs,
                        deselected_channels,
                        show_original=show_original.value,
                        show_new=show_new.value,
                        plot_width=plot_width.value,
                        plot_height=plot_height.value,
                        y_max=y_max_input.value,
                        show_combined_means=show_combined.value,
                        show_aggregated_plots=show_aggregated.value
                    )
                    # Store all figures for export
                    last_plotted_figures.extend([main_fig, agg_fig1, agg_fig2])
                    print("Plots generated successfully.")
                else:
                    # plot_block_layout returns only main_fig
                    main_fig = plot_block_layout(
                        psds_dict,
                        channel_means,
                        freqs,
                        deselected_channels,
                        show_original=show_original.value,
                        show_new=show_new.value,
                        plot_width=plot_width.value,
                        plot_height=plot_height.value,
                        y_max=y_max_input.value,
                        show_combined_means=show_combined.value,
                        show_aggregated_plots=show_aggregated.value
                    )
                    # Store main figure for export
                    last_plotted_figures.append(main_fig)
                    print("Plot generated successfully.")
            except Exception as e:
                print(f"Failed to plot data: {e}")

    def export_callback(b):
        with output_area:
            # Avoid clearing output to keep export messages visible
            # clear_output()
            if not last_plotted_figures:
                print("No figures available to export. Please plot first.")
                return
            try:
                export_dir = os.getcwd()  # You can modify this to allow user to choose export directory
                for idx, fig in enumerate(last_plotted_figures):
                    # Generate unique export path for each figure
                    if len(last_plotted_figures) == 1:
                        filename = f"{export_name.value}"
                    else:
                        filename = f"{export_name.value}_{idx+1}"
                    export_path = os.path.join(export_dir, filename)
                    export_format_selected = export_format.value

                    if export_format_selected.lower() in ["png", "jpeg", "svg", "pdf"]:
                        fig.savefig(f"{export_path}.{export_format_selected}", format=export_format_selected, dpi=300)
                        print(f"Figure {idx+1} exported as {filename}.{export_format_selected}")
                    elif export_format_selected.lower() == "ppt" and HAS_PPTX:
                        if idx == 0:
                            ppt = Presentation()
                        slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # Blank slide
                        temp_img = f"{export_path}.png"
                        fig.savefig(temp_img, format="png", dpi=300)
                        slide.shapes.add_picture(temp_img, Inches(1), Inches(1), Inches(8), Inches(5))
                        os.remove(temp_img)
                        if idx == len(last_plotted_figures) - 1:
                            ppt.save(f"{export_path}.pptx")
                            print(f"All figures exported as {filename}.pptx")
                    else:
                        print(f"Unsupported format or missing dependencies for {export_format_selected}.")
            except Exception as e:
                print(f"Failed to export figures: {e}")

    load_button.on_click(load_psd_callback)
    plot_button.on_click(plot_callback)
    export_button.on_click(export_callback)

    display(
        VBox([
            file_chooser,
            HBox([load_button, plot_button, export_button]),
            HBox([plot_width, plot_height, y_max_input]),
            HBox([export_name, export_format]),
            HBox([show_original, show_new, show_combined, show_aggregated]),
            channel_toggles,
            output_area
        ])
    )

build_gui()
