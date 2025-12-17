import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math
import re

import ipywidgets as widgets
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed, PPT export will be disabled.")

import mne
from mne.utils import set_log_level
from tqdm.notebook import tqdm
from scipy import signal

# -----------------------------------------------------------------
# Suppress less-serious MNE messages
set_log_level("ERROR")  # "WARNING", "ERROR", or "CRITICAL" as desired
# -----------------------------------------------------------------

##############################################################################
# 1) HELPER FUNCTIONS
##############################################################################


def _welch_longest_clean_segment_1d(x, sfreq, nperseg, noverlap, window):
    """
    NaN-aware Welch on 1D array x using only the longest contiguous run
    of finite samples >= nperseg. Returns (freqs, psd) or (None, None).
    """
    x = np.asarray(x, dtype=float)
    mask = np.isfinite(x)
    if mask.sum() < nperseg:
        return None, None

    finite_idx = np.where(mask)[0]
    starts = [finite_idx[0]]
    ends = []

    for i in range(1, len(finite_idx)):
        if finite_idx[i] != finite_idx[i-1] + 1:
            ends.append(finite_idx[i-1])
            starts.append(finite_idx[i])
    ends.append(finite_idx[-1])

    lengths = [e - s + 1 for s, e in zip(starts, ends)]
    best_idx = int(np.argmax(lengths))
    s = starts[best_idx]
    e = ends[best_idx]
    seg = x[s:e+1]

    if seg.size < nperseg:
        return None, None

    freqs, psd = signal.welch(
        seg,
        fs=sfreq,
        window=window,
        nperseg=nperseg,
        noverlap=noverlap,
        scaling="density",
        axis=0,
    )
    return freqs, psd


def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    if test_band_count_threshold is None:
        test_band_count_threshold = len(test_bands) // 2

    mean_psd = np.nanmean(psd_array, axis=0)
    excluded_traces = []
    kept_traces = []

    low_band_indices = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0] for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # 1) Low-band outlier
        if np.any(trace[low_band_indices] > low_band_threshold * mean_psd[low_band_indices]):
            excluded_traces.append(i)
            continue

        # 2) Repeated band outliers
        suprathreshold_count = 0
        for indices in band_indices:
            if np.any(trace[indices] > test_band_threshold * mean_psd[indices]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces

def plot_psds_with_exclusion(
    ax,
    psd_array,
    freqs,
    kept_traces,
    excluded_traces,
    original_mean_psd,
    title,
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    alpha_excluded=0.05,
    alpha_kept=0.7,
    title_fontsize=8,
    axis_label_fontsize=8,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6
):
    if len(title) > max_title_length:
        title = title[:max_title_length] + "..."

    if show_kept:
        for idx_i, idx in enumerate(kept_traces):
            label = "Kept Trace" if idx_i == 0 else None
            ax.plot(freqs, psd_array[idx], color=color_kept, alpha=alpha_kept, label=label)

    if show_original_mean:
        ax.plot(freqs, original_mean_psd, color=color_old_mean, linewidth=2, label="Original Mean")

    if show_new_mean and len(kept_traces) > 0:
        new_mean_psd = np.nanmean(psd_array[kept_traces], axis=0)
        ax.plot(freqs, new_mean_psd, color=color_new_mean, linewidth=2, label="New Mean")

    if show_excluded:
        for idx_j, idx in enumerate(excluded_traces):
            label = "Excluded Trace" if idx_j == 0 else None
            ax.plot(freqs, psd_array[idx], color=color_excluded, alpha=alpha_excluded, zorder=10, label=label)

    if show_vertical_lines and vertical_lines is not None:
        for vfreq in vertical_lines:
            ax.axvline(
                vfreq,
                color=vertical_line_color,
                linestyle=vertical_line_style,
                alpha=vertical_line_alpha
            )

    ax.set_xlabel("Frequency (Hz)", fontsize=axis_label_fontsize)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_label_fontsize)
    ax.set_title(title, fontsize=title_fontsize)
    ax.legend(loc="upper right", fontsize=legend_fontsize)
    ax.tick_params(axis='both', labelsize=tick_label_fontsize)

def plot_psds_with_dropped_traces(
    psds_dict,
    rows_of_psds,
    low_band=(1,3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None,
    x_min=None,
    x_max=None,
    y_min=None,
    y_max=None,
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    num_cols=4,
    title_fontsize=8,
    axis_label_fontsize=8,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6
):
    figures = []
    kept_indices_dict = {}

    for row_idx, row in enumerate(rows_of_psds, start=1):
        valid_keys = [k for k in row if k in psds_dict]
        num_plots = len(valid_keys)
        if num_plots == 0:
            continue

        num_rows = math.ceil(num_plots / num_cols)
        fig, axes = plt.subplots(num_rows, num_cols, figsize=(num_cols * 4, num_rows * 4))
        plt.subplots_adjust(hspace=0.5)

        if num_rows * num_cols == 1:
            axes = [axes]
        else:
            axes = axes.flatten()

        for idx, key in enumerate(valid_keys):
            ax = axes[idx]
            psd_output = psds_dict[key]
            psd_data, freq_data = psd_output.get('psd'), psd_output.get('freqs')
            if psd_data is None or freq_data is None:
                ax.text(0.5, 0.5, f"No PSD data for {key}", ha='center', va='center')
                continue

            original_mean_psd = np.nanmean(psd_data, axis=0)

            kept_traces, excluded_traces = exclude_traces(
                psd_data,
                freq_data,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold
            )
            kept_indices_dict[key] = kept_traces

            plot_psds_with_exclusion(
                ax=ax,
                psd_array=psd_data,
                freqs=freq_data,
                kept_traces=kept_traces,
                excluded_traces=excluded_traces,
                original_mean_psd=original_mean_psd,
                title=key,
                show_kept=show_kept,
                show_excluded=show_excluded,
                show_original_mean=show_original_mean,
                show_new_mean=show_new_mean,
                show_vertical_lines=show_vertical_lines,
                color_kept=color_kept,
                color_excluded=color_excluded,
                color_old_mean=color_old_mean,
                color_new_mean=color_new_mean,
                alpha_excluded=0.05,
                alpha_kept=0.7,
                title_fontsize=title_fontsize,
                axis_label_fontsize=axis_label_fontsize,
                legend_fontsize=legend_fontsize,
                tick_label_fontsize=tick_label_fontsize,
                max_title_length=max_title_length,
                vertical_lines=vertical_lines,
                vertical_line_color=vertical_line_color,
                vertical_line_style=vertical_line_style,
                vertical_line_alpha=vertical_line_alpha
            )

            # x/y range
            if x_min is not None or x_max is not None:
                ax.set_xlim(x_min, x_max)
            if y_min is not None or y_max is not None:
                ax.set_ylim(y_min, y_max)

        # Turn off any unused subplots
        for ax in axes[num_plots:]:
            ax.axis('off')

        plt.tight_layout(rect=[0, 0.03, 1, 0.95])
        plt.suptitle(f"Row {row_idx}", fontsize=title_fontsize + 2)
        figures.append(fig)
        plt.show()

    return figures, kept_indices_dict

def group_keys_by_rows(psd_keys, row_size=4):
    """
    Sort PSD keys by block number and channel name, then group them row_size per row.

    Supports keys like:
      - "block1:Ch1"
      - "block2:Eye1_Ch3"
      - "Eye1_Ch1", "Eye2_Ch5", "EEG"
    """

    def sort_key(key):
        # Parse optional block prefix
        if ':' in key:
            block, ch = key.split(':', 1)
            try:
                block_num = int(block.replace('block', ''))
            except Exception:
                block_num = 0
        else:
            block_num = 0
            ch = key

        ch_num = 99999

        # Eye1/Eye2 explicit patterns
        m1 = re.search(r'Eye1_?Ch(\d+)', ch)
        m2 = re.search(r'Eye2_?Ch(\d+)', ch)
        if m1:
            ch_num = int(m1.group(1))
        elif m2:
            ch_num = 100 + int(m2.group(1))  # ensure Eye2 after Eye1
        elif ch.strip().upper() == 'EEG':
            ch_num = 1000
        else:
            # Generic "ChX" fallback
            m3 = re.search(r'Ch(\d+)', ch)
            if m3:
                try:
                    ch_num = int(m3.group(1))
                except Exception:
                    ch_num = 99999

        return block_num, ch_num

    sorted_keys = sorted(psd_keys, key=sort_key)
    rows = []
    for i in range(0, len(sorted_keys), row_size):
        rows.append(sorted_keys[i:i+row_size])
    return rows

##############################################################################
# 2) Segment & PSD Logic
##############################################################################

def load_epochs(filepath):
    if not os.path.isfile(filepath):
        raise FileNotFoundError(f"File not found: {filepath}")
    ep = mne.read_epochs(filepath, preload=True)
    if not isinstance(ep, mne.Epochs):
        data = ep.get_data()
        info = ep.info
        ep = mne.EpochsArray(data, info)
    return ep


def _clean_channel_matrix_nonfinite(
    ch_data,
    policy="drop_epochs",
    fill_value=0.0,
):
    """
    Clean a (n_epochs, n_times) array for non-finite values (NaN/Inf).

    Parameters
    ----------
    ch_data : np.ndarray
        Shape (n_epochs, n_times).
    policy : str
        One of:
          - 'drop_epochs'   : drop epochs that contain any non-finite values
          - 'interpolate'   : linearly interpolate non-finite values within each epoch
          - 'fill_constant' : replace non-finite values with `fill_value`
          - 'fill_median'   : replace non-finite values with per-epoch median of finite samples
          - 'error'         : raise ValueError if any non-finite values are present
    fill_value : float
        Used for fill policies and as fallback for interpolation failure.

    Returns
    -------
    cleaned : np.ndarray
        Cleaned array (may have fewer epochs if dropping).
    keep_mask : np.ndarray
        Boolean mask over input epochs indicating kept epochs (or all True if not dropping).
    stats : dict
        Summary stats: n_nonfinite, n_bad_epochs, n_dropped_epochs, policy
    """
    ch_data = np.asarray(ch_data)
    if ch_data.ndim != 2:
        raise ValueError(f"Expected 2D array (n_epochs, n_times); got shape {ch_data.shape}")

    finite = np.isfinite(ch_data)
    n_nonfinite = int((~finite).sum())
    n_epochs = ch_data.shape[0]

    if n_nonfinite == 0:
        return ch_data, np.ones(n_epochs, dtype=bool), {
            "n_nonfinite": 0,
            "n_bad_epochs": 0,
            "n_dropped_epochs": 0,
            "policy": policy,
        }

    if policy == "error":
        raise ValueError(f"Found {n_nonfinite} non-finite samples (NaN/Inf).")

    if policy == "drop_epochs":
        keep_mask = finite.all(axis=1)
        cleaned = ch_data[keep_mask, :]
        n_dropped = int((~keep_mask).sum())
        return cleaned, keep_mask, {
            "n_nonfinite": n_nonfinite,
            "n_bad_epochs": n_dropped,
            "n_dropped_epochs": n_dropped,
            "policy": policy,
        }

    # Non-dropping policies
    cleaned = ch_data.astype(float, copy=True)
    keep_mask = np.ones(n_epochs, dtype=bool)

    if policy == "fill_constant":
        cleaned[~finite] = float(fill_value)
        return cleaned, keep_mask, {
            "n_nonfinite": n_nonfinite,
            "n_bad_epochs": int((~finite.all(axis=1)).sum()),
            "n_dropped_epochs": 0,
            "policy": policy,
        }

    if policy == "fill_median":
        for ei in range(n_epochs):
            row = cleaned[ei, :]
            f = np.isfinite(row)
            if f.all():
                continue
            vals = row[f]
            rep = float(np.median(vals)) if vals.size else float(fill_value)
            row[~f] = rep
            cleaned[ei, :] = row
        return cleaned, keep_mask, {
            "n_nonfinite": n_nonfinite,
            "n_bad_epochs": int((~finite.all(axis=1)).sum()),
            "n_dropped_epochs": 0,
            "policy": policy,
        }

    if policy == "interpolate":
        n_time = cleaned.shape[1]
        x = np.arange(n_time, dtype=float)
        for ei in range(n_epochs):
            row = cleaned[ei, :]
            f = np.isfinite(row)
            if f.all():
                continue
            n_f = int(f.sum())
            if n_f == 0:
                row[:] = float(fill_value)
            elif n_f == 1:
                row[:] = float(row[f][0])
            else:
                row[~f] = np.interp(x[~f], x[f], row[f])
            cleaned[ei, :] = row
        return cleaned, keep_mask, {
            "n_nonfinite": n_nonfinite,
            "n_bad_epochs": int((~finite.all(axis=1)).sum()),
            "n_dropped_epochs": 0,
            "policy": policy,
        }

    raise ValueError(f"Unknown policy: {policy}")


def compute_welch_psd(
    epochs,
    fmin=0.0,
    fmax=100.0,
    window_s=2.0,
    overlap_pct=50.0,
    nan_policy="drop_epochs",
    nan_fill_value=0.0
):
    """
    NaN-aware Welch PSD for an mne.Epochs object using the longest contiguous
    clean segment per epoch (per channel).

    For each channel and epoch:
      - Find the longest run of finite samples (no NaN/Inf).
      - If that run is at least nperseg samples long, apply Welch on that segment.
      - If not, the PSD row for that epoch is set to NaN.

    fmin/fmax define the frequency window to keep. NaNs are never interpolated
    or filled; epochs that are too contaminated simply yield all-NaN PSD rows.
    """
    out = {}
    sfreq = float(epochs.info["sfreq"])
    data = epochs.get_data()  # (n_ep, n_ch, n_time)
    n_ep, n_ch, n_time = data.shape

    # Derive Welch parameters with clamping to epoch length (safer than error)
    nperseg = int(round(window_s * sfreq))
    nperseg = max(1, nperseg)
    noverlap = int(round(nperseg * overlap_pct / 100.0))

    if nperseg > n_time:
        print(
            f"WARNING: Requested window_s={window_s}s -> nperseg={nperseg} exceeds "
            f"epoch length={n_time} samples. Clamping nperseg to epoch length."
        )
        nperseg = int(n_time)
        noverlap = int(round(nperseg * overlap_pct / 100.0))

    if noverlap >= nperseg:
        noverlap = max(0, nperseg - 1)

    window = signal.hamming(nperseg, sym=False)

    # Optional global summary of non-finite samples (for logging only)
    nonfinite_total = int((~np.isfinite(data)).sum())
    if nonfinite_total > 0:
        bad_epochs_any = int((~np.isfinite(data).all(axis=2)).any(axis=1).sum())
        print(f"[NaN-aware Welch] Detected {nonfinite_total} non-finite samples "
              f"across {bad_epochs_any}/{n_ep} epochs. Longest-clean-segment rule applied.")

    for ch_idx, ch_name in enumerate(epochs.ch_names):
        ch_data = data[:, ch_idx, :]  # (n_ep, n_time)
        freqs_ref = None
        psd_rows = []

        for ep_idx in range(n_ep):
            x = ch_data[ep_idx]
            freqs_ep, psd_ep = _welch_longest_clean_segment_1d(
                x, sfreq, nperseg, noverlap, window
            )
            if freqs_ep is None or psd_ep is None:
                # Not enough clean data
                if freqs_ref is None:
                    psd_rows.append(None)
                else:
                    psd_rows.append(np.full_like(freqs_ref, np.nan))
                continue

            # First epoch with valid PSD => set fmin/fmax and reference freqs
            if freqs_ref is None:
                freq_mask = (freqs_ep >= fmin) & (freqs_ep <= fmax)
                freqs_ref = freqs_ep[freq_mask]
                psd_sel = psd_ep[freq_mask]
            else:
                freq_mask = (freqs_ep >= fmin) & (freqs_ep <= fmax)
                psd_sel = psd_ep[freq_mask]
                if psd_sel.shape != freqs_ref.shape:
                    raise RuntimeError(
                        f"Frequency grid mismatch in channel '{ch_name}'. "
                        "Check Welch parameters."
                    )
            psd_rows.append(psd_sel)

        if freqs_ref is None:
            out[ch_name] = {"freqs": None, "psd": None}
            continue

        psd_array = []
        for row in psd_rows:
            if row is None:
                psd_array.append(np.full_like(freqs_ref, np.nan))
            else:
                psd_array.append(row)
        psd_array = np.vstack(psd_array)

        out[ch_name] = {"freqs": freqs_ref, "psd": psd_array}

    return out

##############################################################################
# 3) MAIN GUI
##############################################################################

def build_exportable_plot_psd_gui():
    """
    PSD Clean GUI w/Segmentation + ipyfilechooser for figure & cleaned PSD export.
    Channel selection is placed in a more prominent position,
    and MNE warnings are silenced via mne.set_log_level('ERROR').
    """

    # ~~~~~ GUI ELEMENTS ~~~~~
    load_psd_button = widgets.Button(description='Load / Compute PSD', button_style='info')
    fif_chooser = FileChooser(os.getcwd(), title='Select Filtered .fif', select_default=False)
    fif_chooser.show_only_files = True
    fif_chooser.filter_pattern = ['*.fif', '*.fif.gz']

    # Choose whether to compute PSD from .fif or load a precomputed PSD pickle
    input_mode_widget = widgets.ToggleButtons(
        options=['Compute from .fif', 'Load PSD .pkl'],
        value='Compute from .fif',
        description='Input:',
        style={'description_width': 'initial'}
    )

    psd_pkl_chooser = FileChooser(os.getcwd(), title='Select PSD results .pkl', select_default=False)
    psd_pkl_chooser.show_only_files = True
    psd_pkl_chooser.filter_pattern = ['*.pkl']


    segment_toggle = widgets.ToggleButtons(
        options=['No', 'Yes'],
        value='No',
        description='Segment 6 blocks?',
        style={'description_width': 'initial'}
    )

    window_length_widget = widgets.FloatText(value=2.0, description='Window (s):', layout=widgets.Layout(width='200px'))
    overlap_widget = widgets.FloatSlider(value=50.0, min=0.0, max=100.0, step=1.0, 
                                         description='Overlap(%):', 
                                         continuous_update=False, 
                                         layout=widgets.Layout(width='300px'))

            # NaN/Inf handling (applies when computing PSD from .fif)
    # In this version, NaNs are handled via a fixed, NaN-aware Welch rule:
    #   - For each epoch & channel, only the longest contiguous clean segment
    #     (no NaN/Inf) is used for PSD.
    #   - If no segment is long enough for Welch, that epoch's PSD row is set to NaN.
    #   - No interpolation or constant filling is applied.
    nan_policy_widget = widgets.Dropdown(
        options=[
            ('NaN-aware Welch (longest clean segment)', 'nan_aware'),
        ],
        value='nan_aware',
        description='NaN handling:',
        style={'description_width': 'initial'},
        layout=widgets.Layout(width='420px'),
        disabled=True
    )
    nan_fill_value_widget = widgets.FloatText(
        value=0.0,
        description='Fill value:',
        layout=widgets.Layout(width='200px')
    )
    # Hide fill value (not used in NaN-aware Welch policy)
    nan_fill_value_widget.layout.display = 'none'

# Instead of placing channel selection far to the right, let's create it here
    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description='Select Channels:',
        layout=widgets.Layout(width='250px', height='200px')
    )

    # For the load messages / progress
    output_area = widgets.Output()

    # We'll store final PSD in a dictionary "loaded_psd" with possibly multiple blocks: "block1:ChA", etc.
    loaded_psd = {}

    # Booleans for plot
    show_kept_cb = widgets.Checkbox(value=True, description='Show Kept Traces')
    show_excluded_cb = widgets.Checkbox(value=True, description='Show Excluded Traces')
    show_original_mean_cb = widgets.Checkbox(value=True, description='Show Original Mean')
    show_new_mean_cb = widgets.Checkbox(value=True, description='Show New Mean')
    show_vertical_lines_cb = widgets.Checkbox(value=True, description='Show Vertical Lines')

    vertical_lines_text = widgets.Text(value="10,15", description='Vertical Lines(Hz):', layout=widgets.Layout(width='300px'))

    # Axis
    x_min_widget = widgets.FloatText(value=None, description='X Min:', layout=widgets.Layout(width='200px'))
    x_max_widget = widgets.FloatText(value=None, description='X Max:', layout=widgets.Layout(width='200px'))
    y_min_widget = widgets.FloatText(value=None, description='Y Min:', layout=widgets.Layout(width='200px'))
    y_max_widget = widgets.FloatText(value=None, description='Y Max:', layout=widgets.Layout(width='200px'))

    # Fonts
    title_fs_widget = widgets.IntText(value=8, description='Title FS:', layout=widgets.Layout(width='200px'))
    axis_fs_widget = widgets.IntText(value=8, description='Axis FS:', layout=widgets.Layout(width='200px'))
    legend_fs_widget = widgets.IntText(value=8, description='Legend FS:', layout=widgets.Layout(width='200px'))
    tick_fs_widget = widgets.IntText(value=8, description='Tick FS:', layout=widgets.Layout(width='200px'))
    max_title_len_widget = widgets.IntText(value=40, description='Max Title:', layout=widgets.Layout(width='200px'))

    # Colors
    color_kept_widget = widgets.ColorPicker(value='lightgray', description='Kept:', layout=widgets.Layout(width='220px'))
    color_excl_widget = widgets.ColorPicker(value='red', description='Excluded:', layout=widgets.Layout(width='220px'))
    color_oldmean_widget = widgets.ColorPicker(value='blue', description='Orig Mean:', layout=widgets.Layout(width='220px'))
    color_newmean_widget = widgets.ColorPicker(value='green', description='New Mean:', layout=widgets.Layout(width='220px'))

    # Exclusion thresholds
    low_band_min_widget = widgets.FloatText(value=1.0, description='LowBand Min:', layout=widgets.Layout(width='200px'))
    low_band_max_widget = widgets.FloatText(value=3.0, description='LowBand Max:', layout=widgets.Layout(width='200px'))
    low_band_thr_widget = widgets.FloatText(value=3.0, description='LowBand Thr:', layout=widgets.Layout(width='200px'))
    test_band_thr_widget = widgets.FloatText(value=10.0, description='TestBand Thr:', layout=widgets.Layout(width='200px'))

    test_band_text = widgets.Textarea(
        value="(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)",
        description='Test Bands:',
        layout=widgets.Layout(width='300px', height='80px')
    )

    num_cols_widget = widgets.IntText(value=4, description='Cols/Row:', layout=widgets.Layout(width='120px'))

    # Buttons
    compute_psd_button = widgets.Button(description='Compute PSD', button_style='success')
    plot_psd_button = widgets.Button(description='Plot PSDs', button_style='success')

    # Export controls
    export_fig_chooser = FileChooser(
        os.getcwd(),
        title='Choose or type a figure export name (.png, .pptx, etc.)',
        select_default=False
    )

    export_button = widgets.Button(description='Export Figures', button_style='warning')

    # Cleaned PSD export
    cleaned_pickle_chooser = FileChooser(
        os.getcwd(),
        title='Choose or type a .pkl for cleaned PSD',
        select_default=False
    )
    cleaned_pickle_chooser.filter_pattern = ['*.pkl']
    export_cleaned_button = widgets.Button(description='Export Cleaned PSD', button_style='info')

    # Output areas
    plot_output_area = widgets.Output()
    export_output_area = widgets.Output()

    current_figures = []
    kept_indices_dict = {}

    # ~~~~ A) LOAD & SEGMENT & PSD ~~~~
    def on_load_psd_clicked(b):
        with output_area:
            clear_output()
            chosen_mode = input_mode_widget.value

            # ---------------------------------------------------------
            # Mode A: Load a precomputed PSD dictionary from .pkl
            # ---------------------------------------------------------
            if chosen_mode == 'Load PSD .pkl':
                chosen_pkl = psd_pkl_chooser.selected
                if not chosen_pkl:
                    print("No PSD .pkl selected.")
                    return
                if not os.path.isfile(chosen_pkl):
                    print(f"File not found: {chosen_pkl}")
                    return
                try:
                    print(f"Loading PSD pickle: {chosen_pkl}")
                    with open(chosen_pkl, 'rb') as f:
                        obj = pickle.load(f)
                except Exception as e:
                    print(f"Error loading PSD pickle: {e}")
                    return

                # Accept either:
                #   - dict[channel_key] = {'psd': array, 'freqs': array}
                #   - dict with a nested field like {'psd_results': {...}}
                if isinstance(obj, dict) and 'psd_results' in obj and isinstance(obj['psd_results'], dict):
                    psd_dict = obj['psd_results']
                else:
                    psd_dict = obj

                if not isinstance(psd_dict, dict):
                    print("ERROR: PSD pickle does not contain a dictionary.")
                    return

                loaded_psd.clear()
                loaded_psd.update(psd_dict)

                all_keys = sorted(loaded_psd.keys())
                channels_dropdown.options = all_keys
                print(f"Loaded PSD entries: {len(all_keys)}")
                print("Channels/Blocks in dropdown =>", all_keys)
                return

            # ---------------------------------------------------------
            # Mode B: Compute PSD from a .fif file (legacy behavior)
            # ---------------------------------------------------------
            chosen_fif = fif_chooser.selected
            if not chosen_fif:
                print("No .fif file selected.")
                return
            if not os.path.isfile(chosen_fif):
                print(f"File not found: {chosen_fif}")
                return
            try:
                print(f"Loading: {chosen_fif}")
                epochs_all = load_epochs(chosen_fif)
                n_total_epochs = len(epochs_all)
                print(f"Loaded => {n_total_epochs} epochs, {len(epochs_all.ch_names)} channels.")
            except Exception as e:
                print(f"Error loading epochs: {e}")
                return

            do_segment = (segment_toggle.value == 'Yes')
            w_len = window_length_widget.value
            overlap_val = overlap_widget.value

            if w_len <= 0:
                print("ERROR: Window length must be >0.")
                return
            if not (0 <= overlap_val <= 100):
                print("ERROR: Overlap must be 0-100.")
                return

            loaded_psd.clear()

            if do_segment:
                print("Segmenting => up to 6 blocks of 120 each (plus leftover).")
                block_size = 120
                n_blocks = n_total_epochs // block_size
                for i in range(n_blocks):
                    start_ep = i*block_size
                    end_ep = (i+1)*block_size
                    block_ = epochs_all[start_ep:end_ep]
                    block_name = f"block{i+1}"
                    print(f"  Computing PSD for {block_name}: epochs {start_ep}..{end_ep-1}")
                    block_psd = compute_welch_psd(block_, window_s=w_len, overlap_pct=overlap_val, nan_policy=nan_policy_widget.value, nan_fill_value=nan_fill_value_widget.value)
                    for ch, info in block_psd.items():
                        new_key = f"{block_name}:{ch}"
                        loaded_psd[new_key] = info

                leftover = n_total_epochs - (n_blocks*block_size)
                if leftover > 0:
                    print(f"  Leftover block => {leftover} epochs.")
                    block_name = f"block{n_blocks+1}"
                    block_ = epochs_all[n_blocks*block_size:]
                    block_psd = compute_welch_psd(block_, window_s=w_len, overlap_pct=overlap_val, nan_policy=nan_policy_widget.value, nan_fill_value=nan_fill_value_widget.value)
                    for ch, info in block_psd.items():
                        new_key = f"{block_name}:{ch}"
                        loaded_psd[new_key] = info
            else:
                print("No segmentation => single block with all epochs.")
                block_psd = compute_welch_psd(epochs_all, window_s=w_len, overlap_pct=overlap_val, nan_policy=nan_policy_widget.value, nan_fill_value=nan_fill_value_widget.value)
                for ch, info in block_psd.items():
                    loaded_psd[ch] = info

            all_keys = sorted(loaded_psd.keys())
            channels_dropdown.options = all_keys
            print("Channels/Blocks in dropdown =>", all_keys)

    load_psd_button.on_click(on_load_psd_clicked)

    # ~~~~ B) PLOT PSD with dropping ~~~~
    def on_plot_psd_clicked(b):
        nonlocal current_figures, kept_indices_dict
        with plot_output_area:
            clear_output()
            if not loaded_psd:
                print("No PSD data loaded. Please load/compute PSD first.")
                return
            selected_keys = list(channels_dropdown.value)
            if not selected_keys:
                print("No channels selected.")
                return

            show_kept = show_kept_cb.value
            show_excluded = show_excluded_cb.value
            show_original_mean = show_original_mean_cb.value
            show_new_mean = show_new_mean_cb.value
            show_vertical_lines = show_vertical_lines_cb.value

            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value

            title_fs = title_fs_widget.value
            axis_fs = axis_fs_widget.value
            legend_fs = legend_fs_widget.value
            tick_fs = tick_fs_widget.value
            max_title_len = max_title_len_widget.value

            color_kept = color_kept_widget.value
            color_excl = color_excl_widget.value
            color_old_mean = color_oldmean_widget.value
            color_new_mean = color_newmean_widget.value

            lb_min = low_band_min_widget.value
            lb_max = low_band_max_widget.value
            low_band_threshold = low_band_thr_widget.value
            test_band_threshold = test_band_thr_widget.value

            # Parse test bands
            test_str = test_band_text.value.strip()
            test_bands_list = []
            if test_str:
                test_str = test_str.replace(" ", "")
                pairs = test_str.split(")")
                for p in pairs:
                    p = p.strip(",").strip().strip("(").strip()
                    if not p:
                        continue
                    vals = p.split(",")
                    if len(vals) == 2:
                        try:
                            lowf = float(vals[0])
                            highf = float(vals[1])
                            test_bands_list.append((lowf, highf))
                        except:
                            pass
                if not test_bands_list:
                    test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]
            else:
                test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]

            num_cols = num_cols_widget.value
            if num_cols < 1:
                num_cols = 4

            vlines_str = vertical_lines_text.value.strip()
            vert_lines_list = []
            if vlines_str:
                parts = vlines_str.split(",")
                for part in parts:
                    part = part.strip()
                    if part:
                        try:
                            freq_val = float(part)
                            vert_lines_list.append(freq_val)
                        except:
                            pass

            # Build dictionary with selected keys only
            psd_data_dict = {}
            for k in selected_keys:
                psd_dict = loaded_psd.get(k, None)
                if psd_dict is not None:
                    psd_data_dict[k] = psd_dict

            rows_of_psds = group_keys_by_rows(selected_keys, row_size=num_cols)

            figs, kept_dict = plot_psds_with_dropped_traces(
                psds_dict=psd_data_dict,
                rows_of_psds=rows_of_psds,
                low_band=(lb_min, lb_max),
                low_band_threshold=low_band_threshold,
                test_bands=test_bands_list,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=None,
                x_min=x_min,
                x_max=x_max,
                y_min=y_min,
                y_max=y_max,
                show_kept=show_kept,
                show_excluded=show_excluded,
                show_original_mean=show_original_mean,
                show_new_mean=show_new_mean,
                show_vertical_lines=show_vertical_lines,
                color_kept=color_kept,
                color_excluded=color_excl,
                color_old_mean=color_old_mean,
                color_new_mean=color_new_mean,
                num_cols=num_cols,
                title_fontsize=title_fs,
                axis_label_fontsize=axis_fs,
                legend_fontsize=legend_fs,
                tick_label_fontsize=tick_fs,
                max_title_length=max_title_len,
                vertical_lines=vert_lines_list,
                vertical_line_color="black",
                vertical_line_style="--",
                vertical_line_alpha=0.6
            )

            current_figures = figs
            kept_indices_dict = kept_dict

            if figs:
                print(f"Plotted {len(figs)} figure(s) for {len(selected_keys)} key(s).")
                print("Kept indices recorded in 'kept_indices_dict'.")
            else:
                print("No figures plotted.")

    plot_psd_button.on_click(on_plot_psd_clicked)

    # ~~~~ C) EXPORT FIGURES (using FileChooser) ~~~~
    def on_export_button_clicked(b):
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No figures to export.")
                return

            chosen_path = export_fig_chooser.selected
            if not chosen_path:
                print("No export file chosen. Please select or type a filename (e.g. MyPlot.png or MySlides.pptx).")
                return

            base, ext = os.path.splitext(chosen_path)
            ext = ext.lower()
            if not ext:
                print("No file extension detected. Please include something like .png or .pptx.")
                return

            if ext in ('.ppt', '.pptx'):
                if not HAS_PPTX:
                    print("python-pptx not installed. Cannot export PPT.")
                    return
                ppt_file = chosen_path
                print(f"Exporting {len(current_figures)} figure(s) to {ppt_file} ...")
                prs = Presentation()
                blank_layout = prs.slide_layouts[6]
                for i, fig in enumerate(current_figures, start=1):
                    temp_png = f"{base}_temp_{i}.png"
                    fig.savefig(temp_png, format='png', dpi=150)
                    slide = prs.slides.add_slide(blank_layout)
                    left = top = Inches(1)
                    slide.shapes.add_picture(temp_png, left, top, width=Inches(8), height=Inches(4.5))
                    os.remove(temp_png)
                prs.save(ppt_file)
                print(f"Done exporting to {ppt_file}")
            else:
                # treat it as an image format
                print(f"Exporting {len(current_figures)} figure(s) as *{ext} images.")
                for i, fig in enumerate(current_figures, start=1):
                    out_file = f"{base}_{i}{ext}"
                    fig.savefig(out_file, format=ext.lstrip('.'), dpi=150)
                    print(f"Saved => {out_file}")
                print("Done exporting images.")

    export_button.on_click(on_export_button_clicked)

    # ~~~~ D) EXPORT CLEANED PSD (using FileChooser) ~~~~
    def on_export_cleaned_clicked(b):
        with export_output_area:
            clear_output()
            if not kept_indices_dict:
                print("No kept indices found. Please plot first.")
                return
            if not loaded_psd:
                print("No PSD data loaded.")
                return

            chosen_pkl_path = cleaned_pickle_chooser.selected
            if not chosen_pkl_path:
                print("No output .pkl chosen. Please select or type a filename (e.g. CleanedPSD.pkl).")
                return

            if not chosen_pkl_path.lower().endswith('.pkl'):
                print("Output path must end with '.pkl'")
                return

            cleaned_data = {}
            for key, psd_info in loaded_psd.items():
                if key not in kept_indices_dict:
                    continue
                kept_ix = kept_indices_dict[key]
                psd_array = psd_info.get('psd', None)
                freqs = psd_info.get('freqs', None)
                if psd_array is None or freqs is None:
                    continue
                if len(kept_ix) == 0:
                    new_psd = None
                else:
                    new_psd = psd_array[kept_ix, :]
                cleaned_data[key] = {"freqs": freqs, "psd": new_psd}

            if not cleaned_data:
                print("No cleaned data generated (maybe all epochs excluded?).")
                return

            try:
                with open(chosen_pkl_path, 'wb') as f:
                    pickle.dump(cleaned_data, f)
                print(f"Saved cleaned PSD => {chosen_pkl_path}")
                print(f"Channels/Blocks in cleaned data => {list(cleaned_data.keys())}")
            except Exception as e:
                print(f"Error saving cleaned PSD: {e}")

    export_cleaned_button.on_click(on_export_cleaned_clicked)

    # ~~~~ LAYOUT ~~~~
    load_box = widgets.VBox([
        widgets.HTML("<h3>Load PSD Input</h3>"),
        input_mode_widget,
        load_psd_button,
        widgets.HTML("<b>.fif source (compute PSD here)</b>"),
        fif_chooser,
        segment_toggle,
        widgets.HBox([widgets.Label("Window:"), window_length_widget,
                       widgets.Label("Overlap:"), overlap_widget]),
        widgets.HBox([nan_policy_widget, nan_fill_value_widget]),
        widgets.HTML("<b>PSD .pkl source (load precomputed)</b>"),
        psd_pkl_chooser,
    ])

    # Put the channels dropdown here so it's visible & accessible
    channel_box = widgets.VBox([
        widgets.HTML("<h3>Select Channels / Blocks</h3>"),
        channels_dropdown
    ])

    # Combining load, channel, and output area in one vertical stack
    top_panel = widgets.VBox([
        load_box,
        channel_box,
        output_area  # messages appear here
    ])

    exclusion_params_box = widgets.HBox([
        widgets.VBox([
            widgets.Label("Low-band Range:"),
            low_band_min_widget,
            low_band_max_widget,
            low_band_thr_widget
        ]),
        widgets.VBox([
            widgets.Label("Test-band Th:"),
            test_band_thr_widget,
            widgets.Label("Test Bands:"),
            test_band_text
        ])
    ])

    axis_range_box = widgets.HBox([
        widgets.VBox([x_min_widget, x_max_widget]),
        widgets.VBox([y_min_widget, y_max_widget])
    ])

    font_size_box = widgets.HBox([
        title_fs_widget, axis_fs_widget, legend_fs_widget, tick_fs_widget, max_title_len_widget
    ])

    color_box = widgets.HBox([
        color_kept_widget, color_excl_widget, color_oldmean_widget, color_newmean_widget
    ])

    plot_box = widgets.VBox([
        widgets.HTML("<h3>Plot PSD with Dropping</h3>"),
        show_kept_cb, show_excluded_cb, show_original_mean_cb, show_new_mean_cb,
        show_vertical_lines_cb, vertical_lines_text,
        widgets.HTML("<b>Exclusion Params:</b>"),
        exclusion_params_box,
        widgets.HTML("<b>Axes & Layout:</b>"),
        axis_range_box,
        widgets.HTML("<b>Font Sizes:</b>"),
        font_size_box,
        widgets.HTML("<b>Colors:</b>"),
        color_box,
        widgets.HTML("<b>Columns per Row:</b>"),
        num_cols_widget,
        plot_psd_button,
        plot_output_area
    ])

    export_box = widgets.VBox([
        widgets.HTML("<h3>Export Figures</h3>"),
        export_fig_chooser,
        export_button,
        widgets.HTML("<h3>Export Cleaned PSD</h3>"),
        cleaned_pickle_chooser,
        export_cleaned_button,
        export_output_area
    ])

    ui = widgets.VBox([
        widgets.HTML("<h2>PSD Clean GUI (compute from .fif or load precomputed PSD .pkl)</h2>"),
        widgets.HBox([top_panel, plot_box]),
        export_box
    ])

    display(ui)

# Build & display the GUI
build_exportable_plot_psd_gui()