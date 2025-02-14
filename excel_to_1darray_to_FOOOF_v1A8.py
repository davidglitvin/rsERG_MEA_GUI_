import os
import re
import pickle
import pandas as pd
import numpy as np
import ipywidgets as widgets
from IPython.display import display, clear_output
import matplotlib.pyplot as plt
from fooof import FOOOFGroup
from scipy.spatial.distance import cdist
from tqdm.notebook import tqdm

# Tkinter for file/directory dialogs
from tkinter import Tk, filedialog

# Optional: For PPT export (install python-pptx if needed)
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

###############################################################################
#                           CORE ANALYSIS FUNCTIONS
###############################################################################

def process_excel(file_paths, include_block=True, include_eye=True):
    """
    Reads .xlsx files and computes PSD data dictionaries.
    Only processes sheets based on filtering:
      - Sheets with block info (B1-B6) are processed only if include_block is True.
      - Sheets with "Ey" (e.g., Ey1, Ey2) are processed only if include_eye is True.
    Returns a dict: { file_path: { sheet_name: PSD_values, 'freq': freq_array } }
    """
    psd_data_dict_all = {}

    for file_path in file_paths:
        try:
            excel_data = pd.ExcelFile(file_path)
            psd_data_dict = {}
            freq_data = None

            for sheet_name in excel_data.sheet_names:
                is_block = bool(re.search(r'B[1-6]', sheet_name))
                is_eye = ("Ey" in sheet_name) and not is_block

                if is_block and not include_block:
                    continue
                if is_eye and not include_eye:
                    continue
                if not is_block and not is_eye:
                    continue

                sheet_df = pd.read_excel(file_path, sheet_name=sheet_name)
                freq_column = sheet_df.iloc[:, 0]
                if freq_data is None:
                    freq_data = freq_column.values

                psd_columns = sheet_df.iloc[:, 1:]
                averaged_psd = psd_columns.mean(axis=1)
                psd_data_dict[sheet_name] = averaged_psd.values

            if freq_data is not None:
                psd_data_dict["freq"] = freq_data
                psd_data_dict_all[file_path] = psd_data_dict
            else:
                print(f"Warning: No sheets met the filtering criteria in {file_path}")

        except Exception as e:
            print(f"An error occurred while processing {file_path}: {e}")

    return psd_data_dict_all

def run_fooof_analysis(psd_data_dict_all, freq_range, amp_threshold, r2_threshold, max_peaks, fitting_mode, peak_width_limits):
    """
    Runs FOOOF analysis on each sheet's PSD data.
    Returns a dictionary (fg_dict) mapping a unique key (sheet + filename) -> FOOOFGroup object.
    """
    fg_dict = {}

    for file_path, psd_data_dict in tqdm(psd_data_dict_all.items(), desc="Fitting PSDs with FOOOF"):
        freq = psd_data_dict["freq"]
        file_suffix = os.path.basename(file_path).replace(".xlsx", "")

        for sheet_name, psd in psd_data_dict.items():
            if sheet_name == "freq":
                continue

            fg = FOOOFGroup(
                peak_width_limits=peak_width_limits,
                max_n_peaks=max_peaks,
                min_peak_height=amp_threshold,
                verbose=False,
                aperiodic_mode=fitting_mode,
            )

            fg.fit(freq, psd[np.newaxis, :], freq_range)
            # Drop any fits below the provided R² threshold
            fg.drop(fg.get_params('r_squared') < r2_threshold)

            dict_key = f"{sheet_name}_{file_suffix}"
            fg_dict[dict_key] = fg

    return fg_dict

def plot_closest_to_mean(fg, sheet_name, export_dir, base_filename, formats,
                         x_min, x_max, y_min, y_max,
                         show_grid=True, x_tick_font_size=8,
                         include_r2=False, include_peak_table=False):
    """
    For FOOOFGroup objects with multiple PSD fits, this function selects
    10 spectra closest to the mean and plots them. Optionally includes r² in
    the title and displays a table of detected peak parameters beneath each plot.
    """
    models = [fg.get_fooof(ind=i) for i in range(len(fg))]
    spectra = [model.power_spectrum for model in models]
    spectra = np.array(spectra)

    if spectra.ndim != 2:
        raise ValueError(f"Unexpected spectra shape: {spectra.shape}. Expected a 2D array.")

    mean_spectrum = np.mean(spectra, axis=0)
    distances = cdist(spectra, mean_spectrum[None, :], metric='euclidean').flatten()
    closest_indices = np.argsort(distances)[:10]

    if "ppt" in formats and PPTX_AVAILABLE:
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank layout

    for i, index in enumerate(closest_indices):
        fm = fg.get_fooof(ind=index, regenerate=True)
        r2_value = fm.get_params('r_squared')
        if include_r2:
            title = f"{sheet_name} - Closest Example {index} (R²: {r2_value:.2f})"
        else:
            title = f"{sheet_name} - Closest Example {index}"
        fm.plot(title=title, plot_peaks='shade')
        plt.xlim(x_min, x_max)
        plt.ylim(y_min, y_max)
        x_min_int = int(np.floor(x_min))
        x_max_int = int(np.ceil(x_max))
        plt.xticks(range(x_min_int, x_max_int+1, 1))
        plt.setp(plt.gca().get_xticklabels(), fontsize=x_tick_font_size)
        plt.grid(show_grid)
        plt.margins(0)

        fig = plt.gcf()
        file_tag = f"{base_filename}_{sheet_name}_ex{i}"

        for fmt in formats:
            if fmt in ["png", "svg", "jpeg"]:
                filename = os.path.join(export_dir, f"{file_tag}.{fmt}")
                fig.savefig(filename, dpi=300, bbox_inches='tight')
            elif fmt == "ppt" and PPTX_AVAILABLE:
                tmp_png = os.path.join(export_dir, f"{file_tag}_ppt_temp.png")
                fig.savefig(tmp_png, dpi=300, bbox_inches='tight')
                slide = prs.slides.add_slide(slide_layout)
                left = top = Inches(1)
                slide.shapes.add_picture(tmp_png, left, top, height=Inches(5))
        plt.show()
        
        # Optionally display peak parameters as a table beneath the plot
        if include_peak_table:
            peak_params = fm.peak_params_
            if peak_params.size > 0:
                df_peaks = pd.DataFrame(peak_params, columns=['Center Frequency (Hz)', 'Amplitude', 'FWHM'])
                print("Detected Peaks:")
                display(df_peaks)
            else:
                print("No peaks detected.")

        plt.clf()

    if "ppt" in formats and PPTX_AVAILABLE:
        pptx_filename = os.path.join(export_dir, f"{base_filename}_{sheet_name}.pptx")
        prs.save(pptx_filename)

def plot_all_psds(fg, sheet_name, export_dir, base_filename, formats,
                  x_min, x_max, y_min, y_max,
                  show_grid=True, x_tick_font_size=8,
                  include_r2=False, include_peak_table=False):
    """
    For FOOOFGroup objects (with one or multiple PSDs), this function plots each PSD.
    Optionally includes r² in the title and displays a table of peak parameters beneath each plot.
    """
    n_fits = len(fg)
    if "ppt" in formats and PPTX_AVAILABLE:
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
    
    for i in range(n_fits):
        fm = fg.get_fooof(ind=i, regenerate=True)
        r2_value = fm.get_params('r_squared')
        if include_r2:
            title = f"{sheet_name} - Example {i} (R²: {r2_value:.2f})"
        else:
            title = f"{sheet_name} - Example {i}"
        fm.plot(title=title, plot_peaks='shade')
        plt.xlim(x_min, x_max)
        plt.ylim(y_min, y_max)
        x_min_int = int(np.floor(x_min))
        x_max_int = int(np.ceil(x_max))
        plt.xticks(range(x_min_int, x_max_int+1, 1))
        plt.setp(plt.gca().get_xticklabels(), fontsize=x_tick_font_size)
        plt.grid(show_grid)
        plt.margins(0)

        fig = plt.gcf()
        file_tag = f"{base_filename}_{sheet_name}_ex{i}"
        for fmt in formats:
            if fmt in ["png", "svg", "jpeg"]:
                filename = os.path.join(export_dir, f"{file_tag}.{fmt}")
                fig.savefig(filename, dpi=300, bbox_inches='tight')
            elif fmt == "ppt" and PPTX_AVAILABLE:
                tmp_png = os.path.join(export_dir, f"{file_tag}_ppt_temp.png")
                fig.savefig(tmp_png, dpi=300, bbox_inches='tight')
                slide = prs.slides.add_slide(slide_layout)
                left = top = Inches(1)
                slide.shapes.add_picture(tmp_png, left, top, height=Inches(5))
        plt.show()
        
        # Optionally display peak parameters as a table beneath the plot
        if include_peak_table:
            peak_params = fm.peak_params_
            if peak_params.size > 0:
                df_peaks = pd.DataFrame(peak_params, columns=['Center Frequency (Hz)', 'Amplitude', 'FWHM'])
                print("Detected Peaks:")
                display(df_peaks)
            else:
                print("No peaks detected.")
                
        plt.clf()
    
    if "ppt" in formats and PPTX_AVAILABLE:
        pptx_filename = os.path.join(export_dir, f"{base_filename}_{sheet_name}.pptx")
        prs.save(pptx_filename)

def save_fooof_group(fg_dict, filename="fooof_groups.pkl"):
    """
    Serializes and saves the fg_dict (mapping keys -> FOOOFGroup) to a pickle file.
    """
    with open(filename, "wb") as f:
        pickle.dump(fg_dict, f)
    print(f"FOOOFGroup dictionary saved to {filename}")

def export_fg_dict_excel(fg_dict, excel_filename):
    """
    Exports FOOOFGroup parameters to an Excel file.
    For each key in fg_dict, each fit's parameters (r², aperiodic and peak parameters)
    are stored as a row in the Excel sheet.
    """
    rows = []
    for key, fg in fg_dict.items():
        n_fits = len(fg)
        r2_params = fg.get_params('r_squared')
        ap_params = fg.get_params('aperiodic_params')
        peak_params = fg.get_params('peak_params')
        
        for i in range(n_fits):
            r2 = r2_params[i] if i < len(r2_params) else None
            ap = ap_params[i] if i < len(ap_params) else None
            peaks = peak_params[i] if i < len(peak_params) else None
            
            row = {
                "dict_key": key,
                "fit_index": i,
                "r_squared": r2,
                "aperiodic_params": str(ap),
                "peak_params": str(peaks)
            }
            rows.append(row)
    df = pd.DataFrame(rows)
    df.to_excel(excel_filename, index=False)
    print(f"FOOOF results exported to Excel file: {excel_filename}")

def display_fg_dict_table(fg_dict):
    """
    Displays FOOOFGroup parameters in an Excel-like table within the notebook.
    Each row represents one fit and its key parameters.
    """
    rows = []
    for key, fg in fg_dict.items():
        n_fits = len(fg)
        r2_params = fg.get_params('r_squared')
        ap_params = fg.get_params('aperiodic_params')
        peak_params = fg.get_params('peak_params')
        for i in range(n_fits):
            r2 = r2_params[i] if i < len(r2_params) else None
            ap = ap_params[i] if i < len(ap_params) else None
            peaks = peak_params[i] if i < len(peak_params) else None
            row = {
                "dict_key": key,
                "fit_index": i,
                "r_squared": r2,
                "aperiodic_params": str(ap),
                "peak_params": str(peaks)
            }
            rows.append(row)
    df = pd.DataFrame(rows)
    print("FOOOF Group Export Table:")
    display(df)

###############################################################################
#                         TKINTER-BASED GUI CODE & WIDGETS
###############################################################################

selected_files = []
export_directory = ""
fg_dict = {}  # Will store final FOOOF results

# Parameter widgets
freq_range_min = widgets.FloatText(description="Freq Min:", value=4.0, layout=widgets.Layout(width="30%"))
freq_range_max = widgets.FloatText(description="Freq Max:", value=45.0, layout=widgets.Layout(width="30%"))
amp_threshold = widgets.FloatText(description="Amplitude Threshold:", value=0.2, layout=widgets.Layout(width="50%"))
r2_threshold = widgets.FloatText(description="R² Threshold:", value=0.5, layout=widgets.Layout(width="50%"))
max_peaks = widgets.IntText(description="Max Peaks:", value=2, layout=widgets.Layout(width="50%"))
fitting_mode = widgets.Dropdown(description="Fitting Mode:", options=["knee", "fixed"], value="knee", layout=widgets.Layout(width="50%"))
peak_width_min = widgets.FloatText(description="Peak Width Min:", value=2.0, layout=widgets.Layout(width="50%"))
peak_width_max = widgets.FloatText(description="Peak Width Max:", value=10.0, layout=widgets.Layout(width="50%"))

# Axis-range widgets
x_axis_min = widgets.FloatText(description="X Axis Min:", value=4.0, layout=widgets.Layout(width="33%"))
x_axis_max = widgets.FloatText(description="X Axis Max:", value=45.0, layout=widgets.Layout(width="33%"))
y_axis_min = widgets.FloatText(description="Y Axis Min:", value=0.0, layout=widgets.Layout(width="33%"))
y_axis_max = widgets.FloatText(description="Y Axis Max:", value=10.0, layout=widgets.Layout(width="33%"))

# NEW: Grid Toggle
grid_checkbox = widgets.Checkbox(value=True, description="Show Grid")

# NEW: X-Tick Font Size (slider)
x_tick_font_size_widget = widgets.IntSlider(
    value=8,
    min=6,
    max=20,
    step=1,
    description="X Tick Font Size:",
    readout=True,
    layout=widgets.Layout(width="50%")
)

# NEW: Option to include r² in the plot title
include_r2_checkbox = widgets.Checkbox(value=True, description="Include R² in Plot Titles")

# NEW: Option to display the peak parameters table beneath each plot
include_peak_table_checkbox = widgets.Checkbox(value=True, description="Include Peak Parameters Table")

# NEW: Sub-array selection checkboxes
include_block_checkbox = widgets.Checkbox(value=True, description="Include Block Sub-Arrays (B1-B6)")
include_eye_checkbox = widgets.Checkbox(value=True, description="Include Eye-only Sub-Arrays (Ey1/Ey2 only)")

# NEW: Option to export figures. If unchecked, plots will only be shown.
export_figures_checkbox = widgets.Checkbox(value=True, description="Export Figures")

# NEW: Choose whether to use closest-to-mean plotting when multiple PSDs exist.
plot_closest_to_mean_checkbox = widgets.Checkbox(value=True, description="Use Closest-to-Mean Plotting (if multiple PSDs)")

# NEW: Options to export FOOOF results (as pickle and/or Excel)
export_fg_pickle_checkbox = widgets.Checkbox(value=True, description="Export FOOOF Results as Pickle")
export_fg_excel_checkbox = widgets.Checkbox(value=False, description="Export FOOOF Results as Excel")

# NEW: Option to display the FOOOF Group export as an Excel-like table in the notebook.
display_fg_export_table_checkbox = widgets.Checkbox(value=True, description="Display FOOOF Group Export Table")

# Figure format checkboxes
fmt_options = ["png", "svg", "jpeg", "ppt"]
format_checkboxes = [widgets.Checkbox(value=False, description=fmt.upper()) for fmt in fmt_options]

# Buttons
file_picker_button = widgets.Button(description="Select Excel Files", button_style="info", icon="folder")
directory_picker_button = widgets.Button(description="Select Output Directory", button_style="info", icon="folder")
run_button = widgets.Button(description="Process and Analyze", button_style="success", icon="check")
output = widgets.Output()

###################################
# 1. FILE SELECTION CALLBACK
###################################
def on_file_picker_button_click(b):
    global selected_files
    root = Tk()
    root.withdraw()  # Hide the root window
    file_paths = filedialog.askopenfilenames(filetypes=[("Excel files", "*.xlsx")])
    selected_files = list(file_paths)
    if selected_files:
        print(f"Selected files: {selected_files}")
    else:
        print("No files selected.")

file_picker_button.on_click(on_file_picker_button_click)

###################################
# 2. DIRECTORY SELECTION CALLBACK
###################################
def on_directory_picker_button_click(b):
    global export_directory
    root = Tk()
    root.withdraw()  # Hide the root window
    directory_path = filedialog.askdirectory()
    if directory_path:
        export_directory = directory_path
        print(f"Selected output directory: {export_directory}")
    else:
        print("No directory selected.")

directory_picker_button.on_click(on_directory_picker_button_click)

######################################################
# 3. MAIN BUTTON: PROCESS, ANALYZE & EXPORT/ PLOT FIGURES & RESULTS
######################################################
def on_run_button_click(b):
    global fg_dict
    with output:
        output.clear_output()

        # 1. Check if files are selected.
        if not selected_files:
            print("Error: Please select Excel files to process.")
            return

        # 2. If exporting figures or results, ensure an export directory is selected.
        if (export_figures_checkbox.value or export_fg_pickle_checkbox.value or export_fg_excel_checkbox.value) and not export_directory:
            print("Error: Please select an output directory for exporting.")
            return

        # 3. Gather figure export formats (if exporting figures is enabled).
        if export_figures_checkbox.value:
            chosen_formats = [c.description.lower() for c in format_checkboxes if c.value]
            if not chosen_formats:
                print("Warning: No figure format selected. Proceeding without exporting figures.")
        else:
            chosen_formats = []  # Only plot; no file export.
        
        # 4. Process Excel files with the sub-array filtering options.
        print(f"Processing files: {selected_files}")
        psd_data_dict_all = process_excel(
            selected_files, 
            include_block=include_block_checkbox.value, 
            include_eye=include_eye_checkbox.value
        )

        if psd_data_dict_all:
            # 5. Run FOOOF analysis.
            freq_range = [freq_range_min.value, freq_range_max.value]
            peak_width_limits = [peak_width_min.value, peak_width_max.value]

            print("Running FOOOF analysis...")
            fg_dict = run_fooof_analysis(
                psd_data_dict_all,
                freq_range=freq_range,
                amp_threshold=amp_threshold.value,
                r2_threshold=r2_threshold.value,
                max_peaks=max_peaks.value,
                fitting_mode=fitting_mode.value,
                peak_width_limits=peak_width_limits,
            )

            # 6. Plot the FOOOF results.
            print("\nPlotting FOOOF results...")
            for dict_key, fg_item in fg_dict.items():
                try:
                    print(f"\nPlotting results for: {dict_key}")
                    base_filename = dict_key
                    if plot_closest_to_mean_checkbox.value and len(fg_item) > 1:
                        plot_closest_to_mean(
                            fg_item, dict_key, export_directory, base_filename, chosen_formats,
                            x_min=x_axis_min.value,
                            x_max=x_axis_max.value,
                            y_min=y_axis_min.value,
                            y_max=y_axis_max.value,
                            show_grid=grid_checkbox.value,
                            x_tick_font_size=x_tick_font_size_widget.value,
                            include_r2=include_r2_checkbox.value,
                            include_peak_table=include_peak_table_checkbox.value
                        )
                    else:
                        plot_all_psds(
                            fg_item, dict_key, export_directory, base_filename, chosen_formats,
                            x_min=x_axis_min.value,
                            x_max=x_axis_max.value,
                            y_min=y_axis_min.value,
                            y_max=y_axis_max.value,
                            show_grid=grid_checkbox.value,
                            x_tick_font_size=x_tick_font_size_widget.value,
                            include_r2=include_r2_checkbox.value,
                            include_peak_table=include_peak_table_checkbox.value
                        )
                except Exception as e:
                    print(f"Error processing {dict_key}: {e}")

            # 7. Export FOOOF results if desired.
            if export_fg_pickle_checkbox.value:
                pickle_filename = os.path.join(export_directory, "fooof_groups.pkl")
                save_fooof_group(fg_dict, filename=pickle_filename)
            if export_fg_excel_checkbox.value:
                excel_filename = os.path.join(export_directory, "fooof_results.xlsx")
                export_fg_dict_excel(fg_dict, excel_filename)
            
            # 8. Optionally display the FOOOF Group export as an Excel-like table.
            if display_fg_export_table_checkbox.value:
                display_fg_dict_table(fg_dict)
        else:
            print("Processing failed. Please check the files and try again.")

run_button.on_click(on_run_button_click)

##############################
# 4. DISPLAY THE GUI
##############################
display(
    widgets.VBox(
        [
            widgets.Label(value="1) Select Excel Files:"),
            file_picker_button,
            widgets.Label(value="2) Select Output Directory (for exporting):"),
            directory_picker_button,
            widgets.Label(value="3) Configure Analysis Parameters:"),
            widgets.HBox([freq_range_min, freq_range_max]),
            amp_threshold,
            r2_threshold,
            max_peaks,
            fitting_mode,
            widgets.HBox([peak_width_min, peak_width_max]),
            widgets.Label(value="(Optional) Specify Common Plot Axis Ranges:"),
            widgets.HBox([x_axis_min, x_axis_max]),
            widgets.HBox([y_axis_min, y_axis_max]),
            widgets.Label(value="(Optional) Grid, X-Tick Font Size, R² and Peak Tables in Titles/Plots:"),
            widgets.HBox([grid_checkbox, x_tick_font_size_widget, include_r2_checkbox, include_peak_table_checkbox]),
            widgets.Label(value="4) Sub-Array Selection:"),
            widgets.HBox([include_block_checkbox, include_eye_checkbox]),
            widgets.Label(value="5) Figure Export Options:"),
            export_figures_checkbox,
            plot_closest_to_mean_checkbox,
            widgets.Label(value="6) Choose Figure Export Formats (if exporting figures is enabled):"),
            widgets.HBox(format_checkboxes),
            widgets.Label(value="7) Export FOOOF Results:"),
            widgets.HBox([export_fg_pickle_checkbox, export_fg_excel_checkbox]),
            widgets.Label(value="8) Display FOOOF Group Export as an Excel-like Table:"),
            display_fg_export_table_checkbox,
            run_button,
            output,
        ]
    )
)
