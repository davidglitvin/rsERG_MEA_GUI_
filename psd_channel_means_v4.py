import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math
import re

import ipywidgets as widgets
from ipywidgets import VBox, HBox  # Properly import VBox and HBox
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. PPT export will be disabled.")

# Check for pandas (for Excel export)
try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    print("Warning: pandas not installed. Excel export will be disabled.")

###############################################################################
# 1) Helper Functions
###############################################################################

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    """
    Identifies which PSD traces (epochs) should be excluded based on two criteria:
      1) Low-band outliers (e.g., trace in 1-3 Hz band > threshold * mean_psd).
      2) Repeated suprathreshold events in test bands >= test_band_count_threshold times.
    Returns (kept_indices, excluded_indices).
    """
    if test_band_count_threshold is None:
        test_band_count_threshold = len(test_bands) // 2

    mean_psd = np.mean(psd_array, axis=0)
    excluded_traces = []
    kept_traces = []

    # Indices for the specified low frequency band
    low_band_indices = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    # Indices for each test band
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0]
        for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # 1) Check low-frequency extreme outlier
        if np.any(trace[low_band_indices] > low_band_threshold * mean_psd[low_band_indices]):
            excluded_traces.append(i)
            continue

        # 2) Check repeated suprathreshold events in test bands
        suprathreshold_count = 0
        for indices in band_indices:
            if np.any(trace[indices] > test_band_threshold * mean_psd[indices]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces

def compute_channel_means(
    psds_dict,
    exclude=False,
    low_band=(1,3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    """
    Compute channel means before and after exclusion.
    Returns a dict {channel: (original_mean, new_mean)}.
    """
    channel_means = {}
    for channel, data in psds_dict.items():
        psd = data.get('psd', None)
        freqs = data.get('freqs', None)
        if psd is None or freqs is None:
            print(f"Channel '{channel}' is missing 'psd' or 'freqs'. Skipping.")
            continue

        original_mean = np.mean(psd, axis=0)
        if exclude:
            kept_traces, _ = exclude_traces(
                psd_array=psd,
                freqs=freqs,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold
            )
            if kept_traces:
                new_mean = np.mean(psd[kept_traces], axis=0)
            else:
                new_mean = np.zeros_like(original_mean)
        else:
            new_mean = original_mean.copy()
        
        channel_means[channel] = (original_mean, new_mean)
    
    return channel_means

def compute_group_mean(channel_means, selected_channels):
    """
    Compute the mean of 'new_mean' across selected channels.
    Returns a single array (the mean across those channels).
    """
    means = []
    for ch in selected_channels:
        if ch in channel_means:
            _, new_m = channel_means[ch]
            means.append(new_m)
    if means:
        return np.mean(means, axis=0)
    else:
        # If nothing is selected, return zeros of the correct shape
        if len(channel_means) > 0:
            first_val = next(iter(channel_means.values()))  # (orig, new)
            return np.zeros_like(first_val[0])
        else:
            return None

def plot_individual_channels(
    ax,
    channels,
    channel_means_dict,
    freqs,
    colors_for_channels,
    show_original_mean=True,
    show_new_mean=True,
    title="",
    axis_fs=10,
    legend_fs=8,
    tick_fs=8
):
    """
    Plots each channel in 'channels' onto Axes 'ax',
    using a distinct color from 'colors_for_channels'.
    If show_original_mean=True, it plots original mean in color (solid).
    If show_new_mean=True, it plots outlier-excluded mean in color (dashed).
    """
    for ch in channels:
        if ch not in channel_means_dict:
            continue
        orig_mean, new_mean = channel_means_dict[ch]
        color = colors_for_channels.get(ch, "blue")

        lbl_orig = f"{ch} Orig"
        lbl_new  = f"{ch} New"

        if show_original_mean:
            ax.plot(freqs, orig_mean, color=color, linestyle='-', label=lbl_orig)
        if show_new_mean:
            ax.plot(freqs, new_mean, color=color, linestyle='--', label=lbl_new)
    
    ax.set_title(title, fontsize=axis_fs)
    ax.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
    ax.tick_params(axis='both', labelsize=tick_fs)
    ax.legend(fontsize=legend_fs, loc="upper right")

def merge_block_keys(raw_dict):
    """
    Merge keys like 'block1:Ch3', 'block2:Ch3' => 'Ch3'.
    Returns a new dict {ChX: {'psd': array, 'freqs': array}}
    """
    block_pattern = re.compile(r'^block(\d+):Ch(\d+)$')
    merged_dict = {}

    for key, data in raw_dict.items():
        match = block_pattern.match(key)
        if match:
            ch_num = match.group(2)
            ch_key = f"Ch{ch_num}"
            if ch_key not in merged_dict:
                merged_dict[ch_key] = {'psd': [], 'freqs': None}
            merged_dict[ch_key]['psd'].append(data['psd'])
            if merged_dict[ch_key]['freqs'] is None:
                merged_dict[ch_key]['freqs'] = data['freqs']
        else:
            # If it's already "ChX" or something else, keep it
            if key not in merged_dict:
                merged_dict[key] = {'psd': [], 'freqs': None}
            merged_dict[key]['psd'].append(data['psd'])
            if merged_dict[key]['freqs'] is None:
                merged_dict[key]['freqs'] = data['freqs']

    # Concatenate
    for ch_key, ch_data in merged_dict.items():
        if len(ch_data['psd']) > 1:
            merged_dict[ch_key]['psd'] = np.concatenate(ch_data['psd'], axis=0)
        else:
            merged_dict[ch_key]['psd'] = ch_data['psd'][0]

    return merged_dict

def export_psd_data(psd_dict, export_path, export_format='pickle'):
    """
    Export PSD data as pickle or Excel.
    """
    if export_format == 'pickle':
        try:
            with open(export_path, 'wb') as f:
                pickle.dump(psd_dict, f)
            print(f"PSD data successfully exported to '{export_path}' as Pickle.")
        except Exception as e:
            print(f"Failed to export PSD data as Pickle: {e}")
    elif export_format == 'excel':
        if not HAS_PANDAS:
            print("pandas is not installed. Cannot export as Excel.")
            return
        try:
            # Convert psd_dict to a DataFrame with MultiIndex
            data = {}
            for ch, vals in psd_dict.items():
                freqs = vals['freqs']
                psd = vals['psd']
                df = pd.DataFrame(psd, columns=freqs)
                data[ch] = df
            # Save each channel as a separate sheet
            with pd.ExcelWriter(export_path) as writer:
                for ch, df in data.items():
                    df.to_excel(writer, sheet_name=ch[:31])  # Excel sheet names limited to 31 chars
            print(f"PSD data successfully exported to '{export_path}' as Excel.")
        except Exception as e:
            print(f"Failed to export PSD data as Excel: {e}")
    else:
        print(f"Unsupported export format: {export_format}")

def export_final_means(data_dict, fig, export_path, export_format='pickle'):
    """
    Export final means data as Pickle or Excel, and figures as image/PPT.
    """
    if export_format in ['pickle', 'excel']:
        if export_format == 'pickle':
            try:
                with open(export_path, 'wb') as f:
                    pickle.dump(data_dict, f)
                print(f"Final means data successfully exported to '{export_path}' as Pickle.")
            except Exception as e:
                print(f"Failed to export final means data as Pickle: {e}")
        elif export_format == 'excel':
            if not HAS_PANDAS:
                print("pandas is not installed. Cannot export as Excel.")
                return
            try:
                df = pd.DataFrame()
                for ch, (orig, new) in data_dict.items():
                    df[f"{ch}_original_mean"] = orig
                    df[f"{ch}_new_mean"] = new
                df.to_excel(export_path, index=False)
                print(f"Final means data successfully exported to '{export_path}' as Excel.")
            except Exception as e:
                print(f"Failed to export final means data as Excel: {e}")
    elif export_format in ['png','svg','jpeg']:
        if fig is None:
            print("No figure provided for export.")
            return
        try:
            fig.savefig(export_path, format=export_format, dpi=150)
            print(f"Final means figure successfully exported to '{export_path}' as {export_format.upper()}.")
        except Exception as e:
            print(f"Failed to export final means figure as {export_format.upper()}: {e}")
    elif export_format == 'ppt':
        if not HAS_PPTX:
            print("python-pptx is not installed. Cannot export as PPT.")
            return
        if fig is None:
            print("No figure provided for export.")
            return
        try:
            ppt = Presentation()
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # Blank slide
            tmp_png = export_path + ".png"
            fig.savefig(tmp_png, format='png', dpi=150)
            slide.shapes.add_picture(tmp_png, Inches(1), Inches(1), Inches(8), Inches(4.5))
            os.remove(tmp_png)
            ppt.save(export_path)
            print(f"Final means figure successfully exported to '{export_path}' as PPT.")
        except Exception as e:
            print(f"Failed to export final means figure as PPT: {e}")
    else:
        print(f"Unsupported export format: {export_format}")

###############################################################################
# 2) Main Interactive GUI
###############################################################################
def build_exportable_plot_psd_gui():
    """
    An interactive GUI to:
      - Drop channels in both merged and blocked plotting modes.
      - Export PSD data as Pickle or Excel before and after computing mean of means.
      - Export final means data and figures.
      - Ensure consistent widget sizes.
    """
    final_data_dict = {}  
    final_fig = None      
    current_figures = []  

    # 1. Load PSD Pickle
    merge_blocks_cb = widgets.Checkbox(
        value=False, 
        description='Merge blocks?',
        tooltip='If checked, blockN:ChX entries will be merged into ChX. Otherwise, each block remains separate.',
        layout=widgets.Layout(width='200px')
    )

    load_psd_button = widgets.Button(
        description='Load PSD Pickle',
        button_style='info',
        tooltip='Load a pickled PSD file',
        layout=widgets.Layout(width='150px')
    )
    psd_file_chooser = FileChooser(
        os.getcwd(),
        title='Select PSD Pickle File',
        select_default=False
    )
    psd_file_chooser.show_only_files = True
    psd_file_chooser.filter_pattern = ['*.pkl']

    # 2. Select Channels (for merged mode)
    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description='Channels:',
        layout=widgets.Layout(width='300px', height='200px')
    )
    # Add "Select All" and "Deselect All" buttons
    select_all_btn = widgets.Button(
        description='Select All',
        button_style='',
        tooltip='Select all channels',
        layout=widgets.Layout(width='100px')
    )
    deselect_all_btn = widgets.Button(
        description='Deselect All',
        button_style='',
        tooltip='Deselect all channels',
        layout=widgets.Layout(width='100px')
    )

    def on_select_all_clicked(b):
        channels_dropdown.value = tuple(channels_dropdown.options)

    def on_deselect_all_clicked(b):
        channels_dropdown.value = ()

    select_all_btn.on_click(on_select_all_clicked)
    deselect_all_btn.on_click(on_deselect_all_clicked)

    # 3. Plot toggles
    show_eye1_cb = widgets.Checkbox(value=True, description='Show Eye1 (Ch1-8)', layout=widgets.Layout(width='200px'))
    show_eye2_cb = widgets.Checkbox(value=True, description='Show Eye2 (Ch9-16)', layout=widgets.Layout(width='200px'))
    show_original_mean_cb = widgets.Checkbox(value=True, description='Show Original Means', layout=widgets.Layout(width='200px'))
    show_new_mean_cb = widgets.Checkbox(value=True, description='Show New Means', layout=widgets.Layout(width='200px'))
    show_kept_cb = widgets.Checkbox(value=True, description='Show Kept Traces', layout=widgets.Layout(width='200px'))
    show_excluded_cb = widgets.Checkbox(value=True, description='Show Excluded Traces', layout=widgets.Layout(width='200px'))

    # 4. Axis Ranges
    x_min_widget = widgets.FloatText(value=None, description='X min (Hz):', layout=widgets.Layout(width='220px'))
    x_max_widget = widgets.FloatText(value=None, description='X max (Hz):', layout=widgets.Layout(width='220px'))
    y_min_widget = widgets.FloatText(value=None, description='Y min:', layout=widgets.Layout(width='220px'))
    y_max_widget = widgets.FloatText(value=None, description='Y max:', layout=widgets.Layout(width='220px'))

    # 5. Font sizes
    title_fs_widget = widgets.IntText(value=10, description='Title FS:', layout=widgets.Layout(width='200px'))
    axis_fs_widget = widgets.IntText(value=8,  description='Axis FS:',  layout=widgets.Layout(width='200px'))
    legend_fs_widget = widgets.IntText(value=8,description='Legend FS:',layout=widgets.Layout(width='200px'))
    tick_fs_widget = widgets.IntText(value=8,  description='Tick FS:',  layout=widgets.Layout(width='200px'))

    # 6. Test bands + thresholds
    test_band_text = widgets.Textarea(
        value="(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)",
        description='Test Bands:',
        layout=widgets.Layout(width='250px', height='100px')
    )
    low_band_threshold_widget = widgets.FloatText(value=3.0, description='Low band thresh:', layout=widgets.Layout(width='250px'))
    test_band_threshold_widget = widgets.FloatText(value=10.0, description='Test band thresh:', layout=widgets.Layout(width='250px'))

    # 7. Plot Button
    plot_psd_button = widgets.Button(
        description='Plot PSDs',
        button_style='success',
        tooltip='Plot the selected PSDs',
        layout=widgets.Layout(width='150px')
    )

    # 8. Export controls (for subplots)
    export_filename_widget = widgets.Text(
        value='MyPSDExport',
        description='File Base Name:',
        layout=widgets.Layout(width='300px')
    )
    export_format_widget = widgets.Dropdown(
        options=['png', 'svg', 'jpeg', 'ppt'],
        value='png',
        description='Format:',
        layout=widgets.Layout(width='200px')
    )
    export_button = widgets.Button(
        description='Export Figures',
        button_style='warning',
        tooltip='Export the plotted figures',
        layout=widgets.Layout(width='150px')
    )

    # 9. Export PSD Data (Pickle/Excel)
    export_data_filename_widget = widgets.Text(
        value='PSDDataExport',
        description='PSD Data Base Name:',
        layout=widgets.Layout(width='300px')
    )
    export_data_format_widget = widgets.Dropdown(
        options=['pickle', 'excel'],
        value='pickle',
        description='Format:',
        layout=widgets.Layout(width='200px')
    )
    export_data_button = widgets.Button(
        description='Export PSD Data',
        button_style='primary',
        tooltip='Export the PSD data as Pickle or Excel',
        layout=widgets.Layout(width='150px')
    )

    # 10. Export Final Means controls
    final_export_basename_widget = widgets.Text(
        value='FinalMeansExport',
        description='Base Name:',
        layout=widgets.Layout(width='250px')
    )
    final_export_format_widget = widgets.Dropdown(
        options=['pickle','excel','png','svg','jpeg','ppt'],
        value='pickle',
        description='Format:',
        layout=widgets.Layout(width='240px')
    )
    final_export_button = widgets.Button(
        description='Export Final Means',
        button_style='info',
        tooltip='Export the final means data or figure',
        layout=widgets.Layout(width='200px')
    )
    final_export_output = widgets.Output()

    # 11. Output areas
    load_output_area = widgets.Output()
    plot_output_area = widgets.Output()
    export_output_area = widgets.Output()
    export_data_output = widgets.Output()

    ########################################################################
    # LOAD PSD BUTTON CALLBACK
    ########################################################################
    def on_load_psd_clicked(b):
        with load_output_area:
            clear_output()
            psd_path = psd_file_chooser.selected
            if not psd_path:
                print("Please select a PSD pickle file.")
                return
            if not os.path.isfile(psd_path):
                print(f"The file does not exist: {psd_path}")
                return
            try:
                with open(psd_path, 'rb') as f:
                    loaded_data = pickle.load(f)
                # Store loaded data in a global variable or closure variable
                build_exportable_plot_psd_gui.loaded_raw_data = loaded_data
                print(f"Successfully loaded PSD data from '{psd_path}'.")
                
                # If user merges, we'll do that inside the plot callback.
                print("Data loaded. Next, choose 'Merge blocks?' or not, then click 'Plot PSDs'.")
            except Exception as e:
                print(f"Failed to load PSD pickle: {e}")

    load_psd_button.on_click(on_load_psd_clicked)

    ########################################################################
    # PLOT PSD BUTTON CALLBACK
    ########################################################################
    def on_plot_psd_clicked(b):
        nonlocal final_data_dict, final_fig, merged_psd
        with plot_output_area:
            clear_output()
            current_figures.clear()
            final_data_dict.clear()
            final_fig = None

            # Access the loaded data
            loaded_raw_data = build_exportable_plot_psd_gui.loaded_raw_data if hasattr(build_exportable_plot_psd_gui, 'loaded_raw_data') else {}
            if not loaded_raw_data:
                print("No PSD data loaded. Please load a PSD file first.")
                return

            # Decide if we merge blocks or keep them separate
            do_merge = merge_blocks_cb.value
            print(f"Merging Blocks? {do_merge}")

            # Parse test bands
            test_bands_str = test_band_text.value.strip().replace(" ", "")
            test_bands_list = []
            if test_bands_str:
                pairs = test_bands_str.split(")")
                for p in pairs:
                    p = p.strip(",").strip("(").strip()
                    if not p:
                        continue
                    vals = p.split(",")
                    if len(vals) == 2:
                        try:
                            lowf = float(vals[0])
                            highf = float(vals[1])
                            test_bands_list.append((lowf, highf))
                        except:
                            pass
                if not test_bands_list:
                    test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]
            else:
                test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]

            low_band = (1,3)
            low_band_threshold = low_band_threshold_widget.value
            test_band_threshold = test_band_threshold_widget.value

            # Axis and font
            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value
            title_fs = title_fs_widget.value
            axis_fs  = axis_fs_widget.value
            legend_fs= legend_fs_widget.value
            tick_fs  = tick_fs_widget.value

            if do_merge:
                # Merged approach => one dictionary with 'Ch1', 'Ch2', etc.
                merged_psd = merge_block_keys(loaded_raw_data)

                # Sort channels by numeric ID
                def chnum(c):
                    return int(c.replace("Ch","")) if c.startswith("Ch") else 9999
                sorted_ch = sorted(merged_psd.keys(), key=chnum)

                # Populate channels dropdown so user can select
                channels_dropdown.options = sorted_ch

                # Based on user selection:
                selected_channels = list(channels_dropdown.value)
                if not selected_channels:
                    print("Please select at least one channel in the dropdown.")
                    return
                # Enforce up to 16 channels
                if len(selected_channels) > 16:
                    print("Please select up to 16 channels in merged mode.")
                    return

                # Eye1 = ch1-8, Eye2 = ch9-16
                def get_ch_number(ch_name):
                    return int(ch_name.replace('Ch','')) if ch_name.startswith('Ch') else None
                eye1_channels = [ch for ch in selected_channels if (get_ch_number(ch) and 1 <= get_ch_number(ch) <= 8)]
                eye2_channels = [ch for ch in selected_channels if (get_ch_number(ch) and 9 <= get_ch_number(ch) <= 16)]

                # Compute means
                channel_means_before = compute_channel_means(
                    merged_psd,
                    exclude=False
                )
                channel_means_after = compute_channel_means(
                    merged_psd,
                    exclude=True,
                    low_band=low_band,
                    low_band_threshold=low_band_threshold,
                    test_bands=test_bands_list,
                    test_band_threshold=test_band_threshold
                )

                import matplotlib.cm as cm
                color_cycle = cm.get_cmap('tab20').colors
                colors_for_channels = {}
                for i, ch in enumerate(selected_channels):
                    colors_for_channels[ch] = color_cycle[i % len(color_cycle)]

                # We'll create up to 2 subplots (one for Eye1, one for Eye2)
                n_sub = 0
                if show_eye1_cb.value and eye1_channels:
                    n_sub += 1
                if show_eye2_cb.value and eye2_channels:
                    n_sub += 1

                if n_sub == 0:
                    print("No Eye1/Eye2 channels selected to plot in merged mode.")
                    return

                fig, axes = plt.subplots(1, n_sub, figsize=(6*n_sub, 5))
                if n_sub == 1:
                    axes = [axes]

                idx = 0

                # Eye1
                if show_eye1_cb.value and eye1_channels:
                    ax_e1 = axes[idx]
                    freqs = merged_psd[eye1_channels[0]]['freqs']
                    plot_individual_channels(
                        ax=ax_e1,
                        channels=eye1_channels,
                        channel_means_dict=channel_means_before, 
                        freqs=freqs,
                        colors_for_channels=colors_for_channels,
                        show_original_mean=show_original_mean_cb.value,
                        show_new_mean=show_new_mean_cb.value,
                        title='Eye1 (Merged)',
                        axis_fs=axis_fs,
                        legend_fs=legend_fs,
                        tick_fs=tick_fs
                    )
                    # Add axis limits
                    if x_min is not None: ax_e1.set_xlim(left=x_min)
                    if x_max is not None: ax_e1.set_xlim(right=x_max)
                    if y_min is not None: ax_e1.set_ylim(bottom=y_min)
                    if y_max is not None: ax_e1.set_ylim(top=y_max)

                    # Plot group mean of Eye1
                    group_mean_e1 = compute_group_mean(channel_means_after, eye1_channels)
                    if group_mean_e1 is not None:
                        ax_e1.plot(freqs, group_mean_e1, color='black', linewidth=2, label='Eye1 Group Mean')
                        ax_e1.legend(fontsize=legend_fs)
                    idx += 1

                # Eye2
                if show_eye2_cb.value and eye2_channels:
                    ax_e2 = axes[idx]
                    freqs = merged_psd[eye2_channels[0]]['freqs']
                    plot_individual_channels(
                        ax=ax_e2,
                        channels=eye2_channels,
                        channel_means_dict=channel_means_before,
                        freqs=freqs,
                        colors_for_channels=colors_for_channels,
                        show_original_mean=show_original_mean_cb.value,
                        show_new_mean=show_new_mean_cb.value,
                        title='Eye2 (Merged)',
                        axis_fs=axis_fs,
                        legend_fs=legend_fs,
                        tick_fs=tick_fs
                    )
                    if x_min is not None: ax_e2.set_xlim(left=x_min)
                    if x_max is not None: ax_e2.set_xlim(right=x_max)
                    if y_min is not None: ax_e2.set_ylim(bottom=y_min)
                    if y_max is not None: ax_e2.set_ylim(top=y_max)

                    # group mean
                    group_mean_e2 = compute_group_mean(channel_means_after, eye2_channels)
                    if group_mean_e2 is not None:
                        ax_e2.plot(freqs, group_mean_e2, color='black', linewidth=2, label='Eye2 Group Mean')
                        ax_e2.legend(fontsize=legend_fs)

                plt.tight_layout()
                plt.show()
                current_figures.append(fig)
                print("Done plotting merged blocks.")

                # Prepare final means data
                final_data_dict = {}
                if show_new_mean_cb.value:
                    for ch in selected_channels:
                        if ch in channel_means_after:
                            final_data_dict[ch] = channel_means_after[ch]
            else:
                # ----------------------------------------------------------------
                # NON-MERGED MODE => block1:Ch1, block2:Ch1, etc.
                # We want for each block x each Eye => 1 subplot
                # So 6 blocks => up to 12 subplots (6 for Eye1, 6 for Eye2),
                # each with 8 lines + group mean line.
                # ----------------------------------------------------------------
                # Identify all blocks
                block_pattern = re.compile(r'^block(\d+):Ch(\d+)$')
                block_map = {}  # {blockNum: { key -> data }}
                for key, data in loaded_raw_data.items():
                    m = block_pattern.match(key)
                    if m:
                        bnum = int(m.group(1))
                        if bnum not in block_map:
                            block_map[bnum] = {}
                        block_map[bnum][key] = data
                    else:
                        # "ChX" with no block => we can treat as "block0" or separate
                        if 0 not in block_map:
                            block_map[0] = {}
                        block_map[0][key] = data

                # Sort block numbers
                all_blocks = sorted(block_map.keys())
                if not all_blocks:
                    print("No block-based keys found.")
                    return

                # We'll create subplots => #rows = len(all_blocks), #cols = 2 (Eye1 & Eye2)
                n_blocks = len(all_blocks)
                fig, axes = plt.subplots(n_blocks, 2, figsize=(12, 4*n_blocks), squeeze=False)
                plt.subplots_adjust(hspace=0.5, wspace=0.4)

                # We'll do outlier detection + means for each block individually
                # Then inside each block, we separate Eye1 vs Eye2 => "blockX:Ch1..Ch8", etc.
                import matplotlib.cm as cm
                color_cycle = cm.get_cmap('tab20').colors

                row_idx = 0
                for bnum in all_blocks:
                    # Build a "local" dictionary of {ChName: {'psd', 'freqs'}} for this block
                    block_dict = {}
                    for k, v in block_map[bnum].items():
                        block_dict[k] = v

                    # Compute means
                    channel_means_before = compute_channel_means(
                        block_dict,
                        exclude=False
                    )
                    channel_means_after = compute_channel_means(
                        block_dict,
                        exclude=True,
                        low_band=low_band,
                        low_band_threshold=low_band_threshold,
                        test_bands=test_bands_list,
                        test_band_threshold=test_band_threshold
                    )

                    # Identify which channels are Eye1 vs Eye2 within this block
                    def parse_block_ch(key):
                        # key looks like "block5:Ch3"
                        m = block_pattern.match(key)
                        if not m:
                            return None
                        chnum = int(m.group(2))
                        return chnum

                    eye1_keys = []
                    eye2_keys = []
                    for k in block_dict.keys():
                        chnum = parse_block_ch(k)
                        if chnum is not None:
                            if 1 <= chnum <= 8:
                                eye1_keys.append(k)
                            elif 9 <= chnum <= 16:
                                eye2_keys.append(k)

                    # Make color mapping
                    all_block_channels = sorted(block_dict.keys())
                    cdict = {}
                    for i, ckey in enumerate(all_block_channels):
                        cdict[ckey] = color_cycle[i % len(color_cycle)]

                    # Eye1 subplot
                    ax_eye1 = axes[row_idx, 0]
                    if show_eye1_cb.value and len(eye1_keys) > 0:
                        freqs_e1 = block_dict[eye1_keys[0]]['freqs']
                        plot_individual_channels(
                            ax=ax_eye1,
                            channels=eye1_keys,
                            channel_means_dict=channel_means_before,
                            freqs=freqs_e1,
                            colors_for_channels=cdict,
                            show_original_mean=show_original_mean_cb.value,
                            show_new_mean=show_new_mean_cb.value,
                            title=f"Block {bnum}: Eye1",
                            axis_fs=axis_fs,
                            legend_fs=legend_fs,
                            tick_fs=tick_fs
                        )
                        # group mean
                        group_mean_e1 = compute_group_mean(channel_means_after, eye1_keys)
                        if group_mean_e1 is not None:
                            ax_eye1.plot(freqs_e1, group_mean_e1, color='black', linewidth=2, label='Group Mean')
                            ax_eye1.legend(fontsize=legend_fs)

                        # Set axis limits
                        if x_min is not None: ax_eye1.set_xlim(left=x_min)
                        if x_max is not None: ax_eye1.set_xlim(right=x_max)
                        if y_min is not None: ax_eye1.set_ylim(bottom=y_min)
                        if y_max is not None: ax_eye1.set_ylim(top=y_max)
                    else:
                        ax_eye1.set_title(f"Block {bnum}: Eye1 (No Data)", fontsize=axis_fs)

                    # Eye2 subplot
                    ax_eye2 = axes[row_idx, 1]
                    if show_eye2_cb.value and len(eye2_keys) > 0:
                        freqs_e2 = block_dict[eye2_keys[0]]['freqs']
                        plot_individual_channels(
                            ax=ax_eye2,
                            channels=eye2_keys,
                            channel_means_dict=channel_means_before,
                            freqs=freqs_e2,
                            colors_for_channels=cdict,
                            show_original_mean=show_original_mean_cb.value,
                            show_new_mean=show_new_mean_cb.value,
                            title=f"Block {bnum}: Eye2",
                            axis_fs=axis_fs,
                            legend_fs=legend_fs,
                            tick_fs=tick_fs
                        )
                        # group mean
                        group_mean_e2 = compute_group_mean(channel_means_after, eye2_keys)
                        if group_mean_e2 is not None:
                            ax_eye2.plot(freqs_e2, group_mean_e2, color='black', linewidth=2, label='Group Mean')
                            ax_eye2.legend(fontsize=legend_fs)

                        # Set axis limits
                        if x_min is not None: ax_eye2.set_xlim(left=x_min)
                        if x_max is not None: ax_eye2.set_xlim(right=x_max)
                        if y_min is not None: ax_eye2.set_ylim(bottom=y_min)
                        if y_max is not None: ax_eye2.set_ylim(top=y_max)
                    else:
                        ax_eye2.set_title(f"Block {bnum}: Eye2 (No Data)", fontsize=axis_fs)

                    row_idx += 1

                plt.tight_layout()
                plt.show()
                current_figures.append(fig)
                print("Done plotting separate blocks.")

                # Prepare final means data
                final_data_dict = {}
                if show_new_mean_cb.value:
                    # Collect means from all channels across blocks
                    for bnum in all_blocks:
                        for ch_key in merged_psd.keys() if do_merge else block_map[bnum].keys():
                            if do_merge:
                                ch = ch_key
                            else:
                                # Extract channel name from key like 'block1:Ch3'
                                ch = ch_key.split(':')[1]
                            if ch in channel_means_after:
                                final_data_dict[ch] = channel_means_after[ch]
                # Else, final_data_dict remains empty

    plot_psd_button.on_click(on_plot_psd_clicked)

    ########################################################################
    # EXPORT BUTTON CALLBACK (for subplots)
    ########################################################################
    def on_export_button_clicked(b):
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No figures to export. Please plot PSDs first.")
                return

            fname_base = export_filename_widget.value.strip()
            export_fmt = export_format_widget.value.lower()
            if not fname_base:
                print("Please provide a valid base filename.")
                return

            if export_fmt == 'ppt' and not HAS_PPTX:
                print("python-pptx not installed. Cannot export to PPT.")
                return

            print(f"Exporting {len(current_figures)} figure(s) as {export_fmt}...")

            if export_fmt == 'ppt':
                prs = Presentation()
                for idx, fig in enumerate(current_figures, start=1):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
                    image_path = f"{fname_base}_{idx}.png"
                    fig.savefig(image_path, format='png', dpi=150)
                    left = top = Inches(1)
                    slide.shapes.add_picture(image_path, left, top, Inches(8), Inches(4.5))
                    os.remove(image_path)
                ppt_filename = f"{fname_base}.pptx"
                prs.save(ppt_filename)
                print(f"Exported all subplots into PPT: {ppt_filename}")
            else:
                for idx, fig in enumerate(current_figures, start=1):
                    save_name = f"{fname_base}_{idx}.{export_fmt}"
                    fig.savefig(save_name, format=export_fmt, dpi=150)
                    print(f"  Saved => {save_name}")

            print("Done exporting subplots.")

    export_button.on_click(on_export_button_clicked)

    ########################################################################
    # EXPORT PSD DATA BUTTON CALLBACK
    ########################################################################
    def on_export_data_button_clicked(b):
        with export_data_output:
            clear_output()
            # Access the loaded data
            loaded_raw_data = build_exportable_plot_psd_gui.loaded_raw_data if hasattr(build_exportable_plot_psd_gui, 'loaded_raw_data') else {}
            if not loaded_raw_data:
                print("No PSD data loaded. Please load a PSD file first.")
                return

            fname_base = export_data_filename_widget.value.strip()
            export_fmt = export_data_format_widget.value.lower()
            if not fname_base:
                print("Please provide a valid base filename for PSD data export.")
                return

            # Decide whether to export merged or blocked data
            do_merge = merge_blocks_cb.value
            if do_merge:
                # Merged data is stored in merged_psd
                if not merged_psd:
                    print("No merged PSD data available. Please plot PSDs first.")
                    return
                psd_to_export = merged_psd
            else:
                psd_to_export = loaded_raw_data

            # Define export path
            export_path = os.path.join(os.getcwd(), f"{fname_base}.{export_fmt}")

            # Perform export
            export_psd_data(psd_to_export, export_path, export_format=export_fmt)

    export_data_button.on_click(on_export_data_button_clicked)

    ########################################################################
    # EXPORT FINAL MEANS BUTTON CALLBACK
    ########################################################################
    def on_final_export_clicked(b):
        nonlocal final_data_dict, final_fig
        with final_export_output:
            clear_output()
            # Access the latest figure (if any)
            fig = current_figures[-1] if current_figures else None

            base_name = final_export_basename_widget.value.strip()
            chosen_fmt = final_export_format_widget.value
            if not base_name:
                print("Please provide a valid base name.")
                return

            if not final_data_dict and fig is None:
                print("No final data or final figure found. Did you plot and select channels?")
                return

            print(f"Exporting final means data/plot as {chosen_fmt}...")

            if chosen_fmt in ['pickle', 'excel']:
                export_path = os.path.join(os.getcwd(), f"{base_name}_final_means.{chosen_fmt}")
                export_final_means(final_data_dict, None, export_path, export_format=chosen_fmt)
            elif chosen_fmt in ['png','svg','jpeg','ppt']:
                # For figure exports, assume the user wants to export the latest figure
                if not current_figures:
                    print("No figures available to export.")
                    return
                # Take the last plotted figure
                fig = current_figures[-1]
                if chosen_fmt == 'ppt' and not HAS_PPTX:
                    print("python-pptx not installed. Cannot export to PPT.")
                    return
                export_path = os.path.join(os.getcwd(), f"{base_name}_final_means.{chosen_fmt}")
                export_final_means(final_data_dict, fig, export_path, export_format=chosen_fmt)
            else:
                print(f"Unsupported export format: {chosen_fmt}")

    final_export_button.on_click(on_final_export_clicked)

    ########################################################################
    # LAYOUT
    ########################################################################
    # Load section
    load_controls = HBox([merge_blocks_cb, load_psd_button, psd_file_chooser], layout=widgets.Layout(width='100%'))
    load_section = VBox([
        widgets.HTML("<h3>Load PSD Data</h3>"),
        load_controls,
        load_output_area
    ], layout=widgets.Layout(width='100%'))

    # Plot section
    plot_controls_top = HBox([
        show_eye1_cb, 
        show_eye2_cb, 
        show_original_mean_cb, 
        show_new_mean_cb,
        show_kept_cb,
        show_excluded_cb
    ], layout=widgets.Layout(width='100%'))
    thresholds_box = HBox([
        low_band_threshold_widget, 
        test_band_threshold_widget
    ], layout=widgets.Layout(width='100%'))
    axis_box = HBox([x_min_widget, x_max_widget, y_min_widget, y_max_widget], layout=widgets.Layout(width='100%'))
    font_box = HBox([title_fs_widget, axis_fs_widget, legend_fs_widget, tick_fs_widget], layout=widgets.Layout(width='100%'))
    channel_selection_box = HBox([
        channels_dropdown,
        VBox([
            select_all_btn,
            deselect_all_btn
        ], layout=widgets.Layout(width='120px', height='80px'))
    ], layout=widgets.Layout(width='100%'))
    plot_section = VBox([
        widgets.HTML("<h3>Plot Options</h3>"),
        plot_controls_top,
        widgets.Label("Test Bands: (e.g. (7,9),(9,11),...)"),
        test_band_text,
        thresholds_box,
        axis_box,
        font_box,
        widgets.HTML("<b>Channel Selection (Merged Mode):</b>"),
        channel_selection_box,
        plot_psd_button,
        plot_output_area
    ], layout=widgets.Layout(width='100%'))

    # Export subplots section
    export_subplots_section = VBox([
        widgets.HTML("<h3>Export Plotted Figures</h3>"),
        HBox([export_filename_widget, export_format_widget, export_button], layout=widgets.Layout(width='100%')),
        export_output_area
    ], layout=widgets.Layout(width='100%'))

    # Export PSD data section
    export_data_section = VBox([
        widgets.HTML("<h3>Export PSD Data</h3>"),
        HBox([export_data_filename_widget, export_data_format_widget, export_data_button], layout=widgets.Layout(width='100%')),
        export_data_output
    ], layout=widgets.Layout(width='100%'))

    # Export final means section
    export_final_section = VBox([
        widgets.HTML("<h3>Export Final Means Data/Figure</h3>"),
        HBox([final_export_basename_widget, final_export_format_widget, final_export_button], layout=widgets.Layout(width='100%')),
        final_export_output
    ], layout=widgets.Layout(width='100%'))

    # Final UI layout
    ui = VBox([
        widgets.HTML("<h2>PSD Computation and Export GUI</h2>"),
        load_section,
        plot_section,
        export_subplots_section,
        export_data_section,
        export_final_section
    ], layout=widgets.Layout(width='100%'))

    display(ui)

    ########################################################################
    # Data containers
    ########################################################################
    build_exportable_plot_psd_gui.loaded_raw_data = {}  # the PSD dictionary exactly as loaded from pickle
    merged_psd = {}       # for merged data (if user wants it)

    ########################################################################
    # EXPORT PSD DATA BUTTON CALLBACK
    ########################################################################
    def on_export_data_button_clicked(b):
        with export_data_output:
            clear_output()
            # Access the loaded data
            loaded_raw_data = build_exportable_plot_psd_gui.loaded_raw_data
            if not loaded_raw_data:
                print("No PSD data loaded. Please load a PSD file first.")
                return

            fname_base = export_data_filename_widget.value.strip()
            export_fmt = export_data_format_widget.value.lower()
            if not fname_base:
                print("Please provide a valid base filename for PSD data export.")
                return

            # Decide whether to export merged or blocked data
            do_merge = merge_blocks_cb.value
            if do_merge:
                # Merged data is stored in merged_psd
                if not merged_psd:
                    print("No merged PSD data available. Please plot PSDs first.")
                    return
                psd_to_export = merged_psd
            else:
                psd_to_export = loaded_raw_data

            # Define export path
            export_path = os.path.join(os.getcwd(), f"{fname_base}.{export_fmt}")

            # Perform export
            export_psd_data(psd_to_export, export_path, export_format=export_fmt)

    export_data_button.on_click(on_export_data_button_clicked)

###############################################################################
# Run the GUI
###############################################################################
build_exportable_plot_psd_gui()

