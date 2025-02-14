import os
import pickle
import numpy as np
import mne
from scipy import signal
from tqdm.notebook import tqdm

import ipywidgets as widgets
from ipywidgets import (
    FloatText,
    FloatSlider,
    IntProgress,
    Button,
    VBox,
    HBox,
    Label,
    Output,
    SelectMultiple,
    ToggleButtons
)
from IPython.display import display, clear_output
import matplotlib.pyplot as plt

from ipyfilechooser import FileChooser
from IPython import get_ipython

# For PPT export:
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

##############################################################################
# 1) WIDGETS & GLOBALS
##############################################################################

# File chooser for .fif input
input_fif_chooser = FileChooser(
    os.getcwd(),
    title='Select Filtered .fif File',
    select_default=False
)
input_fif_chooser.show_only_files = True
input_fif_chooser.filter_pattern = ['*.fif', '*.fif.gz']

# Segment Toggle: Whether to attempt to split into up to 6 blocks (each 120 epochs)
segment_toggle = ToggleButtons(
    options=['No', 'Yes'],
    value='No',
    description='Segment into 6 blocks?',
    style={'description_width': 'initial'}
)

# Window length and overlap
window_length_widget = FloatText(
    value=2.0,
    description='Window (s):',
    layout=widgets.Layout(width='250px')
)
overlap_widget = FloatSlider(
    value=50.0,
    min=0.0,
    max=100.0,
    step=1.0,
    description='Overlap (%):',
    continuous_update=False,
    layout=widgets.Layout(width='300px')
)

# Output pickle chooser
output_pickle_chooser = FileChooser(
    os.getcwd(),
    title='Select Output Pickle File',
    select_default=False
)
output_pickle_chooser.show_only_files = False
output_pickle_chooser.filter_pattern = ['*.pkl']
output_pickle_chooser.default_filename = 'psd_results.pkl'

# Compute PSD button and progress
compute_psd_button = Button(
    description='Compute PSD',
    button_style='success',
    tooltip='Compute PSD on loaded data and save results'
)
psd_progress = IntProgress(
    value=0,
    min=0,
    max=100,
    step=1,
    description='Progress:',
    bar_style='info',
    orientation='horizontal',
    layout=widgets.Layout(width='400px')
)

# Output area
psd_output_area = Output()

# Global dictionary to store PSD results
psd_results = {}  # Will be { "block1": { "Ch1": {"freqs", "psd"}, ...}, "block2": ...}

# For remembering the last plot figure (for exporting)
last_psd_figure = None

# Channel selection & plot mode
plot_channel_widget = SelectMultiple(
    options=[],
    description='Select Channels:',
    layout=widgets.Layout(width='350px', height='220px')
)
plot_mode_widget = ToggleButtons(
    options=[('Mean & +Std', 'mean_std'), ('Individual + Mean', 'individual_mean')],
    description='Plot Mode:',
    value='mean_std',
    button_style='info',
    style={'description_width': 'initial'},
    layout=widgets.Layout(width='300px')
)

# Axis Range
x_min_widget = FloatText(value=0.0, description='X min (Hz):', layout=widgets.Layout(width='230px'))
x_max_widget = FloatText(value=45.0, description='X max (Hz):', layout=widgets.Layout(width='230px'))
y_min_widget = FloatText(value=0.0, description='Y min:', layout=widgets.Layout(width='230px'))
y_max_widget = FloatText(value=1.0, description='Y max:', layout=widgets.Layout(width='230px'))

# Plot Button
plot_psd_button = Button(
    description='Plot PSD',
    button_style='info',
    tooltip='Plot PSD side by side for each block'
)

##############################################################################
# 1B) FIGURE EXPORT WIDGETS
##############################################################################
export_format_toggle = ToggleButtons(
    options=[('PPTX','pptx'), ('PNG','png'), ('JPEG','jpeg'), ('SVG','svg')],
    value='png',
    description='Export Format:',
    style={'description_width': 'initial'}
)

figure_export_chooser = FileChooser(
    os.getcwd(),
    title='Select directory or type a file name for export',
    select_default=False
)
# We allow both directory and file selection, but typically user can type 'myplot'
# and we will auto-append the extension.

save_figure_button = Button(
    description='Save Figure',
    button_style='warning',
    tooltip='Save the last plotted figure in the selected format'
)

save_figure_output = Output()

##############################################################################
# 2) HELPER FUNCTIONS
##############################################################################

def load_epochs_from_fif(filepath):
    """Load .fif or .fif.gz epochs and return mne.Epochs object."""
    if not os.path.isfile(filepath):
        raise FileNotFoundError(f"File does not exist: {filepath}")
    epochs_loaded = mne.read_epochs(filepath, preload=True)
    if not isinstance(epochs_loaded, mne.Epochs):
        data_array = epochs_loaded.get_data()
        info = epochs_loaded.info
        epochs_obj = mne.EpochsArray(data_array, info)
        return epochs_obj
    return epochs_loaded

def compute_welch_psd_for_epochs(epochs, window_s=2.0, overlap_pct=50.0, fmin=0.0, fmax=45.0):
    """
    Compute PSD for all channels in a single mne.Epochs object using Welch's method.
    Returns dict {channel_name: {"freqs": <1D array>, "psd": <2D array (n_epochs, n_freqs)>}}
    """
    psd_dict = {}
    sfreq = epochs.info['sfreq']
    nperseg = int(window_s * sfreq)
    noverlap = int(nperseg * (overlap_pct / 100.0))

    data = epochs.get_data()  # shape: (n_epochs, n_channels, n_times)
    n_epochs, n_channels, n_times = data.shape

    if nperseg > n_times:
        raise ValueError(f"nperseg ({nperseg}) cannot exceed each epoch’s length in samples ({n_times}).")

    # Precompute window
    window_fn = signal.hamming(nperseg, sym=False)

    for ch_idx, ch_name in enumerate(epochs.ch_names):
        ch_data = data[:, ch_idx, :]  # shape => (n_epochs, n_times)

        freqs, psd_all = signal.welch(
            ch_data,
            fs=sfreq,
            window=window_fn,
            nperseg=nperseg,
            noverlap=noverlap,
            scaling='density',
            axis=1
        )
        # Filter freq range
        mask = (freqs >= fmin) & (freqs <= fmax)
        freqs_f = freqs[mask]
        psd_f = psd_all[:, mask]

        # if no freqs remain, skip
        if freqs_f.size == 0:
            psd_dict[ch_name] = {"freqs": None, "psd": None}
        else:
            psd_dict[ch_name] = {"freqs": freqs_f, "psd": psd_f}
    return psd_dict

##############################################################################
# 3) COMPUTE PSD CALLBACK
##############################################################################

def on_compute_psd_clicked(b):
    global psd_results
    psd_results = {}
    with psd_output_area:
        clear_output()
        
        # 1. Validate & load .fif
        fif_path = input_fif_chooser.selected
        if not fif_path:
            print("ERROR: No .fif file selected.")
            return
        if not os.path.isfile(fif_path):
            print(f"ERROR: File not found => {fif_path}")
            return
        
        try:
            print(f"Loading epochs from: {fif_path}")
            epochs = load_epochs_from_fif(fif_path)
            print(f"Loaded. n_epochs={len(epochs)}, n_channels={len(epochs.ch_names)}")
        except Exception as e:
            print(f"ERROR loading epochs: {e}")
            return
        
        # 2. Check segment toggle
        segment_choice = segment_toggle.value  # "Yes" or "No"
        
        # 3. Parse PSD parameters
        window_s = window_length_widget.value
        overlap_pct = overlap_widget.value
        if window_s <= 0:
            print("ERROR: Window length must be positive.")
            return
        if not (0 <= overlap_pct <= 100):
            print("ERROR: Overlap must be in [0, 100].")
            return
        
        # 4. Output pickle file
        out_pkl = output_pickle_chooser.selected
        if not out_pkl or not out_pkl.endswith('.pkl'):
            print("ERROR: Output file must be .pkl.")
            return
        
        # 5. Segmentation logic
        n_total_epochs = len(epochs)
        if segment_choice == "Yes":
            # If at least 720 epochs, do 6 full blocks of 120 (plus leftover).
            # Otherwise, do as many 120-epoch blocks as possible (plus leftover if any).
            # If fewer than 120, we just do one block with all epochs.
            if n_total_epochs >= 720:
                # 6 full blocks + leftover if any
                print("Segmenting into 6 blocks of 120 epochs each (plus leftover if any).")
                n_blocks = 6
                block_size = 120
                for i in range(n_blocks):
                    start_ep = i * block_size
                    end_ep = (i + 1) * block_size
                    block_epochs = epochs[start_ep:end_ep]  # subselect
                    block_name = f"block{i+1}"
                    print(f"  Computing PSD for {block_name} => epochs {start_ep}..{end_ep-1}")
                    try:
                        block_psd = compute_welch_psd_for_epochs(
                            block_epochs,
                            window_s=window_s,
                            overlap_pct=overlap_pct,
                            fmin=0.0,
                            fmax=45.0
                        )
                        psd_results[block_name] = block_psd
                    except Exception as e:
                        print(f"ERROR computing PSD for block {block_name}: {e}")
                        psd_results[block_name] = {}
                leftover = n_total_epochs - (n_blocks * block_size)
                if leftover > 0:
                    # One leftover block
                    start_ep = n_blocks * block_size
                    end_ep = n_total_epochs
                    block_name = f"block{n_blocks+1}_leftover"
                    print(f"  Computing PSD for leftover block => epochs {start_ep}..{end_ep-1}")
                    try:
                        leftover_epochs = epochs[start_ep:end_ep]
                        block_psd = compute_welch_psd_for_epochs(
                            leftover_epochs,
                            window_s=window_s,
                            overlap_pct=overlap_pct,
                            fmin=0.0,
                            fmax=45.0
                        )
                        psd_results[block_name] = block_psd
                    except Exception as e:
                        print(f"ERROR computing PSD for leftover block: {e}")
                        psd_results[block_name] = {}
                        
            else:
                # Not enough for 6 full blocks
                if n_total_epochs < 120:
                    # Just do one block for everything
                    print(f"Fewer than 120 epochs ({n_total_epochs}), treating entire data as single block.")
                    try:
                        single_psd = compute_welch_psd_for_epochs(
                            epochs,
                            window_s=window_s,
                            overlap_pct=overlap_pct,
                            fmin=0.0,
                            fmax=45.0
                        )
                        psd_results["block1"] = single_psd
                    except Exception as e:
                        print(f"ERROR computing PSD for single block: {e}")
                        return
                else:
                    # Some number of 120-epoch blocks plus leftover
                    block_size = 120
                    n_possible = n_total_epochs // block_size  # full blocks
                    print(f"Segmenting into {n_possible} x 120-epoch blocks (plus leftover if any).")
                    for i in range(n_possible):
                        start_ep = i * block_size
                        end_ep = (i + 1) * block_size
                        block_name = f"block{i+1}"
                        print(f"  Computing PSD for {block_name} => epochs {start_ep}..{end_ep-1}")
                        try:
                            block_epochs = epochs[start_ep:end_ep]
                            block_psd = compute_welch_psd_for_epochs(
                                block_epochs,
                                window_s=window_s,
                                overlap_pct=overlap_pct,
                                fmin=0.0,
                                fmax=45.0
                            )
                            psd_results[block_name] = block_psd
                        except Exception as e:
                            print(f"ERROR computing PSD for {block_name}: {e}")
                            psd_results[block_name] = {}
                    leftover = n_total_epochs - (n_possible * block_size)
                    if leftover > 0:
                        # leftover block
                        start_ep = n_possible * block_size
                        end_ep = n_total_epochs
                        block_name = f"block{n_possible+1}_leftover"
                        print(f"  Computing PSD for leftover block => epochs {start_ep}..{end_ep-1}")
                        try:
                            leftover_epochs = epochs[start_ep:end_ep]
                            block_psd = compute_welch_psd_for_epochs(
                                leftover_epochs,
                                window_s=window_s,
                                overlap_pct=overlap_pct,
                                fmin=0.0,
                                fmax=45.0
                            )
                            psd_results[block_name] = block_psd
                        except Exception as e:
                            print(f"ERROR computing PSD for leftover block: {e}")
                            psd_results[block_name] = {}
        else:
            # No segmentation => treat entire epochs as one block
            print("No segmentation => single block with all epochs.")
            try:
                single_psd = compute_welch_psd_for_epochs(
                    epochs,
                    window_s=window_s,
                    overlap_pct=overlap_pct,
                    fmin=0.0,
                    fmax=45.0
                )
                psd_results["block1"] = single_psd
            except Exception as e:
                print(f"ERROR computing PSD for single block: {e}")
                return
        
        # 6. Check the structure of psd_results
        if not psd_results:
            print("ERROR: No PSD results were computed.")
            return
        
        # 7. We gather all channels from the FIRST block to populate the channel widget
        first_block_name = list(psd_results.keys())[0]
        block_data = psd_results[first_block_name]  # e.g. {ch: {freqs, psd}}
        # only keep valid channels
        valid_channels = [ch for ch, dd in block_data.items() if dd["freqs"] is not None]
        plot_channel_widget.options = valid_channels
        print(f"\nPopulated channel widget with {len(valid_channels)} valid channels from {first_block_name}.\n")
        
        # 8. Save PSD results to pickle
        try:
            with open(out_pkl, 'wb') as f:
                pickle.dump(psd_results, f)
            print(f"Saved PSD results to '{out_pkl}'")
        except Exception as e:
            print(f"ERROR saving PSD results: {e}")
            return
        
        # 9. Done
        print("Done computing PSD! Ready to plot.")

##############################################################################
# 4) PLOT PSD
##############################################################################

def on_plot_psd_clicked(b):
    global last_psd_figure
    with psd_output_area:
        clear_output()
        
        if not psd_results:
            print("ERROR: No PSD results found. Please compute PSD first.")
            return
        
        selected_channels = list(plot_channel_widget.value)
        if not selected_channels:
            print("ERROR: No channels selected for plotting.")
            return
        
        plot_mode = plot_mode_widget.value  # "mean_std" or "individual_mean"
        
        x_min = x_min_widget.value
        x_max = x_max_widget.value
        y_min = y_min_widget.value
        y_max = y_max_widget.value
        
        if x_min >= x_max or y_min >= y_max:
            print("ERROR: Invalid axis limits.")
            return
        
        # We expect psd_results = { "block1": {ch => {...}}, "block2": {...}, ...}
        block_names = sorted(psd_results.keys())  # e.g. ["block1"] or ["block1","block2",...]
        
        n_blocks = len(block_names)
        n_channels = len(selected_channels)
        
        print(f"Plotting {n_channels} channel(s) across {n_blocks} block(s).")
        
        fig, axes = plt.subplots(n_channels, n_blocks, figsize=(4*n_blocks, 2.5*n_channels), squeeze=False)
        
        for row_idx, channel in enumerate(selected_channels):
            for col_idx, block_name in enumerate(block_names):
                ax = axes[row_idx, col_idx]
                
                block_data = psd_results[block_name]  # dict {ch => {"freqs","psd"}}
                if channel not in block_data or block_data[channel]["freqs"] is None:
                    ax.text(0.5,0.5, f"No data\n{block_name}\n{channel}", va='center', ha='center')
                    ax.set_xlim(0,1)
                    ax.set_ylim(0,1)
                    continue
                
                freqs = block_data[channel]["freqs"]
                psd_array = block_data[channel]["psd"]  # shape => (n_epochs_in_block, n_freqs)
                
                psd_mean = np.nanmean(psd_array, axis=0)
                psd_std = np.nanstd(psd_array, axis=0)
                
                if plot_mode == 'mean_std':
                    ax.plot(freqs, psd_mean, color='blue', label='Mean')
                    ax.fill_between(freqs, psd_mean, psd_mean+psd_std, color='blue', alpha=0.3, label='+Std')
                else:
                    # "individual_mean"
                    for ep_idx in range(psd_array.shape[0]):
                        ax.plot(freqs, psd_array[ep_idx], color='lightgray', linewidth=0.5)
                    ax.plot(freqs, psd_mean, color='blue', label='Mean')
                
                if row_idx == 0:
                    ax.set_title(block_name, fontsize=10)
                if row_idx == n_channels-1:
                    ax.set_xlabel("Freq (Hz)")
                if col_idx == 0:
                    ax.set_ylabel(f"{channel}\nPSD (V²/Hz)")
                
                ax.set_xlim(x_min, x_max)
                ax.set_ylim(y_min, y_max)
        
        plt.tight_layout()
        plt.show()
        
        # Store the figure globally so we can export it
        last_psd_figure = fig
        
        print("Done plotting PSD side by side.")

##############################################################################
# 5) FIGURE SAVE CALLBACK
##############################################################################

def on_save_figure_clicked(b):
    """
    Save the last plotted figure in the chosen format and location.
    For PPTX, we embed the figure as an image into a single-slide presentation.
    """
    global last_psd_figure
    with save_figure_output:
        clear_output()
        
        if last_psd_figure is None:
            print("ERROR: There is no figure to save. Please plot first.")
            return
        
        # User-chosen directory or file
        chosen_path = figure_export_chooser.selected
        if not chosen_path:
            print("ERROR: No export path selected.")
            return
        
        # Ensure we know if user picked a file or just a directory
        export_dir = ""
        export_name = ""
        if os.path.isdir(chosen_path):
            # User selected a directory. We'll create a file name in that directory
            export_dir = chosen_path
            export_name = "myfigure"
        else:
            # A file was selected or typed in the text field
            export_dir = os.path.dirname(chosen_path)
            export_name = os.path.splitext(os.path.basename(chosen_path))[0]
        
        export_format = export_format_toggle.value  # 'pptx', 'png', 'jpeg', 'svg'
        
        # Final file path
        out_file = os.path.join(export_dir, export_name + '.' + export_format)
        
        if export_format in ('png', 'jpeg', 'svg'):
            # Directly save as an image (or svg)
            last_psd_figure.savefig(out_file, format=export_format, dpi=300)
            print(f"Figure saved as {out_file}")
            
        elif export_format == 'pptx':
            if not PPTX_AVAILABLE:
                print("ERROR: python-pptx is not installed. Install it via 'pip install python-pptx'.")
                return
            # 1) Save figure temporarily as PNG
            temp_img = os.path.join(export_dir, export_name + '_temp.png')
            last_psd_figure.savefig(temp_img, format='png', dpi=300)
            # 2) Create a PPTX and insert the image
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # empty layout
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(1)
            pic = slide.shapes.add_picture(temp_img, left, top, width=Inches(6))
            prs.save(out_file)
            # 3) Remove temp file if you want
            try:
                os.remove(temp_img)
            except OSError:
                pass
            print(f"Figure exported to {out_file}")
        else:
            print("Unknown export format selected.")

##############################################################################
# 6) GUI LAYOUT
##############################################################################

# Step 1: Compute PSD
compute_box = VBox([
    Label("STEP 1: Select Filtered .fif"),
    input_fif_chooser,
    segment_toggle,
    HBox([window_length_widget, overlap_widget]),
    output_pickle_chooser,
    compute_psd_button,
    psd_progress,
    psd_output_area
])

# Step 2: Plot PSD
plot_box = VBox([
    Label("STEP 2: Plot PSD (Side by Side if multiple blocks)"),
    HBox([Label("Channels:"), plot_channel_widget]),
    HBox([Label("Plot Mode:"), plot_mode_widget]),
    HBox([
        Label("X Range:"),
        x_min_widget,
        x_max_widget,
        Label("Y Range:"),
        y_min_widget,
        y_max_widget
    ]),
    plot_psd_button
])

# Step 3: Export Figure
export_box = VBox([
    Label("STEP 3: Export Figure"),
    export_format_toggle,
    figure_export_chooser,
    save_figure_button,
    save_figure_output
])

main_ui = VBox([
    widgets.HTML("<h2>PSD Computation with Optional Segmentation</h2>"),
    compute_box,
    plot_box,
    export_box
])

# Attach callbacks
compute_psd_button.on_click(on_compute_psd_clicked)
plot_psd_button.on_click(on_plot_psd_clicked)
save_figure_button.on_click(on_save_figure_clicked)

# Display UI
display(main_ui)
