import os
import pandas as pd
import numpy as np
import ipywidgets as widgets
from IPython.display import display, clear_output
import matplotlib.pyplot as plt
from fooof import FOOOFGroup
from scipy.spatial.distance import cdist
from tqdm.notebook import tqdm

# Tkinter for file/directory dialogs
from tkinter import Tk, filedialog

# For pickling and unpickling the fg_dict
import pickle

# Optional: For PPT export (install python-pptx if needed)
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

###############################################################################
#                           CORE ANALYSIS FUNCTIONS
###############################################################################

def process_excel(file_paths):
    """
    Reads .xlsx files and computes PSD data dictionaries.
    Returns a dict: { file_path: { sheet_name: PSD_values, 'freq': freq_array } }
    """
    psd_data_dict_all = {}

    for file_path in file_paths:
        try:
            excel_data = pd.ExcelFile(file_path)
            psd_data_dict = {}
            freq_data = None

            for sheet_name in excel_data.sheet_names:
                sheet_df = pd.read_excel(file_path, sheet_name=sheet_name)
                freq_column = sheet_df.iloc[:, 0]
                if freq_data is None:
                    freq_data = freq_column.values

                psd_columns = sheet_df.iloc[:, 1:]
                averaged_psd = psd_columns.mean(axis=1)
                psd_data_dict[sheet_name] = averaged_psd.values

            psd_data_dict["freq"] = freq_data
            psd_data_dict_all[file_path] = psd_data_dict
        except Exception as e:
            print(f"An error occurred while processing {file_path}: {e}")

    return psd_data_dict_all

def run_fooof_analysis(psd_data_dict_all, freq_range, amp_threshold, r2_threshold, max_peaks, fitting_mode, peak_width_limits):
    """
    Runs FOOOF analysis on each sheet's PSD data. Returns a dictionary (fg_dict)
    mapping a unique key (sheet + filename) -> FOOOFGroup object.
    """
    fg_dict = {}

    for file_path, psd_data_dict in tqdm(psd_data_dict_all.items(), desc="Fitting PSDs with FOOOF"):
        freq = psd_data_dict["freq"]
        # Use filename minus .xlsx for suffix
        file_suffix = os.path.basename(file_path).replace(".xlsx", "")

        for sheet_name, psd in psd_data_dict.items():
            if sheet_name == "freq":
                continue

            fg = FOOOFGroup(
                peak_width_limits=peak_width_limits,
                max_n_peaks=max_peaks,
                min_peak_height=amp_threshold,
                verbose=False,
                aperiodic_mode=fitting_mode,
            )

            fg.fit(freq, psd[np.newaxis, :], freq_range)
            # Drop any fits below R² threshold
            fg.drop(fg.get_params('r_squared') < r2_threshold)

            dict_key = f"{sheet_name}_{file_suffix}"
            fg_dict[dict_key] = fg

    return fg_dict

def plot_closest_to_mean(fg, sheet_name, export_dir, base_filename, formats):
    """
    Finds 10 spectra closest to the mean and plots them.
    Saves figures in the specified formats and directory.
    """
    models = [fg.get_fooof(ind=i) for i in range(len(fg))]
    spectra = [model.power_spectrum for model in models]
    spectra = np.array(spectra)

    if spectra.ndim != 2:
        raise ValueError(f"Unexpected spectra shape: {spectra.shape}. Expected a 2D array.")

    mean_spectrum = np.mean(spectra, axis=0)
    distances = cdist(spectra, mean_spectrum[None, :], metric='euclidean').flatten()
    closest_indices = np.argsort(distances)[:10]

    # Optional: Prepare a PPTX if 'ppt' in requested formats
    if "ppt" in formats and PPTX_AVAILABLE:
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank layout

    for i, index in enumerate(closest_indices):
        fm = fg.get_fooof(ind=index, regenerate=True)
        title = f"{sheet_name} - Closest Example {index}"
        fm.plot(title=title, plot_peaks='shade')

        fig = plt.gcf()
        file_tag = f"{base_filename}_{sheet_name}_ex{i}"

        # Export to chosen formats
        for fmt in formats:
            if fmt in ["png", "svg", "jpeg"]:
                filename = os.path.join(export_dir, f"{file_tag}.{fmt}")
                fig.savefig(filename, dpi=300, bbox_inches='tight')
            elif fmt == "ppt" and PPTX_AVAILABLE:
                tmp_png = os.path.join(export_dir, f"{file_tag}_ppt_temp.png")
                fig.savefig(tmp_png, dpi=300, bbox_inches='tight')
                slide = prs.slides.add_slide(slide_layout)
                left = top = Inches(1)
                pic = slide.shapes.add_picture(tmp_png, left, top, height=Inches(5))
        plt.show()
        plt.clf()

    # Finalize PPT if requested
    if "ppt" in formats and PPTX_AVAILABLE:
        pptx_filename = os.path.join(export_dir, f"{base_filename}_{sheet_name}.pptx")
        prs.save(pptx_filename)

###############################################################################
#                    SAVE/LOAD FOOOFGroup PICKLES
###############################################################################

def save_fooof_group(fg_dict, filename="fooof_groups.pkl"):
    """
    Serializes and saves the fg_dict (mapping keys -> FOOOFGroup) to a pickle file.
    """
    with open(filename, "wb") as f:
        pickle.dump(fg_dict, f)
    print(f"FOOOFGroup dictionary saved to {filename}")

def load_fooof_group(filename="fooof_groups.pkl"):
    """
    Loads the fg_dict from a pickle file.
    """
    with open(filename, "rb") as f:
        fg_dict = pickle.load(f)
    print(f"FOOOFGroup dictionary loaded from {filename}")
    return fg_dict

###############################################################################
#                         TKINTER-BASED GUI CODE
###############################################################################

selected_files = []
export_directory = ""
fg_dict = {}  # Will store final FOOOF results

# Parameter widgets
freq_range_min = widgets.FloatText(description="Freq Min:", value=4.0, layout=widgets.Layout(width="30%"))
freq_range_max = widgets.FloatText(description="Freq Max:", value=45.0, layout=widgets.Layout(width="30%"))
amp_threshold = widgets.FloatText(description="Amplitude Threshold:", value=0.2, layout=widgets.Layout(width="50%"))
r2_threshold = widgets.FloatText(description="R² Threshold:", value=0.5, layout=widgets.Layout(width="50%"))
max_peaks = widgets.IntText(description="Max Peaks:", value=2, layout=widgets.Layout(width="50%"))
fitting_mode = widgets.Dropdown(description="Fitting Mode:", options=["knee", "fixed"], value="knee", layout=widgets.Layout(width="50%"))
peak_width_min = widgets.FloatText(description="Peak Width Min:", value=2.0, layout=widgets.Layout(width="50%"))
peak_width_max = widgets.FloatText(description="Peak Width Max:", value=10.0, layout=widgets.Layout(width="50%"))

# Figure format checkboxes
fmt_options = ["png", "svg", "jpeg", "ppt"]
format_checkboxes = [widgets.Checkbox(value=False, description=fmt.upper()) for fmt in fmt_options]

# Buttons
file_picker_button = widgets.Button(description="Select Excel Files", button_style="info", icon="folder")
directory_picker_button = widgets.Button(description="Select Output Directory", button_style="info", icon="folder")
run_button = widgets.Button(description="Process and Analyze", button_style="success", icon="check")
output = widgets.Output()

###################################
# 1. FILE SELECTION CALLBACK
###################################
def on_file_picker_button_click(b):
    global selected_files
    root = Tk()
    root.withdraw()  # Hide the root window
    file_paths = filedialog.askopenfilenames(filetypes=[("Excel files", "*.xlsx")])
    selected_files = list(file_paths)
    if selected_files:
        print(f"Selected files: {selected_files}")
    else:
        print("No files selected.")

file_picker_button.on_click(on_file_picker_button_click)

###################################
# 2. DIRECTORY SELECTION CALLBACK
###################################
def on_directory_picker_button_click(b):
    global export_directory
    root = Tk()
    root.withdraw()  # Hide the root window
    directory_path = filedialog.askdirectory()
    if directory_path:
        export_directory = directory_path
        print(f"Selected output directory: {export_directory}")
    else:
        print("No directory selected.")

directory_picker_button.on_click(on_directory_picker_button_click)

######################################################
# 3. MAIN BUTTON: PROCESS, ANALYZE & EXPORT FIGURES
######################################################
def on_run_button_click(b):
    global fg_dict
    with output:
        output.clear_output()

        # 1. Check if we have selected files
        if not selected_files:
            print("Error: Please select Excel files to process.")
            return

        # 2. Check if we have an export directory
        if not export_directory:
            print("Error: Please select an output directory.")
            return

        # 3. Gather figure formats
        chosen_formats = [c.description.lower() for c in format_checkboxes if c.value]
        if not chosen_formats:
            print("Warning: No figure format selected. Proceeding without figure export.")
        
        # 4. Process Excel files
        print(f"Processing files: {selected_files}")
        psd_data_dict_all = process_excel(selected_files)

        if psd_data_dict_all:
            # 5. Run FOOOF
            freq_range = [freq_range_min.value, freq_range_max.value]
            peak_width_limits = [peak_width_min.value, peak_width_max.value]

            print("Running FOOOF analysis...")
            fg_dict = run_fooof_analysis(
                psd_data_dict_all,
                freq_range=freq_range,
                amp_threshold=amp_threshold.value,
                r2_threshold=r2_threshold.value,
                max_peaks=max_peaks.value,
                fitting_mode=fitting_mode.value,
                peak_width_limits=peak_width_limits,
            )

            # 6. Plot and export
            print("\nAnalyzing quality of FOOOF fits (closest to mean)...")
            for dict_key, fg in fg_dict.items():
                try:
                    print(f"\nPlotting 10 examples closest to mean for: {dict_key}")
                    base_filename = dict_key
                    plot_closest_to_mean(fg, dict_key, export_directory, base_filename, chosen_formats)
                except ValueError as e:
                    print(f"Error processing {dict_key}: {e}")

            # 7. Save the entire fg_dict to fooof_groups.pkl in the chosen directory
            pickle_filename = os.path.join(export_directory, "fooof_groups.pkl")
            save_fooof_group(fg_dict, filename=pickle_filename)

            print("All exports complete. You can now verify 'fooof_groups.pkl' or load it using 'load_fooof_group'.")
        else:
            print("Processing failed. Please check the files and try again.")

run_button.on_click(on_run_button_click)

##############################
# 4. DISPLAY THE GUI
##############################
display(
    widgets.VBox(
        [
            widgets.Label(value="1) Select Excel Files:"),
            file_picker_button,
            widgets.Label(value="2) Select Output Directory:"),
            directory_picker_button,
            widgets.Label(value="3) Configure Analysis Parameters:"),
            widgets.HBox([freq_range_min, freq_range_max]),
            amp_threshold,
            r2_threshold,
            max_peaks,
            fitting_mode,
            widgets.HBox([peak_width_min, peak_width_max]),
            widgets.Label(value="4) Choose Figure Export Formats:"),
            widgets.HBox(format_checkboxes),
            run_button,
            output,
        ]
    )
)
