import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math
import re

import ipywidgets as widgets
from ipywidgets import VBox, HBox
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed. PPT export will be disabled.")

# Check for pandas (for Excel export)
try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    print("Warning: pandas not installed. Excel export will be disabled.")

###############################################################################
# Helper Functions
###############################################################################

def custom_load_psd(file_path, loader_method="pickle"):
    """
    Attempt to load a PSD dictionary from file_path using the chosen loader method.
    """
    if loader_method == "pickle":
        with open(file_path, "rb") as f:
            return pickle.load(f)
    elif loader_method == "dill":
        import dill
        with open(file_path, "rb") as f:
            return dill.load(f)
    elif loader_method == "joblib":
        import joblib
        return joblib.load(file_path)
    elif loader_method == "torch":
        import torch
        return torch.load(file_path)
    else:
        raise ValueError(f"Unknown loader method: {loader_method}")

def compute_combined_means_for_block(psds_dict, block, eye_channels, deselected_channels):
    """
    Returns the average PSD across all selected channels in eye_channels for the given block.
    """
    combined_mean = None
    valid_channels = 0
    for ch in eye_channels:
        channel_key = f"{block}:{ch}"
        if channel_key not in psds_dict or ch in deselected_channels:
            continue
        psd = psds_dict[channel_key]['psd']
        mean_psd = np.mean(psd, axis=0)
        if combined_mean is None:
            combined_mean = mean_psd.copy()
        else:
            combined_mean += mean_psd
        valid_channels += 1

    if valid_channels > 0:
        combined_mean /= valid_channels

    return combined_mean

def plot_aggregated_means(
    combined_means,
    freqs,
    eye,
    title_fs=10,
    axis_fs=8,
    legend_fs=8,
    y_max=80,
    mean_line_color='black',
    show_mean_of_means=True
):
    """
    Plots the aggregated PSD traces. Each entry in combined_means is plotted as a dashed line.
    Then optionally plots the average of all those lines ('Mean of Means') in a solid line.
    """
    fig, ax = plt.subplots(figsize=(12, 6))

    mean_of_means = None
    count = 0

    for key, mean in combined_means.items():
        if mean is None:
            continue
        ax.plot(freqs, mean, linestyle='--', label=f"{eye}:{key} Mean")
        if mean_of_means is None:
            mean_of_means = mean.copy()
        else:
            mean_of_means += mean
        count += 1

    if mean_of_means is not None and count > 0:
        mean_of_means /= count
        if show_mean_of_means:
            ax.plot(freqs, mean_of_means, linestyle='-', color=mean_line_color, linewidth=2, label="Mean of Means")

    ax.set_title(f"{eye} Aggregated Means", fontsize=title_fs)
    ax.set_xlim([0, 45])
    ax.set_ylim([0, y_max])
    ax.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
    ax.legend(fontsize=legend_fs)
    plt.tight_layout()
    display(fig)  # display() ensures it appears in the notebook cell
    return fig  # Return the figure for potential exporting

def plot_block_layout(
    psds_dict,
    channel_means,
    freqs,
    deselected_channels,
    show_original=True,
    show_new=True,
    title_fs=10,
    axis_fs=8,
    legend_fs=8,
    plot_width=12,
    plot_height=18,
    y_max=80,
    mean_line_color='blue',
    channel_colors=None,
    show_combined_means=False,
    show_aggregated_plots=True,
    aggregation_mode='block'
):
    """
    Main plotting function:
      1) Plots each block (row) with Eye1 (Ch1-8) in one column, Eye2 (Ch9-16) in the other column.
      2) Optionally plots aggregated plots. 
         - If aggregation_mode='block', each aggregated trace is the average of all selected channels per block.
         - If aggregation_mode='channel', each aggregated trace is the average of that channel across blocks.
    """

    if channel_colors is None:
        color_cycle = plt.cm.get_cmap('tab20').colors
        unique_channels = sorted({key.split(':Ch')[-1] for key in psds_dict.keys()})
        channel_colors = {f"Ch{ch}": color_cycle[i % len(color_cycle)] for i, ch in enumerate(unique_channels)}

    # Distinct blocks and channel groupings
    blocks = sorted(set([key.split(':')[0] for key in psds_dict.keys()]))
    eye1_channels = [f"Ch{i}" for i in range(1, 9)]
    eye2_channels = [f"Ch{i}" for i in range(9, 17)]

    # Create subplots for all blocks
    fig, axes = plt.subplots(len(blocks), 2, figsize=(plot_width, plot_height))
    if len(blocks) == 1:
        axes = [axes]  # so we can index axes[i][0], axes[i][1]

    combined_eye1 = {}
    combined_eye2 = {}

    # =============== MAIN BLOCK-BY-BLOCK PLOTS ===============
    for i, block in enumerate(blocks):
        # Eye 1 (Ch1-8)
        ax_eye1 = axes[i][0] if len(blocks) > 1 else axes[0]
        block_mean_eye1 = []
        for channel in eye1_channels:
            channel_full = f"{block}:{channel}"
            if channel in deselected_channels or channel_full not in channel_means:
                continue
            original_mean, new_mean = channel_means[channel_full]
            color = channel_colors.get(channel, "black")
            if show_original:
                ax_eye1.plot(freqs, original_mean, linestyle='-', color=color, label=f"{channel_full} Original")
            if show_new:
                ax_eye1.plot(freqs, new_mean, linestyle='--', color=color, label=f"{channel_full} New")
            block_mean_eye1.append(new_mean)
        if show_combined_means and block_mean_eye1:
            block_combined_mean = np.mean(block_mean_eye1, axis=0)
            ax_eye1.plot(freqs, block_combined_mean, linestyle='-', color=mean_line_color, linewidth=2,
                         label=f"{block} Combined Eye 1")
        ax_eye1.set_title(f"{block} Eye 1 (Ch1-8)", fontsize=title_fs)
        ax_eye1.set_xlim([0, 45])
        ax_eye1.set_ylim([0, y_max])
        ax_eye1.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
        ax_eye1.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
        ax_eye1.legend(fontsize=legend_fs, loc='upper right')

        # Eye 2 (Ch9-16)
        ax_eye2 = axes[i][1] if len(blocks) > 1 else axes[1]
        block_mean_eye2 = []
        for channel in eye2_channels:
            channel_full = f"{block}:{channel}"
            if channel in deselected_channels or channel_full not in channel_means:
                continue
            original_mean, new_mean = channel_means[channel_full]
            color = channel_colors.get(channel, "black")
            if show_original:
                ax_eye2.plot(freqs, original_mean, linestyle='-', color=color, label=f"{channel_full} Original")
            if show_new:
                ax_eye2.plot(freqs, new_mean, linestyle='--', color=color, label=f"{channel_full} New")
            block_mean_eye2.append(new_mean)
        if show_combined_means and block_mean_eye2:
            block_combined_mean = np.mean(block_mean_eye2, axis=0)
            ax_eye2.plot(freqs, block_combined_mean, linestyle='-', color=mean_line_color, linewidth=2,
                         label=f"{block} Combined Eye 2")
        ax_eye2.set_title(f"{block} Eye 2 (Ch9-16)", fontsize=title_fs)
        ax_eye2.set_xlim([0, 45])
        ax_eye2.set_ylim([0, y_max])
        ax_eye2.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
        ax_eye2.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
        ax_eye2.legend(fontsize=legend_fs, loc='upper right')

    plt.tight_layout()
    display(fig)  # show the main block-layout figure

    # =============== AGGREGATED PLOTS ===============
    aggregated_fig1 = None
    aggregated_fig2 = None
    if show_aggregated_plots:
        if aggregation_mode == 'block':
            # Each block's average (across channels) is one line
            for block in blocks:
                # Eye 1 average for this block
                means = []
                for ch in eye1_channels:
                    ch_full = f"{block}:{ch}"
                    if ch_full in channel_means and ch not in deselected_channels:
                        _, new_mean = channel_means[ch_full]
                        means.append(new_mean)
                if len(means) > 0:
                    combined_eye1[block] = np.mean(means, axis=0)

                # Eye 2 average
                means = []
                for ch in eye2_channels:
                    ch_full = f"{block}:{ch}"
                    if ch_full in channel_means and ch not in deselected_channels:
                        _, new_mean = channel_means[ch_full]
                        means.append(new_mean)
                if len(means) > 0:
                    combined_eye2[block] = np.mean(means, axis=0)

        elif aggregation_mode == 'channel':
            # Each channel's average across blocks is one line
            for ch in eye1_channels:
                if ch in deselected_channels:
                    continue
                means = []
                for block in blocks:
                    ch_full = f"{block}:{ch}"
                    if ch_full in channel_means:
                        _, new_mean = channel_means[ch_full]
                        means.append(new_mean)
                if len(means) > 0:
                    combined_eye1[ch] = np.mean(means, axis=0)

            for ch in eye2_channels:
                if ch in deselected_channels:
                    continue
                means = []
                for block in blocks:
                    ch_full = f"{block}:{ch}"
                    if ch_full in channel_means:
                        _, new_mean = channel_means[ch_full]
                        means.append(new_mean)
                if len(means) > 0:
                    combined_eye2[ch] = np.mean(means, axis=0)

        aggregated_fig1 = plot_aggregated_means(
            combined_eye1, freqs, "Eye 1",
            title_fs=title_fs, axis_fs=axis_fs, legend_fs=legend_fs, y_max=y_max,
            mean_line_color=mean_line_color
        )
        aggregated_fig2 = plot_aggregated_means(
            combined_eye2, freqs, "Eye 2",
            title_fs=title_fs, axis_fs=axis_fs, legend_fs=legend_fs, y_max=y_max,
            mean_line_color=mean_line_color
        )
        return fig, aggregated_fig1, aggregated_fig2

    return fig

def export_figure(figs, export_path, export_format="png"):
    """
    Saves one or more figures to disk in the desired format.
    If export_format == "ppt" (and python-pptx is installed), it exports to a PPTX file.
    """
    if not isinstance(figs, (list, tuple)):
        figs = [figs]  # Ensure figs is a list

    for i, fig in enumerate(figs):
        file_path = f"{export_path}_{i+1}.{export_format}"
        if export_format.lower() in ["png", "jpeg", "svg", "pdf"]:
            fig.savefig(file_path, format=export_format, dpi=300)
            print(f"Figure exported as {file_path}")
        elif export_format.lower() == "ppt" and HAS_PPTX:
            ppt = Presentation()
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # Blank slide
            temp_img = f"{export_path}_{i+1}.png"
            fig.savefig(temp_img, format="png", dpi=300)
            slide.shapes.add_picture(temp_img, Inches(1), Inches(1), Inches(8), Inches(5))
            ppt.save(f"{export_path}.pptx")
            os.remove(temp_img)
            print(f"Figure exported as {export_path}.pptx")
        else:
            print(f"Unsupported format or missing dependencies for {export_format}.")

###############################################################################
# Main GUI
###############################################################################

def build_gui():
    psds_dict = {}
    channel_means = {}
    freqs = None
    deselected_channels = []

    file_chooser = FileChooser(os.getcwd())
    load_button = widgets.Button(description="Load PSD", button_style='info')
    plot_button = widgets.Button(description="Plot", button_style='success')
    export_button = widgets.Button(description="Export Figure", button_style='warning')
    output_area = widgets.Output()

    # Widgets for plot customization
    channel_toggles = widgets.VBox()
    plot_width = widgets.FloatText(value=12, description="Plot Width")
    plot_height = widgets.FloatText(value=18, description="Plot Height")
    y_max_input = widgets.FloatText(value=80, description="Y Max")
    mean_line_color = widgets.ColorPicker(value='blue', description="Mean Line Color")
    export_format = widgets.Dropdown(
        options=["png", "jpeg", "svg", "pdf", "ppt"], 
        value="png", 
        description="Export Format"
    )
    export_name = widgets.Text(value="plot_export", description="File Name")
    show_original = widgets.Checkbox(value=True, description="Show Original Traces")
    show_new = widgets.Checkbox(value=True, description="Show New Traces")
    show_combined = widgets.Checkbox(value=False, description="Show Combined Means")
    show_aggregated = widgets.Checkbox(value=True, description="Show Aggregated Plots")

    # Widget to select how to aggregate the data
    aggregation_mode = widgets.Dropdown(
        options=["block", "channel"],
        value="block",
        description="Aggregation Mode",
        style={'description_width': 'initial'}
    )

    # **New** widget to specify the library used to load the pickle
    pickle_loader_method = widgets.Dropdown(
        options=["pickle", "dill", "joblib", "torch"],
        value="pickle",
        description="Loader Method",
        style={'description_width': 'initial'}
    )

    def update_channel_toggles():
        nonlocal deselected_channels
        unique_channels = sorted({key.split(':Ch')[-1] for key in psds_dict.keys()})
        toggles = []
        for ch in unique_channels:
            toggle = widgets.Checkbox(value=True, description=f"Ch{ch}")
            def toggle_change(change, ch=ch):
                if not change['new']:
                    if ch not in deselected_channels:
                        deselected_channels.append(f"Ch{ch}")
                else:
                    if f"Ch{ch}" in deselected_channels:
                        deselected_channels.remove(f"Ch{ch}")
            toggle.observe(toggle_change, names='value')
            toggles.append(toggle)
        channel_toggles.children = toggles

    def load_psd_callback(_):
        nonlocal psds_dict, freqs, channel_means
        with output_area:
            clear_output()
            file_path = file_chooser.selected
            if not file_path or not os.path.exists(file_path):
                print("Please select a valid file.")
                return
            try:
                # Use our custom loader with selected library
                loader = pickle_loader_method.value
                psds_dict = custom_load_psd(file_path, loader_method=loader)

                if not psds_dict:
                    print("Loaded PSD data is empty.")
                    return
                freqs = next(iter(psds_dict.values()))['freqs']

                # Build the "channel_means" dict: (original_mean, new_mean)
                channel_means = {
                    key: (
                        np.mean(data['psd'], axis=0),  # Original mean
                        np.mean(data['psd'], axis=0)   # New mean (placeholder)
                    ) for key, data in psds_dict.items()
                }
                print(f"Loaded PSD data using {loader}.")
                update_channel_toggles()
            except Exception as e:
                print(f"Failed to load file: {e}")

    def plot_callback(_):
        with output_area:
            clear_output()
            if not psds_dict:
                print("No PSD data loaded.")
                return
            try:
                figs = plot_block_layout(
                    psds_dict,
                    channel_means,
                    freqs,
                    deselected_channels,
                    show_original=show_original.value,
                    show_new=show_new.value,
                    plot_width=plot_width.value,
                    plot_height=plot_height.value,
                    y_max=y_max_input.value,
                    mean_line_color=mean_line_color.value,
                    show_combined_means=show_combined.value,
                    show_aggregated_plots=show_aggregated.value,
                    aggregation_mode=aggregation_mode.value
                )
                return figs
            except Exception as e:
                print(f"Failed to plot data: {e}")

    def export_callback(_):
        with output_area:
            if not psds_dict:
                print("No PSD data loaded. Please load and plot first.")
                return
            try:
                export_dir = os.getcwd()
                filename = export_name.value
                export_path = os.path.join(export_dir, filename)
                figs = plot_callback(None)  # Re-plot and capture the figures
                if isinstance(figs, tuple):
                    export_figure(list(figs), export_path, export_format.value)
                else:
                    export_figure(figs, export_path, export_format.value)
            except Exception as e:
                print(f"Failed to export figure: {e}")

    load_button.on_click(load_psd_callback)
    plot_button.on_click(plot_callback)
    export_button.on_click(export_callback)

    display(
        VBox([
            file_chooser,
            HBox([load_button, plot_button, export_button]),
            HBox([pickle_loader_method]),   # <-- Let user choose the loader
            HBox([plot_width, plot_height, y_max_input]),
            HBox([mean_line_color, export_name, export_format]),
            HBox([show_original, show_new, show_combined, show_aggregated]),
            aggregation_mode,
            widgets.Label("Select Channels to Include:"),
            channel_toggles,
            output_area
        ])
    )

build_gui()
