
"""
psd_channel_means_v1b_NaN.py

GUI to work *after* psd_clean_v2d3_NaN:

- Input: a "cleaned" PSD pickle produced by psd_clean_v2d3_NaN
  (dict: key -> {"freqs": 1D array, "psd": 2D array (epochs x freqs)}).
- NaN-aware: uses np.nanmean so it behaves nicely if PSD contains NaNs.
- Channel naming is flexible and works with:
    - Eye-averaged channels: "Eye1_avg", "Eye2_avg",
    - Blocked names: "block1:Eye1_avg", "block1:Eye2_avg", etc.,
    - Or legacy "Ch1".."Ch16" style names.

Main features:
- Load cleaned PSD pickle via ipyfilechooser.
- Select channels (keys) to screen.
- Visualize original vs "new" (post-exclusion) channel means per Eye1/Eye2.
- Interactively select which channels contribute to final Eye1/Eye2 group means.
- Plot "means only" figure (Eye1 vs Eye2 final means of means).
- Export:
    - All figures (screening + final) as images or PPTX,
    - Final means data/plot as .pkl, .xlsx, .png/.svg/.jpg, or PPTX.

Usage in Jupyter:

    import sys
    sys.path.append("/Users/davidlitvin/Multielectrode array rsERG/Updated_MEA_Pipeline")
    import psd_channel_means_v1b_NaN
    psd_channel_means_v1b_NaN.launch_psd_channel_means_nan_gui()
"""

import os
import pickle
import numpy as np
import matplotlib.pyplot as plt

import ipywidgets as widgets
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Optional PPTX export
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed; PPT export will be disabled.")


# ---------------------------------------------------------------------
# Helper functions (NaN-aware versions)
# ---------------------------------------------------------------------

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=None,
    test_band_threshold=10.0,
    test_band_count_threshold=None,
):
    """
    NaN-aware trace exclusion.

    Parameters
    ----------
    psd_array : array, shape (n_epochs, n_freqs)
        PSD per epoch (rows) for one channel; may contain NaNs.
    freqs : array, shape (n_freqs,)
    low_band : tuple
        e.g. (1, 3) Hz for low-frequency outliers.
    low_band_threshold : float
        Multiple of mean PSD in low_band used to flag an epoch.
    test_bands : list of (fmin, fmax) or None
        Frequency bands used to detect repeated suprathreshold events.
    test_band_threshold : float
        Multiple of mean PSD in each test band to count as suprathreshold.
    test_band_count_threshold : int or None
        Number of suprathreshold bands needed to exclude an epoch.
        If None, defaults to len(test_bands)//2.

    Returns
    -------
    kept_traces     : list of epoch indices
    excluded_traces : list of epoch indices
    """
    if test_bands is None:
        test_bands = [(7, 9), (9, 11), (11, 13),
                      (13, 15), (15, 17), (17, 19), (19, 21)]
    if test_band_count_threshold is None:
        test_band_count_threshold = max(1, len(test_bands) // 2)

    # NaN-robust mean PSD across epochs
    with np.errstate(invalid="ignore"):
        mean_psd = np.nanmean(psd_array, axis=0)

    excluded_traces = []
    kept_traces = []

    low_band_idx = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0]
        for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # If entire trace is NaN, mark as excluded
        if np.all(np.isnan(trace)):
            excluded_traces.append(i)
            continue

        # 1) Low-frequency outlier (NaN-safe)
        if low_band_idx.size > 0:
            t_seg = trace[low_band_idx]
            m_seg = mean_psd[low_band_idx]
            # Only compare where both are finite
            valid = np.isfinite(t_seg) & np.isfinite(m_seg)
            if np.any(valid):
                if np.any(t_seg[valid] > low_band_threshold * m_seg[valid]):
                    excluded_traces.append(i)
                    continue

        # 2) Repeated suprathreshold events across test bands
        suprathreshold_count = 0
        for idx in band_indices:
            if idx.size == 0:
                continue
            t_seg = trace[idx]
            m_seg = mean_psd[idx]
            valid = np.isfinite(t_seg) & np.isfinite(m_seg)
            if np.any(valid):
                if np.any(t_seg[valid] > test_band_threshold * m_seg[valid]):
                    suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces


def compute_channel_means(
    psds_dict,
    exclude=False,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=None,
    test_band_threshold=10.0,
    test_band_count_threshold=None,
):
    """
    Compute channel means before and (optionally) after exclusion.

    Parameters
    ----------
    psds_dict : dict
        key -> {"freqs": array, "psd": array (epochs x freqs)}
    exclude : bool
        If True, apply exclude_traces to each channel.
    Returns
    -------
    channel_means : dict
        channel -> (original_mean, new_mean)
        where new_mean is either the same as original_mean (exclude=False)
        or the mean over kept epochs (exclude=True).
    """
    if test_bands is None:
        test_bands = [(7, 9), (9, 11), (11, 13),
                      (13, 15), (15, 17), (17, 19), (19, 21)]

    channel_means = {}
    for channel, data in psds_dict.items():
        psd = data.get("psd", None)
        freqs = data.get("freqs", None)
        if psd is None or freqs is None:
            print(f"Channel '{channel}' missing 'psd' or 'freqs'; skipping.")
            continue
        if not isinstance(psd, np.ndarray) or psd.ndim != 2:
            print(f"Channel '{channel}' PSD has invalid shape {getattr(psd, 'shape', None)}; skipping.")
            continue

        # NaN-safe original mean
        with np.errstate(invalid="ignore"):
            original_mean = np.nanmean(psd, axis=0)

        if exclude:
            kept_traces, _ = exclude_traces(
                psd_array=psd,
                freqs=freqs,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold,
            )
            if kept_traces:
                with np.errstate(invalid="ignore"):
                    new_mean = np.nanmean(psd[kept_traces], axis=0)
            else:
                new_mean = np.full_like(original_mean, np.nan)
        else:
            new_mean = original_mean.copy()

        channel_means[channel] = (original_mean, new_mean)

    return channel_means


def compute_group_mean(channel_means, selected_channels):
    """
    Compute the mean of 'new_mean' across selected channels.
    Uses np.nanmean to be robust to NaNs.
    """
    means = []
    for ch in selected_channels:
        if ch in channel_means:
            _, new_m = channel_means[ch]
            means.append(new_m)
    if means:
        stacked = np.vstack(means)
        with np.errstate(invalid="ignore"):
            return np.nanmean(stacked, axis=0)
    else:
        # If nothing is selected, try to return zeros of correct shape
        if channel_means:
            first_val = next(iter(channel_means.values()))
            return np.zeros_like(first_val[0])
        return None


def plot_individual_channels(
    ax,
    channels,
    channel_means_dict,
    freqs,
    colors_for_channels,
    show_original_mean=True,
    show_new_mean=True,
    title="",
    axis_fs=10,
    legend_fs=8,
    tick_fs=8,
):
    """
    Plot each channel (original/new mean) with a distinct color.
    """
    for ch in channels:
        if ch not in channel_means_dict:
            continue
        orig_mean, new_mean = channel_means_dict[ch]
        color = colors_for_channels.get(ch, "blue")

        if show_original_mean:
            ax.plot(freqs, orig_mean, color=color, linestyle="-", label=f"{ch} Orig")
        if show_new_mean:
            ax.plot(freqs, new_mean, color=color, linestyle="--", label=f"{ch} New")

    ax.set_title(title, fontsize=axis_fs)
    ax.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
    ax.tick_params(axis="both", labelsize=tick_fs)
    ax.legend(fontsize=legend_fs, loc="upper right")


def _classify_eye_group(ch_name):
    """
    Classify channel name into 'Eye1', 'Eye2', or None.

    Handles patterns like:
        'Eye1_avg', 'Eye2_avg',
        'block1:Eye1_avg', 'block2:Eye2_avg',
        or classic 'Ch1'..'Ch16' (Eye1=1..8, Eye2=9..16).
    """
    # Strip block prefix if present
    if ":" in ch_name:
        _, core = ch_name.split(":", 1)
    else:
        core = ch_name
    lower = core.lower()

    # EyeAvg style
    if "eye1" in lower:
        return "Eye1"
    if "eye2" in lower:
        return "Eye2"

    # Legacy "Ch1..Ch16" style (keep for backward compatibility)
    if core.startswith("Ch"):
        try:
            num = int(core.replace("Ch", ""))
        except ValueError:
            num = None
        if num is not None:
            if 1 <= num <= 8:
                return "Eye1"
            if 9 <= num <= 16:
                return "Eye2"

    return None


# ---------------------------------------------------------------------
# Main GUI
# ---------------------------------------------------------------------

def build_exportable_plot_psd_gui():
    """
    Build and display an interactive GUI to:
      1) Load a cleaned PSD pickle (from psd_clean_v2d3_NaN).
      2) Screen channel means per Eye1/Eye2.
      3) Interactively choose which channels contribute to final Eye1/Eye2
         group means.
      4) Plot a "means only" figure (Eye1 vs Eye2).
      5) Export all figures & final means data/plot.
    """
    # Shared state
    loaded_psd = {}        # key -> {"freqs", "psd"}
    current_figures = []   # list of matplotlib Figures
    final_data_dict = {}   # e.g., {"Eye1_MeanOfMeans", "Eye1_freqs", ...}
    final_fig = None       # the means-only figure

    # === 1) LOAD PSD PICKLE ===

    load_psd_button = widgets.Button(
        description="Load PSD Pickle",
        button_style="info",
        tooltip="Load a cleaned PSD .pkl from psd_clean_v2d3_NaN",
    )
    psd_file_chooser = FileChooser(
        os.getcwd(),
        title="Select cleaned PSD Pickle (.pkl)",
        select_default=False,
    )
    psd_file_chooser.show_only_files = True
    psd_file_chooser.filter_pattern = ["*.pkl"]

    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description="Channels / Keys:",
        layout=widgets.Layout(width="300px", height="200px"),
    )

    load_output_area = widgets.Output()

    def on_load_psd_clicked(b):
        with load_output_area:
            clear_output()
            path = psd_file_chooser.selected
            if not path:
                print("Please select a PSD .pkl file.")
                return
            if not os.path.isfile(path):
                print(f"File not found: {path}")
                return
            try:
                with open(path, "rb") as f:
                    loaded = pickle.load(f)
                if not isinstance(loaded, dict):
                    print("ERROR: Pickle does not contain a dict. Aborting.")
                    return
                loaded_psd.clear()
                loaded_psd.update(loaded)
                print(f"Loaded PSD pickle from:\n  {path}")
                print(f"Available keys: {list(loaded_psd.keys())}")

                # Populate dropdown
                channels_dropdown.options = sorted(loaded_psd.keys())
                if channels_dropdown.options:
                    # Select all by default
                    channels_dropdown.value = tuple(channels_dropdown.options)
            except Exception as e:
                print(f"Error loading PSD pickle: {e}")

    load_psd_button.on_click(on_load_psd_clicked)

    # === 2) PLOTTING & CONTROL WIDGETS ===

    # Eye toggles
    show_eye1_cb = widgets.Checkbox(value=True, description="Show Eye1")
    show_eye2_cb = widgets.Checkbox(value=True, description="Show Eye2")

    # Means toggles
    show_original_mean_cb = widgets.Checkbox(
        value=True, description="Show original (all epochs) channel means"
    )
    show_new_mean_cb = widgets.Checkbox(
        value=True, description="Show new (post-exclusion) channel means"
    )

    # Exclusion thresholds
    low_band_threshold_widget = widgets.FloatText(
        value=3.0, description="Low-band thr (×mean):",
        layout=widgets.Layout(width="180px"),
    )
    test_band_threshold_widget = widgets.FloatText(
        value=10.0, description="Test-band thr (×mean):",
        layout=widgets.Layout(width="200px"),
    )

    test_band_text = widgets.Textarea(
        value="(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)",
        description="Test bands:",
        layout=widgets.Layout(width="300px", height="80px"),
    )

    # Axis ranges
    x_min_widget = widgets.FloatText(
        value=None, description="X min (Hz):",
        layout=widgets.Layout(width="150px"),
    )
    x_max_widget = widgets.FloatText(
        value=None, description="X max (Hz):",
        layout=widgets.Layout(width="150px"),
    )
    y_min_widget = widgets.FloatText(
        value=None, description="Y min:",
        layout=widgets.Layout(width="150px"),
    )
    y_max_widget = widgets.FloatText(
        value=None, description="Y max:",
        layout=widgets.Layout(width="150px"),
    )

    # Fonts
    title_fs_widget = widgets.IntText(
        value=10, description="Title FS:",
        layout=widgets.Layout(width="150px"),
    )
    axis_fs_widget = widgets.IntText(
        value=8, description="Axis FS:",
        layout=widgets.Layout(width="150px"),
    )
    legend_fs_widget = widgets.IntText(
        value=8, description="Legend FS:",
        layout=widgets.Layout(width="150px"),
    )
    tick_fs_widget = widgets.IntText(
        value=8, description="Tick FS:",
       layout=widgets.Layout(width="150px"),
    )

    # Colors (for display; kept/excluded colors not central here but we keep them)
    color_kept_widget = widgets.ColorPicker(
        value="lightgray",
        description="Kept trace:",
        layout=widgets.Layout(width="180px"),
    )
    color_excl_widget = widgets.ColorPicker(
        value="red",
        description="Excluded trace:",
        layout=widgets.Layout(width="180px"),
    )
    color_oldmean_widget = widgets.ColorPicker(
        value="blue",
        description="Orig mean:",
        layout=widgets.Layout(width="180px"),
    )
    color_newmean_widget = widgets.ColorPicker(
        value="green",
        description="New mean:",
        layout=widgets.Layout(width="180px"),
    )

    plot_psd_button = widgets.Button(
        description="Plot channel means",
        button_style="success",
    )

    plot_output_area = widgets.Output()

    # === 3) EXPORT CONTROLS ===

    export_fig_chooser = FileChooser(
        os.getcwd(),
        title="Export all figures as images or PPTX (e.g. group_psd.png / group_psd.pptx):",
        select_default=False,
    )
    export_button = widgets.Button(
        description="Export figures",
        button_style="warning",
    )
    export_output_area = widgets.Output()

    final_export_chooser = FileChooser(
        os.getcwd(),
        title="Export final Eye1/Eye2 means (e.g. final_means.pkl/.xlsx/.png/.pptx):",
        select_default=False,
    )
    final_export_button = widgets.Button(
        description="Export final means",
        button_style="info",
    )
    final_export_output = widgets.Output()

    # ------------------------------------------------------------------
    # Plot callback
    # ------------------------------------------------------------------

    def on_plot_psd_clicked(b):
        nonlocal current_figures, final_data_dict, final_fig
        with plot_output_area:
            clear_output()
            current_figures.clear()
            final_data_dict.clear()
            final_fig = None

            if not loaded_psd:
                print("No PSD data loaded. Use 'Load PSD Pickle' first.")
                return

            selected_keys = list(channels_dropdown.value)
            if not selected_keys:
                print("No channels/keys selected.")
                return

            # Parse test bands
            tb_str = test_band_text.value.strip()
            test_bands = []
            if tb_str:
                tb_str = tb_str.replace(" ", "")
                pieces = tb_str.split(")")
                for p in pieces:
                    p = p.strip(",").strip()
                    p = p.strip("(").strip()
                    if not p:
                        continue
                    vals = p.split(",")
                    if len(vals) == 2:
                        try:
                            f1 = float(vals[0])
                            f2 = float(vals[1])
                            test_bands.append((f1, f2))
                        except Exception:
                            pass
            if not test_bands:
                test_bands = [(7, 9), (9, 11), (11, 13),
                              (13, 15), (15, 17), (17, 19), (19, 21)]

            low_band = (1, 3)
            low_thr = low_band_threshold_widget.value
            test_thr = test_band_threshold_widget.value

            # Determine Eye groups based on channel names
            eye1_keys = []
            eye2_keys = []
            for ch in selected_keys:
                cls = _classify_eye_group(ch)
                if cls == "Eye1":
                    eye1_keys.append(ch)
                elif cls == "Eye2":
                    eye2_keys.append(ch)

            # If classification fails but there are only 2 keys, treat them as Eye1/Eye2
            if not eye1_keys and not eye2_keys and len(selected_keys) == 2:
                eye1_keys = [selected_keys[0]]
                eye2_keys = [selected_keys[1]]

            show_eye1 = show_eye1_cb.value
            show_eye2 = show_eye2_cb.value
            show_orig = show_original_mean_cb.value
            show_new = show_new_mean_cb.value

            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value

            title_fs = title_fs_widget.value
            axis_fs = axis_fs_widget.value
            legend_fs = legend_fs_widget.value
            tick_fs = tick_fs_widget.value

            # 1) Compute channel means (before/after exclusion)
            channel_means_before = compute_channel_means(
                psds_dict=loaded_psd,
                exclude=False,
            )
            channel_means_after = compute_channel_means(
                psds_dict=loaded_psd,
                exclude=True,
                low_band=low_band,
                low_band_threshold=low_thr,
                test_bands=test_bands,
                test_band_threshold=test_thr,
                test_band_count_threshold=None,
            )

            # 2) Screening figure: original/new per channel, grouped Eye1/Eye2

            # Build color mapping per channel for screening
            import matplotlib.cm as cm
            cmap = cm.get_cmap("tab20")
            all_screen_keys = eye1_keys + eye2_keys if (eye1_keys or eye2_keys) else selected_keys
            colors_for_channels = {}
            for i, ch in enumerate(sorted(all_screen_keys)):
                colors_for_channels[ch] = cmap(i % len(cmap.colors))

            n_subplots = 0
            if show_eye1 and eye1_keys:
                n_subplots += 1
            if show_eye2 and eye2_keys:
                n_subplots += 1

            if n_subplots == 0:
                print("No Eye1/Eye2 groups to show; will still compute final means.")
            else:
                fig_screen, axes_screen = plt.subplots(
                    1, n_subplots, figsize=(6 * n_subplots, 5)
                )
                if n_subplots == 1:
                    axes_screen = [axes_screen]

                ax_idx = 0
                if show_eye1 and eye1_keys:
                    ax1 = axes_screen[ax_idx]
                    first_key = eye1_keys[0]
                    freqs = loaded_psd[first_key]["freqs"]
                    plot_individual_channels(
                        ax=ax1,
                        channels=eye1_keys,
                        channel_means_dict=channel_means_before,
                        freqs=freqs,
                        colors_for_channels=colors_for_channels,
                        show_original_mean=show_orig,
                        show_new_mean=show_new,
                        title="Screening Eye1",
                        axis_fs=axis_fs,
                        legend_fs=legend_fs,
                        tick_fs=tick_fs,
                    )
                    if x_min is not None:
                        ax1.set_xlim(left=x_min)
                    if x_max is not None:
                        ax1.set_xlim(right=x_max)
                    if y_min is not None:
                        ax1.set_ylim(bottom=y_min)
                    if y_max is not None:
                        ax1.set_ylim(top=y_max)
                    ax_idx += 1

                if show_eye2 and eye2_keys and ax_idx < len(axes_screen):
                    ax2 = axes_screen[ax_idx]
                    first_key = eye2_keys[0]
                    freqs = loaded_psd[first_key]["freqs"]
                    plot_individual_channels(
                        ax=ax2,
                        channels=eye2_keys,
                        channel_means_dict=channel_means_before,
                        freqs=freqs,
                        colors_for_channels=colors_for_channels,
                        show_original_mean=show_orig,
                        show_new_mean=show_new,
                        title="Screening Eye2",
                        axis_fs=axis_fs,
                        legend_fs=legend_fs,
                        tick_fs=tick_fs,
                    )
                    if x_min is not None:
                        ax2.set_xlim(left=x_min)
                    if x_max is not None:
                        ax2.set_xlim(right=x_max)
                    if y_min is not None:
                        ax2.set_ylim(bottom=y_min)
                    if y_max is not None:
                        ax2.set_ylim(top=y_max)

                plt.tight_layout()
                plt.show()
                current_figures.append(fig_screen)

            # 3) Interactive selection of channels contributing to final group means

            group_sel_widgets = {}
            if eye1_keys:
                group_sel_widgets["Eye1"] = widgets.SelectMultiple(
                    options=eye1_keys,
                    value=tuple(eye1_keys),
                    description="Eye1 channels:",
                    layout=widgets.Layout(width="300px", height="150px"),
                )
            if eye2_keys:
                group_sel_widgets["Eye2"] = widgets.SelectMultiple(
                    options=eye2_keys,
                    value=tuple(eye2_keys),
                    description="Eye2 channels:",
                    layout=widgets.Layout(width="300px", height="150px"),
                )

            if group_sel_widgets:
                display(widgets.HTML("<h3>Select channels to include in final Eye1/Eye2 group means</h3>"))
                display(widgets.HBox(list(group_sel_widgets.values())))
            else:
                print("No Eye1/Eye2 classification available; final means will be skipped.")
                return

            final_button = widgets.Button(
                description="Plot final Eye1/Eye2 group means",
                button_style="success",
            )
            display(final_button)
            out_final = widgets.Output()
            display(out_final)

            def on_final_button_clicked(bb):
                nonlocal final_data_dict, final_fig, current_figures
                with out_final:
                    clear_output()

                    final_eye1 = list(group_sel_widgets.get("Eye1", widgets.SelectMultiple()).value)
                    final_eye2 = list(group_sel_widgets.get("Eye2", widgets.SelectMultiple()).value)

                    # Compute group means using channel_means_after (post-exclusion)
                    group_mean_eye1 = compute_group_mean(channel_means_after, final_eye1) if final_eye1 else None
                    group_mean_eye2 = compute_group_mean(channel_means_after, final_eye2) if final_eye2 else None

                    # 4) Plot final Eye1/Eye2 group means as separate subplots
                    n_subplots = 0
                    if show_eye1 and final_eye1:
                        n_subplots += 1
                    if show_eye2 and final_eye2:
                        n_subplots += 1

                    if n_subplots == 0:
                        print("No final Eye1/Eye2 channels selected or toggles disabled.")
                    else:
                        fig_final, axes_final = plt.subplots(
                            1, n_subplots, figsize=(6 * n_subplots, 5)
                        )
                        if n_subplots == 1:
                            axes_final = [axes_final]
                        idx = 0

                        if show_eye1 and final_eye1:
                            ax1 = axes_final[idx]
                            freqs1 = loaded_psd[final_eye1[0]]["freqs"]
                            # Show per-channel "new" means
                            import matplotlib.cm as cm
                            cmap = cm.get_cmap("tab20")
                            colors_for_channels = {}
                            for i, ch in enumerate(sorted(final_eye1)):
                                colors_for_channels[ch] = cmap(i % len(cmap.colors))
                            plot_individual_channels(
                                ax=ax1,
                                channels=final_eye1,
                                channel_means_dict=channel_means_after,
                                freqs=freqs1,
                                colors_for_channels=colors_for_channels,
                                show_original_mean=False,
                                show_new_mean=True,
                                title="Final Eye1 (per-channel new means)",
                                axis_fs=axis_fs,
                                legend_fs=legend_fs,
                                tick_fs=tick_fs,
                            )
                            if group_mean_eye1 is not None:
                                ax1.plot(freqs1, group_mean_eye1, color="black", linewidth=2, label="Eye1 group mean")
                                ax1.legend(fontsize=legend_fs, loc="upper right")
                            if x_min is not None:
                                ax1.set_xlim(left=x_min)
                            if x_max is not None:
                                ax1.set_xlim(right=x_max)
                            if y_min is not None:
                                ax1.set_ylim(bottom=y_min)
                            if y_max is not None:
                                ax1.set_ylim(top=y_max)
                            idx += 1

                        if show_eye2 and final_eye2 and idx < len(axes_final):
                            ax2 = axes_final[idx]
                            freqs2 = loaded_psd[final_eye2[0]]["freqs"]
                            import matplotlib.cm as cm
                            cmap = cm.get_cmap("tab20")
                            colors_for_channels = {}
                            for i, ch in enumerate(sorted(final_eye2)):
                                colors_for_channels[ch] = cmap(i % len(cmap.colors))
                            plot_individual_channels(
                                ax=ax2,
                                channels=final_eye2,
                                channel_means_dict=channel_means_after,
                                freqs=freqs2,
                                colors_for_channels=colors_for_channels,
                                show_original_mean=False,
                                show_new_mean=True,
                                title="Final Eye2 (per-channel new means)",
                                axis_fs=axis_fs,
                                legend_fs=legend_fs,
                                tick_fs=tick_fs,
                            )
                            if group_mean_eye2 is not None:
                                ax2.plot(freqs2, group_mean_eye2, color="black", linewidth=2, label="Eye2 group mean")
                                ax2.legend(fontsize=legend_fs, loc="upper right")
                            if x_min is not None:
                                ax2.set_xlim(left=x_min)
                            if x_max is not None:
                                ax2.set_xlim(right=x_max)
                            if y_min is not None:
                                ax2.set_ylim(bottom=y_min)
                            if y_max is not None:
                                ax2.set_ylim(top=y_max)

                        plt.tight_layout()
                        plt.show()
                        current_figures.append(fig_final)

                    # 5) "Means only" figure: Eye1 vs Eye2 curves
                    if (group_mean_eye1 is None) and (group_mean_eye2 is None):
                        print("No group means to plot as 'means only'.")
                        return

                    fig_mom, ax_mom = plt.subplots(figsize=(6, 4))
                    if group_mean_eye1 is not None and final_eye1:
                        freqs1 = loaded_psd[final_eye1[0]]["freqs"]
                        ax_mom.plot(freqs1, group_mean_eye1, color="red", linewidth=2, label="Eye1 group mean")
                        final_data_dict["Eye1_MeanOfMeans"] = group_mean_eye1
                        final_data_dict["Eye1_freqs"] = freqs1
                    if group_mean_eye2 is not None and final_eye2:
                        freqs2 = loaded_psd[final_eye2[0]]["freqs"]
                        ax_mom.plot(freqs2, group_mean_eye2, color="blue", linewidth=2, label="Eye2 group mean")
                        final_data_dict["Eye2_MeanOfMeans"] = group_mean_eye2
                        final_data_dict["Eye2_freqs"] = freqs2

                    ax_mom.set_title("Final Eye1/Eye2 means of means", fontsize=title_fs)
                    ax_mom.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
                    ax_mom.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
                    ax_mom.tick_params(axis="both", labelsize=tick_fs)
                    ax_mom.legend(fontsize=legend_fs)
                    if x_min is not None:
                        ax_mom.set_xlim(left=x_min)
                    if x_max is not None:
                        ax_mom.set_xlim(right=x_max)
                    if y_min is not None:
                        ax_mom.set_ylim(bottom=y_min)
                    if y_max is not None:
                        ax_mom.set_ylim(top=y_max)
                    plt.tight_layout()
                    plt.show()
                    current_figures.append(fig_mom)
                    final_fig = fig_mom

            final_button.on_click(on_final_button_clicked)

    plot_psd_button.on_click(on_plot_psd_clicked)

    # ------------------------------------------------------------------
    # Export callbacks
    # ------------------------------------------------------------------

    def on_export_button_clicked(b):
        nonlocal current_figures
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No figures to export (plot first).")
                return
            path = export_fig_chooser.selected
            if not path:
                print("Please select or type an export filename (e.g. group_psd.png or group_psd.pptx).")
                return

            base, ext = os.path.splitext(path)
            ext = ext.lower()
            if not ext:
                print("Please provide a file extension (.png, .svg, .jpg, .pptx, etc.)")
                return

            if ext in (".ppt", ".pptx"):
                if not HAS_PPTX:
                    print("python-pptx not installed; cannot export PPT.")
                    return
                prs = Presentation()
                blank_layout = prs.slide_layouts[6]
                for i, fig in enumerate(current_figures, start=1):
                    tmp = f"{base}_tmp_{i}.png"
                    fig.savefig(tmp, dpi=150)
                    slide = prs.slides.add_slide(blank_layout)
                    left = top = Inches(1)
                    slide.shapes.add_picture(tmp, left, top, width=Inches(8), height=Inches(4.5))
                    os.remove(tmp)
                prs.save(path)
                print(f"Exported {len(current_figures)} figure(s) to PPTX:\n  {path}")
            else:
                for i, fig in enumerate(current_figures, start=1):
                    out_file = f"{base}_fig{i}{ext}"
                    fig.savefig(out_file, dpi=150)
                    print(f"Saved: {out_file}")
                print("Image export done.")

    export_button.on_click(on_export_button_clicked)

    def on_final_export_clicked(b):
        nonlocal final_data_dict, final_fig
        with final_export_output:
            clear_output()
            path = final_export_chooser.selected
            if not path:
                print("Please select or type an export filename (e.g. final_means.pkl).")
                return

            base, ext = os.path.splitext(path)
            ext = ext.lower()
            if not ext:
                print("Please include an extension (.pkl, .xlsx, .png, .pptx, etc.).")
                return

            if not final_data_dict and final_fig is None:
                print("No final data/figure found. Run the final group means step first.")
                return

            if ext == ".pkl":
                try:
                    with open(path, "wb") as f:
                        pickle.dump(final_data_dict, f)
                    print(f"Saved final_data_dict => {path}")
                except Exception as e:
                    print(f"Error saving pickle: {e}")
            elif ext in (".xls", ".xlsx"):
                try:
                    import pandas as pd
                    df = {}
                    for k, v in final_data_dict.items():
                        arr = np.asarray(v).ravel()
                        df[k] = arr
                    df = pd.DataFrame(df)
                    df.to_excel(path, index=False)
                    print(f"Saved final_data_dict to Excel => {path}")
                except Exception as e:
                    print(f"Error saving Excel: {e}")
            elif ext in (".png", ".jpg", ".jpeg", ".svg"):
                if final_fig is None:
                    print("No final figure to export.")
                    return
                try:
                    final_fig.savefig(path, dpi=150)
                    print(f"Saved final figure => {path}")
                except Exception as e:
                    print(f"Error saving figure: {e}")
            elif ext in (".ppt", ".pptx"):
                if not HAS_PPTX:
                    print("python-pptx not installed; cannot export PPT.")
                    return
                if final_fig is None:
                    print("No final figure to export.")
                    return
                try:
                    prs = Presentation()
                    blank_layout = prs.slide_layouts[6]
                    tmp = f"{base}_final_tmp.png"
                    final_fig.savefig(tmp, dpi=150)
                    slide = prs.slides.add_slide(blank_layout)
                    left = top = Inches(1)
                    slide.shapes.add_picture(tmp, left, top, width=Inches(8), height=Inches(4.5))
                    os.remove(tmp)
                    prs.save(path)
                    print(f"Exported final figure to PPTX => {path}")
                except Exception as e:
                    print(f"Error exporting PPTX: {e}")
            else:
                print(f"Unsupported extension '{ext}'. Use .pkl, .xlsx, .png, .jpg, .svg, .pptx.")

    final_export_button.on_click(on_final_export_clicked)

    # ------------------------------------------------------------------
    # Layout
    # ------------------------------------------------------------------

    load_box = widgets.VBox(
        [
            widgets.HTML("<h3>1) Load cleaned PSD pickle (from psd_clean_v2d3_NaN)</h3>"),
            widgets.HBox([load_psd_button, psd_file_chooser]),
            load_output_area,
        ]
    )

    controls_left = widgets.VBox(
        [
            widgets.HTML("<h3>2) Channel selection</h3>"),
            channels_dropdown,
        ]
    )

    controls_middle = widgets.VBox(
        [
            widgets.HTML("<h3>3) Eye & mean display</h3>"),
            show_eye1_cb,
            show_eye2_cb,
            show_original_mean_cb,
            show_new_mean_cb,
            widgets.HTML("<b>Exclusion thresholds</b>"),
            low_band_threshold_widget,
            test_band_threshold_widget,
            widgets.HTML("<b>Test bands</b>"),
            test_band_text,
        ]
    )

    controls_right = widgets.VBox(
        [
            widgets.HTML("<h3>4) Axes & fonts</h3>"),
            widgets.HBox([x_min_widget, x_max_widget]),
            widgets.HBox([y_min_widget, y_max_widget]),
            widgets.HBox([title_fs_widget, axis_fs_widget]),
            widgets.HBox([legend_fs_widget, tick_fs_widget]),
            widgets.HTML("<b>Colors (mostly cosmetic here)</b>"),
            widgets.HBox([color_kept_widget, color_excl_widget]),
            widgets.HBox([color_oldmean_widget, color_newmean_widget]),
            widgets.HTML("<br>"),
            plot_psd_button,
        ]
    )

    controls_box = widgets.HBox([controls_left, controls_middle, controls_right])

    export_box = widgets.VBox(
        [
            widgets.HTML("<h3>5) Export figures and final means</h3>"),
            widgets.HTML("<b>All figures (screening + final):</b>"),
            export_fig_chooser,
            export_button,
            export_output_area,
            widgets.HTML("<b>Final Eye1/Eye2 means data/plot:</b>"),
            final_export_chooser,
            final_export_button,
            final_export_output,
        ]
    )

    ui = widgets.VBox(
        [
            widgets.HTML("<h2>PSD Channel Means (NaN-aware, EyeAvg pipeline) — v1b_NaN</h2>"),
            load_box,
            controls_box,
            widgets.HTML("<h3>6) Plots</h3>"),
            plot_output_area,
            export_box,
        ]
    )

    display(ui)


def launch_psd_channel_means_nan_gui():
    """Convenience wrapper to launch the GUI in Jupyter."""
    build_exportable_plot_psd_gui()
