
"""
psd_clean_v2d3_NaN.py

NaN-aware PSD cleaning / plotting GUI for Epochs FIF files (e.g. 2‑channel EyeAvg FIF with NaNs).
Layout is organized into three vertical "levels":

  1. Loading & PSD computation (top row)
  2. Display / exclusion settings (middle row)
  3. Plot output + export controls (bottom rows)

Usage (in Jupyter):

    import sys
    sys.path.append("/Users/davidlitvin/Multielectrode array rsERG/Updated_MEA_Pipeline")
    import psd_clean_v2d3_NaN
    psd_clean_v2d3_NaN.launch_psd_clean_nan_gui()

Requirements:
    - numpy, matplotlib, scipy
    - mne
    - ipywidgets
    - ipyfilechooser
    - python-pptx (optional, only for PPTX export)
"""

import os
import pickle
import math

import numpy as np
import matplotlib.pyplot as plt
from scipy import signal

import mne
from mne.utils import set_log_level

import ipywidgets as widgets
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed, PPT export will be disabled.")

# Keep MNE logs quieter
set_log_level("WARNING")


# ---------------------------------------------------------------------
# 1) Helper functions for NaN-aware PSD and exclusion
# ---------------------------------------------------------------------

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=None,
    test_band_threshold=10.0,
    test_band_count_threshold=None,
):
    """
    NaN-aware trace exclusion.

    Parameters
    ----------
    psd_array : array, shape (n_epochs, n_freqs)
    freqs     : array, shape (n_freqs,)
    low_band  : tuple
        e.g. (1, 3) Hz for low-frequency outliers.
    low_band_threshold : float
        Multiple of mean PSD in low_band used to flag an epoch.
    test_bands : list of (fmin, fmax) or None
        Frequency bands used to detect repeated suprathreshold events.
    test_band_threshold : float
        Multiple of mean PSD in each test band to count as suprathreshold.
    test_band_count_threshold : int or None
        Number of suprathreshold bands needed to exclude an epoch.
        If None, defaults to len(test_bands)//2.

    Returns
    -------
    kept_traces     : list of epoch indices
    excluded_traces : list of epoch indices
    """
    if test_bands is None:
        test_bands = [(7, 9), (9, 11), (11, 13),
                      (13, 15), (15, 17), (17, 19), (19, 21)]
    if test_band_count_threshold is None:
        test_band_count_threshold = max(1, len(test_bands) // 2)

    # NaN-robust mean PSD across epochs
    with np.errstate(invalid="ignore"):
        mean_psd = np.nanmean(psd_array, axis=0)

    excluded_traces = []
    kept_traces = []

    low_band_idx = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    band_indices = [
        np.where((freqs >= b[0]) & (freqs <= b[1]))[0]
        for b in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # If this epoch's PSD is entirely NaN, drop it
        if np.all(np.isnan(trace)):
            excluded_traces.append(i)
            continue

        # 1) Low-frequency outlier
        if low_band_idx.size > 0:
            if np.any(trace[low_band_idx] > low_band_threshold * mean_psd[low_band_idx]):
                excluded_traces.append(i)
                continue

        # 2) Repeated suprathreshold events across test bands
        suprathreshold_count = 0
        for idx in band_indices:
            if idx.size == 0:
                continue
            if np.any(trace[idx] > test_band_threshold * mean_psd[idx]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces


def plot_psds_with_exclusion(
    ax,
    psd_array,
    freqs,
    kept_traces,
    excluded_traces,
    original_mean_psd,
    title,
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    alpha_kept=0.6,
    alpha_excluded=0.8,
    line_width_kept=1.0,
    line_width_excluded=1.5,
    line_width_old_mean=2.0,
    line_width_new_mean=2.0,
    title_fontsize=10,
    axis_label_fontsize=10,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6,
):
    """
    Plot PSDs for one channel/block with kept vs excluded epochs highlighted.
    """
    # Truncate title if long
    if len(title) > max_title_length:
        title = title[: max_title_length - 3] + "..."

    # 1) Kept traces
    if show_kept and kept_traces:
        for idx in kept_traces:
            ax.plot(
                freqs,
                psd_array[idx],
                color=color_kept,
                alpha=alpha_kept,
                linewidth=line_width_kept,
            )

    # 2) Excluded traces
    if show_excluded and excluded_traces:
        for idx in excluded_traces:
            ax.plot(
                freqs,
                psd_array[idx],
                color=color_excluded,
                alpha=alpha_excluded,
                linewidth=line_width_excluded,
            )

    # 3) Original mean PSD
    if show_original_mean and original_mean_psd is not None:
        ax.plot(
            freqs,
            original_mean_psd,
            color=color_old_mean,
            linewidth=line_width_old_mean,
            label="Original Mean",
        )

    # 4) New mean using only kept traces
    if show_new_mean and kept_traces:
        new_mean_psd = np.nanmean(psd_array[kept_traces], axis=0)
        ax.plot(
            freqs,
            new_mean_psd,
            color=color_new_mean,
            linewidth=line_width_new_mean,
            label="New Mean",
        )

    # 5) Optional vertical lines for reference frequencies
    if show_vertical_lines and vertical_lines is not None:
        for vfreq in vertical_lines:
            ax.axvline(
                vfreq,
                color=vertical_line_color,
                linestyle=vertical_line_style,
                alpha=vertical_line_alpha,
            )

    ax.set_xlabel("Frequency (Hz)", fontsize=axis_label_fontsize)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_label_fontsize)
    ax.set_title(title, fontsize=title_fontsize)
    ax.legend(loc="upper right", fontsize=legend_fontsize)
    ax.tick_params(axis="both", labelsize=tick_label_fontsize)


def plot_psds_with_dropped_traces(
    psds_dict,
    rows_of_psds,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=None,
    test_band_threshold=10.0,
    test_band_count_threshold=None,
    x_min=None,
    x_max=None,
    y_min=None,
    y_max=None,
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    num_cols=4,
    title_fontsize=10,
    axis_label_fontsize=10,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6,
):
    """
    Wrapper that loops across PSD entries (channels/blocks),
    applies the exclusion, and arranges subplots in rows/cols.
    """
    if test_bands is None:
        test_bands = [(7, 9), (9, 11), (11, 13),
                      (13, 15), (15, 17), (17, 19), (19, 21)]

    figures = []
    kept_indices_dict = {}

    if vertical_lines is None:
        vertical_lines = [(b[0] + b[1]) / 2.0 for b in test_bands]

    for row_idx, row_keys in enumerate(rows_of_psds):
        n_plots = len(row_keys)
        if n_plots == 0:
            continue

        n_rows = math.ceil(n_plots / num_cols)
        fig, axes = plt.subplots(
            nrows=n_rows, ncols=num_cols, figsize=(4 * num_cols, 3 * n_rows)
        )
        axes = np.atleast_1d(axes).flatten()

        for i, key in enumerate(row_keys):
            if key not in psds_dict:
                continue
            info = psds_dict[key]
            freqs = info.get("freqs", None)
            psd_data = info.get("psd", None)
            if freqs is None or psd_data is None:
                continue
            if not isinstance(freqs, np.ndarray) or not isinstance(psd_data, np.ndarray):
                continue

            ax = axes[i]

            # NaN-robust original mean
            original_mean_psd = np.nanmean(psd_data, axis=0)

            kept_traces, excluded_traces = exclude_traces(
                psd_data,
                freqs,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold,
            )
            kept_indices_dict[key] = kept_traces

            plot_psds_with_exclusion(
                ax=ax,
                psd_array=psd_data,
                freqs=freqs,
                kept_traces=kept_traces,
                excluded_traces=excluded_traces,
                original_mean_psd=original_mean_psd,
                title=key,
                show_kept=show_kept,
                show_excluded=show_excluded,
                show_original_mean=show_original_mean,
                show_new_mean=show_new_mean,
                show_vertical_lines=show_vertical_lines,
                color_kept=color_kept,
                color_excluded=color_excluded,
                color_old_mean=color_old_mean,
                color_new_mean=color_new_mean,
                title_fontsize=title_fontsize,
                axis_label_fontsize=axis_label_fontsize,
                legend_fontsize=legend_fontsize,
                tick_label_fontsize=tick_label_fontsize,
                max_title_length=max_title_length,
                vertical_lines=vertical_lines,
                vertical_line_color=vertical_line_color,
                vertical_line_style=vertical_line_style,
                vertical_line_alpha=vertical_line_alpha,
            )

            if x_min is not None or x_max is not None:
                ax.set_xlim(x_min, x_max)
            if y_min is not None or y_max is not None:
                ax.set_ylim(y_min, y_max)

        # Turn off unused axes
        for ax in axes[n_plots:]:
            ax.axis("off")

        plt.tight_layout(rect=[0, 0.03, 1, 0.95])
        plt.suptitle(f"Row {row_idx}", fontsize=title_fontsize + 2)
        figures.append(fig)
        plt.show()

    return figures, kept_indices_dict


def group_keys_by_rows(psd_keys, row_size=4):
    """
    Sort keys by block index (if any) and channel index, then chunk into rows.
    Accepts keys like "Ch1", "Eye1_avg", "block1:Eye1_avg", etc.
    """
    def sort_key(key):
        # Block part
        if ":" in key:
            block, ch = key.split(":", 1)
            try:
                block_num = int(block.replace("block", ""))
            except ValueError:
                block_num = 0
        else:
            block_num = 0
            ch = key

        # Channel numeric part if present (Ch1, Eye1_avg)
        ch_num = 99999
        if ch.startswith("Ch"):
            try:
                ch_num = int(ch.replace("Ch", ""))
            except ValueError:
                ch_num = 99999
        elif ch.lower().startswith("eye"):
            digits = "".join(c for c in ch if c.isdigit())
            if digits:
                ch_num = int(digits)

        return (block_num, ch_num, ch)

    sorted_keys = sorted(psd_keys, key=sort_key)
    rows = []
    for i in range(0, len(sorted_keys), row_size):
        rows.append(sorted_keys[i:i + row_size])
    return rows


def load_epochs(filepath):
    """
    Load mne.Epochs from .fif or .fif.gz. If needed, reconstruct from
    data/info using EpochsArray.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found: {filepath}")
    if not (filepath.endswith(".fif") or filepath.endswith(".fif.gz")):
        raise ValueError("Input file must be .fif or .fif.gz")

    epochs = mne.read_epochs(filepath, preload=True)
    if not isinstance(epochs, mne.Epochs):
        data = epochs.get_data()
        info = epochs.info
        epochs = mne.EpochsArray(data, info)
    return epochs


def _welch_longest_clean_segment_1d(x, sfreq, nperseg, noverlap, window):
    """
    NaN-aware Welch on 1D array x using only the longest contiguous run
    of finite samples >= nperseg. Returns (freqs, psd) or (None, None).
    """
    x = np.asarray(x, dtype=float)
    mask = np.isfinite(x)
    if mask.sum() < nperseg:
        return None, None

    finite_idx = np.where(mask)[0]
    starts = [finite_idx[0]]
    ends = []

    for i in range(1, len(finite_idx)):
        if finite_idx[i] != finite_idx[i-1] + 1:
            ends.append(finite_idx[i-1])
            starts.append(finite_idx[i])
    ends.append(finite_idx[-1])

    lengths = [e - s + 1 for s, e in zip(starts, ends)]
    best_idx = int(np.argmax(lengths))
    s = starts[best_idx]
    e = ends[best_idx]
    seg = x[s:e+1]

    if seg.size < nperseg:
        return None, None

    freqs, psd = signal.welch(
        seg,
        fs=sfreq,
        window=window,
        nperseg=nperseg,
        noverlap=noverlap,
        scaling="density",
        axis=0,
    )
    return freqs, psd


def compute_welch_psd(
    epochs,
    fmin=0.0,
    fmax=45.0,
    window_s=2.0,
    overlap_pct=50.0,
):
    """
    NaN-aware Welch PSD for an mne.Epochs object.

    For each channel and epoch, finds the longest contiguous run of
    finite samples and applies a standard Welch PSD to that segment.
    """
    out = {}
    sfreq = epochs.info["sfreq"]
    data = epochs.get_data()  # (n_ep, n_ch, n_time)
    n_ep, n_ch, n_time = data.shape

    nperseg = int(window_s * sfreq)
    noverlap = int(nperseg * overlap_pct / 100.0)

    if nperseg > n_time:
        raise ValueError(f"nperseg ({nperseg}) > n_times per epoch ({n_time}).")

    window = signal.hamming(nperseg, sym=False)

    for ch_idx, ch_name in enumerate(epochs.ch_names):
        ch_data = data[:, ch_idx, :]  # (n_ep, n_time)
        freqs_ref = None
        psd_rows = []

        for ep_idx in range(n_ep):
            x = ch_data[ep_idx]
            freqs_ep, psd_ep = _welch_longest_clean_segment_1d(
                x, sfreq, nperseg, noverlap, window
            )
            if freqs_ep is None or psd_ep is None:
                # Not enough clean data
                if freqs_ref is None:
                    psd_rows.append(None)
                else:
                    psd_rows.append(np.full_like(freqs_ref, np.nan))
                continue

            # First epoch with valid PSD => set fmin/fmax and reference freqs
            if freqs_ref is None:
                freq_mask = (freqs_ep >= fmin) & (freqs_ep <= fmax)
                freqs_ref = freqs_ep[freq_mask]
                psd_sel = psd_ep[freq_mask]
            else:
                freq_mask = (freqs_ep >= fmin) & (freqs_ep <= fmax)
                psd_sel = psd_ep[freq_mask]
                if psd_sel.shape != freqs_ref.shape:
                    raise RuntimeError(
                        f"Frequency grid mismatch in channel '{ch_name}'. "
                        "Check Welch parameters."
                    )
            psd_rows.append(psd_sel)

        if freqs_ref is None:
            out[ch_name] = {"freqs": None, "psd": None}
            continue

        psd_array = []
        for row in psd_rows:
            if row is None:
                psd_array.append(np.full_like(freqs_ref, np.nan))
            else:
                psd_array.append(row)
        psd_array = np.vstack(psd_array)

        out[ch_name] = {"freqs": freqs_ref, "psd": psd_array}

    return out


# ---------------------------------------------------------------------
# 2) Main GUI
# ---------------------------------------------------------------------

def build_exportable_plot_psd_gui():
    """
    Build a vertical, 3‑row style GUI:

      Row 1: Loading & PSD computation
      Row 2: Display / exclusion settings
      Row 3: Plot output + export controls
    """
    # ---------- STORAGE ----------
    loaded_psd = {}         # key -> {"freqs": array, "psd": array}
    current_figures = []    # list of matplotlib Figure objects
    kept_indices_dict = {}  # key -> list of epoch indices



    # ---------- ROW 1: LOADING / PSD COMPUTATION ----------

    fif_chooser = FileChooser(
        os.getcwd(),
        title="1) Select input Epochs .fif / .fif.gz (NaNs allowed)",
        select_default=False,
    )
    fif_chooser.show_only_files = True
    fif_chooser.filter_pattern = ["*.fif", "*.fif.gz"]

    segment_toggle = widgets.ToggleButtons(
        options=["No", "Yes"],
        value="No",
        description="Segment into 6 blocks of 120 epochs?",
        style={"description_width": "initial"},
    )

    window_length_widget = widgets.FloatText(
        value=2.0,
        description="Welch window (s):",
        layout=widgets.Layout(width="220px"),
    )
    overlap_widget = widgets.FloatSlider(
        value=50.0,
        min=0.0,
        max=95.0,
        step=1.0,
        description="Overlap (%):",
        continuous_update=False,
        layout=widgets.Layout(width="350px"),
    )

    load_psd_button = widgets.Button(
        description="Load & Compute PSD",
        button_style="success",
        layout=widgets.Layout(width="200px"),
    )
    recompute_psd_button = widgets.Button(
        description="Recompute PSD (same file)",
        button_style="info",
        layout=widgets.Layout(width="220px"),
    )

    load_output_area = widgets.Output(layout=widgets.Layout(border="1px solid #ccc"))

    # ---------- ROW 2: DISPLAY / EXCLUSION SETTINGS ----------

    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description="Channels / Blocks:",
        layout=widgets.Layout(width="260px", height="200px"),
    )

    # Show toggles
    show_kept_cb = widgets.Checkbox(
        value=True, description="Show kept traces", indent=False
    )
    show_excluded_cb = widgets.Checkbox(
        value=True, description="Show excluded traces", indent=False
    )
    show_original_mean_cb = widgets.Checkbox(
        value=True, description="Show original mean", indent=False
    )
    show_new_mean_cb = widgets.Checkbox(
        value=True, description="Show new mean (kept only)", indent=False
    )
    show_vertical_lines_cb = widgets.Checkbox(
        value=True, description="Show band centers", indent=False
    )

    # Colors
    color_kept_widget = widgets.ColorPicker(
        value="#D3D3D3",
        description="Kept color:",
        layout=widgets.Layout(width="250px"),
    )
    color_excl_widget = widgets.ColorPicker(
        value="#FF0000",
        description="Excluded color:",
        layout=widgets.Layout(width="250px"),
    )
    color_old_mean_widget = widgets.ColorPicker(
        value="#0000FF",
        description="Original mean:",
        layout=widgets.Layout(width="250px"),
    )
    color_new_mean_widget = widgets.ColorPicker(
        value="#00AA00",
        description="New mean:",
        layout=widgets.Layout(width="250px"),
    )

    # Axis ranges
    x_min_widget = widgets.FloatText(
        value=0.0,
        description="X min (Hz):",
        layout=widgets.Layout(width="160px"),
    )
    x_max_widget = widgets.FloatText(
        value=45.0,
        description="X max (Hz):",
        layout=widgets.Layout(width="160px"),
    )
    y_min_widget = widgets.FloatText(
        value=0.0,
        description="Y min:",
        layout=widgets.Layout(width="160px"),
    )
    y_max_widget = widgets.FloatText(
        value=1.0,
        description="Y max:",
        layout=widgets.Layout(width="160px"),
    )

    # Fonts
    title_fs_widget = widgets.IntSlider(
        value=10,
        min=6,
        max=20,
        step=1,
        description="Title FS:",
        continuous_update=False,
        layout=widgets.Layout(width="260px"),
    )
    axis_fs_widget = widgets.IntSlider(
        value=10,
        min=6,
        max=20,
        step=1,
        description="Axis label FS:",
        continuous_update=False,
        layout=widgets.Layout(width="260px"),
    )
    legend_fs_widget = widgets.IntSlider(
        value=8,
        min=6,
        max=20,
        step=1,
        description="Legend FS:",
        continuous_update=False,
        layout=widgets.Layout(width="260px"),
    )
    tick_fs_widget = widgets.IntSlider(
        value=8,
        min=6,
        max=20,
        step=1,
        description="Tick FS:",
        continuous_update=False,
        layout=widgets.Layout(width="260px"),
    )
    max_title_len_widget = widgets.IntSlider(
        value=40,
        min=10,
        max=80,
        step=1,
        description="Max title length:",
        continuous_update=False,
        layout=widgets.Layout(width="260px"),
    )

    # Exclusion bands
    low_band_min_widget = widgets.FloatText(
        value=1.0,
        description="Low band min (Hz):",
        layout=widgets.Layout(width="190px"),
    )
    low_band_max_widget = widgets.FloatText(
        value=3.0,
        description="Low band max (Hz):",
        layout=widgets.Layout(width="190px"),
    )
    low_band_thresh_widget = widgets.FloatText(
        value=3.0,
        description="Low band thresh (x mean):",
        layout=widgets.Layout(width="230px"),
    )

    test_band_min_widget = widgets.FloatText(
        value=7.0,
        description="Test band start (Hz):",
        layout=widgets.Layout(width="210px"),
    )
    test_band_max_widget = widgets.FloatText(
        value=21.0,
        description="Test band end (Hz):",
        layout=widgets.Layout(width="210px"),
    )
    test_band_step_widget = widgets.FloatText(
        value=2.0,
        description="Test band step (Hz):",
        layout=widgets.Layout(width="210px"),
    )
    test_band_thresh_widget = widgets.FloatText(
        value=10.0,
        description="Test band thresh (x mean):",
        layout=widgets.Layout(width="240px"),
    )
    test_band_count_widget = widgets.IntText(
        value=3,
        description="#bands over thresh to exclude:",
        layout=widgets.Layout(width="270px"),
    )

    plot_psd_button = widgets.Button(
        description="Plot PSDs", button_style="success",
        layout=widgets.Layout(width="180px"),
    )

    # ---------- ROW 3: PLOT OUTPUT & EXPORT ----------

    plot_output_area = widgets.Output(layout=widgets.Layout(border="1px solid #ccc"))

    export_fig_chooser = FileChooser(
        os.getcwd(),
        title="Figure export (e.g. MyPsd.png or MyPsd.pptx):",
        select_default=False,
    )

    export_button = widgets.Button(
        description="Export figures", button_style="warning",
        layout=widgets.Layout(width="160px"),
    )

    cleaned_pickle_chooser = FileChooser(
        os.getcwd(),
        title="Export cleaned PSD .pkl:",
        select_default=False,
    )
    cleaned_pickle_chooser.filter_pattern = ["*.pkl"]

    export_cleaned_button = widgets.Button(
        description="Export cleaned PSD", button_style="info",
        layout=widgets.Layout(width="180px"),
    )

    export_output_area = widgets.Output(layout=widgets.Layout(border="1px solid #ccc"))

    # ------------------------------------------------------------------
    # Callbacks
    # ------------------------------------------------------------------

    def _compute_psd_for_current_file():
        """
        Internal helper to compute Welch PSD from the selected FIF using
        current window/overlap and segmentation settings.
        Updates 'loaded_psd' in outer scope.
        """
        nonlocal loaded_psd

        with load_output_area:
            clear_output()
            fif_path = fif_chooser.selected
            if not fif_path:
                print("No .fif file selected.")
                return
            if not os.path.isfile(fif_path):
                print(f"File not found: {fif_path}")
                return

            print(f"Loading epochs from: {fif_path}")
            epochs_all = load_epochs(fif_path)
            print(f"  Loaded {len(epochs_all)} epochs, {len(epochs_all.ch_names)} channels.")
            window_s = window_length_widget.value
            overlap_pct = overlap_widget.value

            if window_s <= 0:
                print("ERROR: window length must be > 0.")
                return
            if not (0 <= overlap_pct < 100):
                print("ERROR: overlap% must be in [0, 100).")
                return

            seg = segment_toggle.value
            loaded_psd = {}

            if seg == "No":
                print("No segmentation: computing PSD on full epoch set.")
                psd_dict = compute_welch_psd(
                    epochs_all, fmin=0.0, fmax=45.0,
                    window_s=window_s, overlap_pct=overlap_pct,
                )
                for ch in epochs_all.ch_names:
                    loaded_psd[ch] = psd_dict[ch]
            else:
                print("Segmenting into blocks of 120 epochs (up to 6 blocks).")
                n_total = len(epochs_all)
                block_size = 120
                n_blocks = min(6, n_total // block_size)

                for bi in range(n_blocks):
                    start_ep = bi * block_size
                    end_ep = (bi + 1) * block_size
                    block_epochs = epochs_all[start_ep:end_ep]
                    block_name = f"block{bi+1}"
                    print(f"  {block_name}: epochs {start_ep}..{end_ep-1}")
                    psd_block_dict = compute_welch_psd(
                        block_epochs,
                        fmin=0.0,
                        fmax=45.0,
                        window_s=window_s,
                        overlap_pct=overlap_pct,
                    )
                    for ch in block_epochs.ch_names:
                        key = f"{block_name}:{ch}"
                        loaded_psd[key] = psd_block_dict[ch]

                leftover = n_total % block_size
                if leftover > 0:
                    print(f"  Leftover {leftover} epochs not segmented.")

            all_keys = list(loaded_psd.keys())
            if not all_keys:
                print("No PSD entries computed.")
                channels_dropdown.options = []
                return

            channels_dropdown.options = all_keys
            # Preselect up to first 4
            channels_dropdown.value = tuple(all_keys[: min(4, len(all_keys))])
            print(f"Computed PSD for {len(all_keys)} channel/block keys.")

    def on_load_psd_clicked(b):
        _compute_psd_for_current_file()

    def on_recompute_psd_clicked(b):
        _compute_psd_for_current_file()

    load_psd_button.on_click(on_load_psd_clicked)
    recompute_psd_button.on_click(on_recompute_psd_clicked)

    def on_plot_psd_clicked(b):
        nonlocal current_figures, kept_indices_dict
        with plot_output_area:
            clear_output()

            if not loaded_psd:
                print("No PSD data loaded. Use 'Load & Compute PSD' first.")
                return

            selected_keys = list(channels_dropdown.value)
            if not selected_keys:
                print("No channels/blocks selected.")
                return

            show_kept = show_kept_cb.value
            show_excluded = show_excluded_cb.value
            show_original_mean = show_original_mean_cb.value
            show_new_mean = show_new_mean_cb.value
            show_vertical_lines = show_vertical_lines_cb.value

            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value

            title_fs = title_fs_widget.value
            axis_fs = axis_fs_widget.value
            legend_fs = legend_fs_widget.value
            tick_fs = tick_fs_widget.value
            max_title_len = max_title_len_widget.value

            color_kept = color_kept_widget.value
            color_excl = color_excl_widget.value
            color_old_mean = color_old_mean_widget.value
            color_new_mean = color_new_mean_widget.value

            low_band = (low_band_min_widget.value, low_band_max_widget.value)
            low_band_thresh = low_band_thresh_widget.value

            tb_min = test_band_min_widget.value
            tb_max = test_band_max_widget.value
            tb_step = test_band_step_widget.value
            if tb_step <= 0:
                print("ERROR: test band step must be > 0.")
                return

            # Build test bands from min/max/step
            test_bands = []
            centers = []
            f_start = tb_min
            while f_start < tb_max:
                f_end = min(f_start + tb_step, tb_max)
                test_bands.append((f_start, f_end))
                centers.append((f_start + f_end) / 2.0)
                f_start = f_end

            test_band_thresh = test_band_thresh_widget.value
            test_band_count = test_band_count_widget.value

            # Build subset dict only for selected keys
            psd_sub_dict = {k: loaded_psd[k] for k in selected_keys if k in loaded_psd}
            rows_of_psds = group_keys_by_rows(selected_keys, row_size=4)

            figs, kept_dict = plot_psds_with_dropped_traces(
                psds_dict=psd_sub_dict,
                rows_of_psds=rows_of_psds,
                low_band=low_band,
                low_band_threshold=low_band_thresh,
                test_bands=test_bands,
                test_band_threshold=test_band_thresh,
                test_band_count_threshold=test_band_count,
                x_min=x_min,
                x_max=x_max,
                y_min=y_min,
                y_max=y_max,
                show_kept=show_kept,
                show_excluded=show_excluded,
                show_original_mean=show_original_mean,
                show_new_mean=show_new_mean,
                show_vertical_lines=show_vertical_lines,
                color_kept=color_kept,
                color_excluded=color_excl,
                color_old_mean=color_old_mean,
                color_new_mean=color_new_mean,
                num_cols=4,
                title_fontsize=title_fs,
                axis_label_fontsize=axis_fs,
                legend_fontsize=legend_fs,
                tick_label_fontsize=tick_fs,
                max_title_length=max_title_len,
                vertical_lines=centers,
                vertical_line_color="black",
                vertical_line_style="--",
                vertical_line_alpha=0.6,
            )

            current_figures = figs
            kept_indices_dict = kept_dict

            if figs:
                print(f"Plotted {len(figs)} figure(s) for {len(selected_keys)} keys.")
                print("Kept epoch indices stored in kept_indices_dict.")
            else:
                print("No figures created (check PSD / selection).")

    plot_psd_button.on_click(on_plot_psd_clicked)

    def on_export_figures_clicked(b):
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No figures to export (plot first).")
                return
            export_path = export_fig_chooser.selected
            if not export_path:
                print("No export path selected.")
                return

            base, ext = os.path.splitext(export_path)
            ext = ext.lower()

            if ext in (".ppt", ".pptx"):
                if not HAS_PPTX:
                    print("python-pptx not installed; cannot export PPTX.")
                    return
                prs = Presentation()
                blank_layout = prs.slide_layouts[6]
                for i, fig in enumerate(current_figures, start=1):
                    temp_png = f"{base}_fig{i}.png"
                    fig.savefig(temp_png, dpi=150)
                    slide = prs.slides.add_slide(blank_layout)
                    left = Inches(1)
                    top = Inches(1)
                    slide.shapes.add_picture(temp_png, left, top, width=Inches(8))
                    os.remove(temp_png)
                prs.save(export_path)
                print(f"Exported {len(current_figures)} figure(s) to PPTX: {export_path}")
            else:
                if ext == "":
                    ext = ".png"
                for i, fig in enumerate(current_figures, start=1):
                    out_file = f"{base}_fig{i}{ext}"
                    fig.savefig(out_file, dpi=150)
                    print(f"Saved: {out_file}")
                print("Figure export done.")

    export_button.on_click(on_export_figures_clicked)

    def on_export_cleaned_psd_clicked(b):
        with export_output_area:
            clear_output()
            if not kept_indices_dict:
                print("No kept_indices_dict (plot first to compute exclusions).")
                return
            if not loaded_psd:
                print("No PSD data loaded.")
                return

            out_pkl = cleaned_pickle_chooser.selected
            if not out_pkl:
                print("No output .pkl selected.")
                return
            if not out_pkl.endswith(".pkl"):
                print("Output must end with '.pkl'.")
                return

            cleaned = {}
            for key, info in loaded_psd.items():
                if key not in kept_indices_dict:
                    continue
                kept_ix = kept_indices_dict[key]
                psd_array = info.get("psd", None)
                freqs = info.get("freqs", None)
                if psd_array is None or freqs is None:
                    continue
                if len(kept_ix) == 0:
                    new_psd = None
                else:
                    new_psd = psd_array[kept_ix, :]
                cleaned[key] = {"freqs": freqs, "psd": new_psd}

            if not cleaned:
                print("No cleaned data produced (maybe all epochs excluded).")
                return

            try:
                with open(out_pkl, "wb") as f:
                    pickle.dump(cleaned, f)
                print(f"Saved cleaned PSD pickle => {out_pkl}")
                print(f"Channels/blocks in cleaned PSD: {list(cleaned.keys())}")
            except Exception as e:
                print(f"Error saving cleaned PSD: {e}")

    export_cleaned_button.on_click(on_export_cleaned_psd_clicked)

    # ------------------------------------------------------------------
    # Layout: 3 vertical "rows"
    # ------------------------------------------------------------------

    row1 = widgets.VBox(
        [
            widgets.HTML("<h3>Row 1 — Load Epochs & Compute NaN‑aware PSD</h3>"),
            fif_chooser,
            widgets.HBox([segment_toggle]),
            widgets.HBox([window_length_widget, overlap_widget]),
            widgets.HBox([load_psd_button, recompute_psd_button]),
            load_output_area,
        ]
    )

    row2_left = widgets.VBox(
        [
            widgets.HTML("<h4>Channel / Block selection</h4>"),
            channels_dropdown,
        ]
    )

    row2_middle = widgets.VBox(
        [
            widgets.HTML("<h4>What to show</h4>"),
            show_kept_cb,
            show_excluded_cb,
            show_original_mean_cb,
            show_new_mean_cb,
            show_vertical_lines_cb,
            widgets.HTML("<h4>Colors</h4>"),
            color_kept_widget,
            color_excl_widget,
            color_old_mean_widget,
            color_new_mean_widget,
        ]
    )

    row2_right = widgets.VBox(
        [
            widgets.HTML("<h4>Axis & fonts</h4>"),
            widgets.HBox([x_min_widget, x_max_widget]),
            widgets.HBox([y_min_widget, y_max_widget]),
            title_fs_widget,
            axis_fs_widget,
            legend_fs_widget,
            tick_fs_widget,
            max_title_len_widget,
            widgets.HTML("<h4>Exclusion bands</h4>"),
            widgets.HBox([low_band_min_widget, low_band_max_widget]),
            low_band_thresh_widget,
            widgets.HTML("<h4>Oscillation test bands</h4>"),
            widgets.HBox([test_band_min_widget, test_band_max_widget]),
            test_band_step_widget,
            test_band_thresh_widget,
            test_band_count_widget,
            widgets.HTML("<br>"),
            plot_psd_button,
        ]
    )

    row2 = widgets.VBox(
        [
            widgets.HTML("<h3>Row 2 — Display & Exclusion Settings</h3>"),
            widgets.HBox([row2_left, row2_middle, row2_right]),
        ]
    )

    row3 = widgets.VBox(
        [
            widgets.HTML("<h3>Row 3 — Plot Output & Export</h3>"),
            plot_output_area,
            widgets.HTML("<h4>Export</h4>"),
            widgets.HTML("Figures:"),
            export_fig_chooser,
            export_button,
            widgets.HTML("Cleaned PSD (.pkl):"),
            cleaned_pickle_chooser,
            export_cleaned_button,
            export_output_area,
        ]
    )

    ui = widgets.VBox(
        [
            widgets.HTML(
                "<h2>PSD Clean GUI (NaN‑aware, 2‑ch EyeAvg FIF compatible) — v2d3_NaN</h2>"
            ),
            row1,
            row2,
            row3,
        ]
    )

    display(ui)


def launch_psd_clean_nan_gui():
    """
    Convenience wrapper to create and display the GUI.
    """
    build_exportable_plot_psd_gui()
